#!/usr/bin/env python3
"""Test slides 1-8 to isolate the breaking point."""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from chuk_mcp_pptx.components.charts import (
    ColumnChart, BarChart, WaterfallChart, LineChart, AreaChart, PieChart
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def add_slide_title_and_description(slide, title: str, description: str, features: list, theme_obj):
    """Add title and description to a slide."""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    title_para.alignment = PP_ALIGN.CENTER


async def test_first_8_slides():
    """Test the first 8 slides step by step."""
    
    print("Testing slides 1-8 to find the breaking point...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    theme_manager = ThemeManager()
    theme_obj = theme_manager.get_theme("dark-blue")
    
    try:
        # SLIDE 1: Title
        print("Creating slide 1...")
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        theme_obj.apply_to_slide(slide1)
        slide1.shapes.title.text = "Test Showcase"
        slide1.placeholders[1].text = "Testing slides 1-8"
        
        # SLIDE 2: Column Chart (Clustered)
        print("Creating slide 2...")
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide2)
        add_slide_title_and_description(slide2, "02. Column Chart (Clustered)", "Test", [], theme_obj)
        
        column_chart = ColumnChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series={"Revenue": [100, 120, 140, 160], "Profit": [20, 25, 30, 35]},
            variant="clustered",
            theme=theme_obj.__dict__
        )
        await column_chart.render(slide2, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 2 completed")
        
        # SLIDE 3: Column Chart (Stacked)
        print("Creating slide 3...")
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide3)
        add_slide_title_and_description(slide3, "03. Column Chart (Stacked)", "Test", [], theme_obj)
        
        stacked_column = ColumnChart(
            categories=["2021", "2022", "2023", "2024"],
            series={"Product A": [25, 30, 35, 40], "Product B": [20, 25, 30, 35], "Product C": [15, 20, 25, 30]},
            variant="stacked",
            theme=theme_obj.__dict__
        )
        await stacked_column.render(slide3, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 3 completed")
        
        # SLIDE 4: Bar Chart (THE PROBLEM SLIDE?)
        print("Creating slide 4...")
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide4)
        add_slide_title_and_description(slide4, "04. Bar Chart (Horizontal)", "Test", [], theme_obj)
        
        bar_chart = BarChart(
            categories=["Artificial Intelligence", "Cloud Computing", "Cybersecurity", "IoT Solutions", "Blockchain", "AR/VR"],
            series={
                "2024 Adoption (%)": [88, 92, 78, 65, 45, 32],
                "2025 Projected (%)": [95, 96, 85, 75, 58, 48]
            },
            variant="clustered",
            theme=theme_obj.__dict__
        )
        await bar_chart.render(slide4, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 4 completed")
        
        # Test save at this point
        prs.save("outputs/test_slides_1_4.pptx")
        print("✅ Saved slides 1-4 successfully")
        
        # SLIDE 5: Waterfall Chart
        print("Creating slide 5...")
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide5)
        add_slide_title_and_description(slide5, "05. Waterfall Chart", "Test", [], theme_obj)
        
        waterfall_chart = WaterfallChart(
            categories=["Starting Revenue", "New Sales", "Upsells", "Marketing Costs", "R&D Investment", "Operating Costs", "Final Profit"],
            values=[100, 65, 25, -15, -22, -18, 135],
            theme=theme_obj.__dict__
        )
        await waterfall_chart.render(slide5, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 5 completed")
        
        # SLIDE 6: Line Chart
        print("Creating slide 6...")
        slide6 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide6)
        add_slide_title_and_description(slide6, "06. Line Chart", "Test", [], theme_obj)
        
        line_chart = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"],
            series={
                "Website Traffic": [1200, 1350, 1480, 1620, 1850, 2100, 2350, 2650],
                "Mobile App": [800, 920, 1050, 1180, 1350, 1520, 1750, 1980],
                "API Calls": [450, 520, 680, 820, 950, 1100, 1280, 1450]
            },
            theme=theme_obj.__dict__
        )
        await line_chart.render(slide6, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 6 completed")
        
        # SLIDE 7: Area Chart
        print("Creating slide 7...")
        slide7 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide7)
        add_slide_title_and_description(slide7, "07. Area Chart", "Test", [], theme_obj)
        
        area_chart = AreaChart(
            categories=["2020", "2021", "2022", "2023", "2024", "2025"],
            series={
                "Enterprise": [45, 52, 68, 82, 95, 110],
                "SMB": [32, 38, 45, 58, 72, 85],
                "Startup": [18, 25, 32, 41, 55, 68]
            },
            variant="stacked",
            transparency=35,
            theme=theme_obj.__dict__
        )
        await area_chart.render(slide7, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 7 completed")
        
        # SLIDE 8: Pie Chart (THE OTHER PROBLEM SLIDE?)
        print("Creating slide 8...")
        slide8 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide8)
        add_slide_title_and_description(slide8, "08. Pie Chart", "Test", [], theme_obj)
        
        pie_chart = PieChart(
            categories=["SaaS Platform", "Enterprise Services", "Mobile Apps", "API Revenue", "Consulting"],
            values=[42, 28, 18, 8, 4],
            explode=0,  # Explode largest slice
            theme=theme_obj.__dict__,
            options={"show_percentages": True, "show_categories": True}
        )
        await pie_chart.render(slide8, left=2.5, top=1.2, width=5.0, height=4.0)
        print("✅ Slide 8 completed")
        
        # Final save
        prs.save("outputs/test_slides_1_8.pptx")
        print("✅ Saved slides 1-8 successfully")
        
    except Exception as e:
        print(f"❌ Error at slide creation: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    asyncio.run(test_first_8_slides())