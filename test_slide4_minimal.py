#!/usr/bin/env python3
"""Test slide 4 with minimal setup to isolate the issue."""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.charts import ColumnChart, BarChart
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


async def test_slide4_in_context():
    """Test slide 4 in the exact context of comprehensive showcase."""
    
    print("Testing slide 4 in comprehensive showcase context...")
    
    # Create presentation with same structure as comprehensive showcase
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    theme_manager = ThemeManager()
    theme_obj = theme_manager.get_theme("dark-blue")
    
    try:
        # SLIDE 1: Title (same as comprehensive)
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        theme_obj.apply_to_slide(slide1)
        slide1.shapes.title.text = f"{theme_obj.name} Chart Showcase"
        slide1.placeholders[1].text = "Comprehensive Gallery • One Chart Per Slide"
        print("✅ Slide 1 completed")
        
        # SLIDE 2: Column Chart (Clustered) - same as comprehensive
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide2)
        
        column_chart = ColumnChart(
            categories=["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024"],
            series={
                "Revenue ($M)": [125, 148, 172, 195],
                "Profit ($M)": [28, 35, 42, 48],
                "Growth (%)": [15, 22, 28, 35]
            },
            variant="clustered",
            theme=theme_obj.__dict__
        )
        await column_chart.render(slide2, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 2 completed")
        
        # SLIDE 3: Column Chart (Stacked) - same as comprehensive
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide3)
        
        stacked_column = ColumnChart(
            categories=["2021", "2022", "2023", "2024"],
            series={
                "Product A ($M)": [25, 30, 35, 40],
                "Product B ($M)": [20, 25, 30, 35],
                "Product C ($M)": [15, 20, 25, 30]
            },
            variant="stacked",
            theme=theme_obj.__dict__
        )
        await stacked_column.render(slide3, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 3 completed")
        
        # SLIDE 4: Bar Chart (THE PROBLEM SLIDE) - MINIMAL VERSION
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide4)
        
        # NO title function - just the chart
        bar_chart = BarChart(
            categories=["AI/ML", "Cloud", "Security"],  # Even fewer categories
            series={
                "Adoption": [88, 92, 78]  # Single series
            },
            theme=theme_obj.__dict__
        )
        await bar_chart.render(slide4, left=2.0, top=2.0, width=6.0, height=3.0)
        print("✅ Slide 4 completed (minimal)")
        
        # Save to test
        prs.save("outputs/test_slide4_context.pptx")
        print("✅ Saved test with slides 1-4")
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()


async def test_slide4_different_chart_type():
    """Test replacing slide 4 bar chart with a different chart type."""
    
    print("\nTesting slide 4 with different chart type...")
    
    prs = Presentation()
    theme_manager = ThemeManager()
    theme_obj = theme_manager.get_theme("dark-blue")
    
    try:
        # Create 4 slides, but replace slide 4 with a simple column chart
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            theme_obj.apply_to_slide(slide)
        
        # SLIDE 4: Use ColumnChart instead of BarChart
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide4)
        
        replacement_chart = ColumnChart(
            categories=["AI/ML", "Cloud", "Security"],
            series={"Adoption": [88, 92, 78]},
            theme=theme_obj.__dict__
        )
        await replacement_chart.render(slide4, left=2.0, top=2.0, width=6.0, height=3.0)
        
        prs.save("outputs/test_slide4_replacement.pptx")
        print("✅ Saved test with ColumnChart replacement on slide 4")
        
    except Exception as e:
        print(f"❌ Error with replacement: {e}")
        import traceback
        traceback.print_exc()


async def main():
    await test_slide4_in_context()
    await test_slide4_different_chart_type()


if __name__ == "__main__":
    asyncio.run(main())