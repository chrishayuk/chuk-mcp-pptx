#!/usr/bin/env python3
"""Test the exact Bar Chart from slide 4 to isolate the issue."""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from chuk_mcp_pptx.components.charts import BarChart
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def add_slide_title_and_description(slide, title: str, description: str, features: list, theme_obj):
    """Add title and description to a slide."""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    title_para.alignment = PP_ALIGN.CENTER
    
    # Description box
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(4), Inches(1.8))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(14)
    desc_para.font.color.rgb = theme_obj.get_color("muted.foreground")
    
    # Features box
    features_box = slide.shapes.add_textbox(Inches(5.5), Inches(5.5), Inches(4), Inches(1.8))
    features_frame = features_box.text_frame
    features_frame.text = "Key Features:\n" + "\n".join(f"• {feature}" for feature in features)
    features_para = features_frame.paragraphs[0]
    features_para.font.size = Pt(12)
    features_para.font.color.rgb = theme_obj.get_color("muted.foreground")


async def test_exact_slide4():
    """Test the exact slide 4 configuration."""
    
    print("Testing exact slide 4 Bar Chart configuration...")
    
    # Create presentation with 3 slides first (to match slide position)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    theme_manager = ThemeManager()
    theme_obj = theme_manager.get_theme("dark-blue")
    
    # Add slides 1-3 first
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    theme_obj.apply_to_slide(slide1)
    slide1.shapes.title.text = "Test Showcase"
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide2)
    
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  
    theme_obj.apply_to_slide(slide3)
    
    print("✅ Added slides 1-3")
    
    try:
        # SLIDE 4: EXACT BAR CHART FROM COMPREHENSIVE SHOWCASE
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide4)
        
        add_slide_title_and_description(
            slide4,
            "03. Bar Chart (Horizontal)",
            "Horizontal orientation provides better readability for long category names and easy comparison.",
            ["Long label support", "Easy comparison", "Professional layout", "Multiple series"],
            theme_obj
        )
        
        bar_chart = BarChart(
            categories=["AI/ML", "Cloud", "Security", "IoT", "Blockchain"],
            series={
                "2024 Adoption (%)": [88, 92, 78, 65, 45],
                "2025 Projected (%)": [95, 96, 85, 75, 58]
            },
            theme=theme_obj.__dict__
        )
        await bar_chart.render(slide4, left=1.0, top=1.2, width=8.0, height=4.0)
        
        print("✅ Slide 4 Bar Chart created")
        
        # Save
        prs.save("outputs/test_exact_slide4.pptx")
        print("✅ Saved test_exact_slide4.pptx")
        
    except Exception as e:
        print(f"❌ Error on slide 4: {e}")
        import traceback
        traceback.print_exc()


async def test_simplified_slide4():
    """Test a simplified version of slide 4."""
    
    print("\nTesting simplified slide 4...")
    
    prs = Presentation()
    theme_manager = ThemeManager()
    theme_obj = theme_manager.get_theme("dark-blue")
    
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide)
        
        # Super simple bar chart
        bar_chart = BarChart(
            categories=["A", "B", "C"],
            series={"Values": [10, 20, 15]},
            theme=theme_obj.__dict__
        )
        await bar_chart.render(slide, left=2.0, top=2.0, width=6.0, height=3.0)
        
        prs.save("outputs/test_simple_slide4.pptx")
        print("✅ Saved test_simple_slide4.pptx")
        
    except Exception as e:
        print(f"❌ Error on simple slide 4: {e}")
        import traceback
        traceback.print_exc()


async def main():
    await test_exact_slide4()
    await test_simplified_slide4()


if __name__ == "__main__":
    asyncio.run(main())