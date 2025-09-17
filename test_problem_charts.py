#!/usr/bin/env python3
"""Test the specific charts causing issues in slides 4 and 8."""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.charts import BarChart, PieChart
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


async def test_slide_4_bar_chart():
    """Test the bar chart from slide 4."""
    
    print("Testing Slide 4 Bar Chart...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)
    
    try:
        # Exact bar chart from slide 4
        bar_chart = BarChart(
            categories=["Artificial Intelligence", "Cloud Computing", "Cybersecurity", "IoT Solutions", "Blockchain", "AR/VR"],
            series={
                "2024 Adoption (%)": [88, 92, 78, 65, 45, 32],
                "2025 Projected (%)": [95, 96, 85, 75, 58, 48]
            },
            variant="clustered",
            theme=theme.__dict__
        )
        
        await bar_chart.render(slide, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Bar chart created successfully")
        
        # Save
        prs.save("outputs/test_slide4_bar.pptx")
        print("✅ File saved: outputs/test_slide4_bar.pptx")
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()


async def test_slide_8_pie_chart():
    """Test the pie chart from slide 8."""
    
    print("\nTesting Slide 8 Pie Chart...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)
    
    try:
        # Exact pie chart from slide 8
        pie_chart = PieChart(
            categories=["SaaS Platform", "Enterprise Services", "Mobile Apps", "API Revenue", "Consulting"],
            values=[42, 28, 18, 8, 4],
            explode=0,  # Explode largest slice
            theme=theme.__dict__,
            options={"show_percentages": True, "show_categories": True}
        )
        
        await pie_chart.render(slide, left=2.5, top=1.2, width=5.0, height=4.0)
        print("✅ Pie chart created successfully")
        
        # Save
        prs.save("outputs/test_slide8_pie.pptx")
        print("✅ File saved: outputs/test_slide8_pie.pptx")
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()


async def main():
    """Test both problematic charts."""
    await test_slide_4_bar_chart()
    await test_slide_8_pie_chart()


if __name__ == "__main__":
    asyncio.run(main())