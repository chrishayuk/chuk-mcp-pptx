#!/usr/bin/env python3
"""
Beautiful Chart Gallery - Comprehensive showcase of all chart components.
Creates a stunning presentation with modern chart designs across multiple themes.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches

# Import chart components
from chuk_mcp_pptx.components.charts import (
    ColumnChart, BarChart, WaterfallChart,
    LineChart, AreaChart, SparklineChart,
    PieChart, DoughnutChart, SunburstChart,
    ScatterChart, BubbleChart,
    RadarChart, ComboChart
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


async def create_beautiful_chart_gallery():
    """Create a beautiful chart gallery showcasing all chart types."""
    
    print("\n📊 Creating Beautiful Chart Gallery")
    print("=" * 70)
    
    # Initialize presentation and theme manager
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    theme_manager = ThemeManager()
    
    # ==========================================================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================================================
    print("\n1. Creating title slide...")
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Beautiful Chart Gallery"
    subtitle.text = "Component-Based Charts with Modern Themes\\nSimilar to shadcn/ui but for PowerPoint"
    
    # Apply dark theme
    theme = theme_manager.get_theme("dark-violet")
    theme.apply_to_slide(slide1)
    
    # ==========================================================================
    # SLIDE 2: COLUMN CHARTS - Dark Blue Theme
    # ==========================================================================
    print("2. Creating column chart variations...")
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    theme = theme_manager.get_theme("dark-blue")
    theme.apply_to_slide(slide2)
    
    # Add title
    title_shape = slide2.shapes.title
    title_shape.text = "Column Chart Variations"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Regular column chart
    column_chart = ColumnChart(
        categories=["Q1", "Q2", "Q3", "Q4"],
        series={
            "Revenue": [120, 142, 165, 190],
            "Profit": [25, 32, 41, 48],
            "Expenses": [95, 110, 124, 142]
        },
        variant="clustered",
        title="Quarterly Performance",
        theme=theme.__dict__,
        options={"show_values": True}
    )
    await column_chart.render(slide2, left=0.5, top=2.0, width=4.0, height=3.0)
    
    # Stacked column chart
    stacked_chart = ColumnChart(
        categories=["North", "South", "East", "West"],
        series={
            "Online": [85, 92, 78, 88],
            "Retail": [65, 71, 69, 75],
            "Wholesale": [45, 38, 52, 41]
        },
        variant="stacked",
        title="Sales by Region & Channel",
        theme=theme.__dict__
    )
    await stacked_chart.render(slide2, left=5.0, top=2.0, width=4.0, height=3.0)
    
    # Waterfall chart
    waterfall = WaterfallChart(
        categories=["Start", "Sales", "Marketing", "R&D", "Costs", "Net"],
        values=[100, 45, -12, -8, -15, 110],
        title="Profit Analysis",
        theme=theme.__dict__
    )
    await waterfall.render(slide2, left=2.5, top=5.5, width=4.0, height=1.8)
    
    # ==========================================================================
    # SLIDE 3: LINE & AREA CHARTS - Cyberpunk Theme
    # ==========================================================================
    print("3. Creating line and area charts...")
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("cyberpunk")
    theme.apply_to_slide(slide3)
    
    title_shape = slide3.shapes.title
    title_shape.text = "Line & Area Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Line chart with smooth curves
    line_chart = LineChart(
        categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
        series={
            "Users": [1200, 1450, 1680, 1920, 2240, 2650],
            "Sessions": [3400, 4100, 4750, 5200, 5890, 6420],
            "Revenue": [450, 520, 680, 750, 890, 1020]
        },
        smooth=True,
        markers=True,
        title="Growth Metrics",
        theme=theme.__dict__,
        options={"show_values": False}
    )
    await line_chart.render(slide3, left=0.5, top=2.0, width=4.0, height=3.0)
    
    # Area chart with transparency
    area_chart = AreaChart(
        categories=["2019", "2020", "2021", "2022", "2023", "2024"],
        series={
            "Platform A": [45, 52, 68, 75, 82, 95],
            "Platform B": [32, 38, 45, 58, 65, 72],
            "Platform C": [18, 25, 32, 41, 48, 55]
        },
        variant="stacked",
        transparency=40,
        title="Market Share Evolution",
        theme=theme.__dict__
    )
    await area_chart.render(slide3, left=5.0, top=2.0, width=4.0, height=3.0)
    
    # Sparklines
    sparkline_data = [45, 48, 52, 49, 55, 58, 62, 60, 65, 68, 72, 75]
    sparkline = SparklineChart(
        values=sparkline_data,
        color="#00ffff",
        theme=theme.__dict__
    )
    await sparkline.render(slide3, left=2.0, top=5.8, width=2.0, height=0.5)
    
    sparkline2 = SparklineChart(
        values=[v * 1.2 for v in sparkline_data],
        color="#ff00ff",
        theme=theme.__dict__
    )
    await sparkline2.render(slide3, left=5.0, top=5.8, width=2.0, height=0.5)
    
    # ==========================================================================
    # SLIDE 4: PIE & DOUGHNUT CHARTS - Corporate Theme
    # ==========================================================================
    print("4. Creating pie and doughnut charts...")
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("corporate")
    theme.apply_to_slide(slide4)
    
    title_shape = slide4.shapes.title
    title_shape.text = "Pie & Doughnut Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Pie chart with exploded slice
    pie_chart = PieChart(
        categories=["Enterprise", "SMB", "Startups", "Government", "Education"],
        values=[42, 28, 18, 8, 4],
        explode_slice=0,  # Explode Enterprise slice
        title="Customer Segments",
        theme=theme.__dict__,
        options={
            "show_percentages": True,
            "show_categories": True,
            "labels_outside": False
        }
    )
    await pie_chart.render(slide4, left=0.5, top=2.0, width=4.0, height=3.5)
    
    # Modern doughnut chart
    doughnut = DoughnutChart(
        categories=["Cloud", "On-Premise", "Hybrid", "Edge"],
        values=[45, 30, 20, 5],
        hole_size=60,
        title="Deployment Models",
        theme=theme.__dict__,
        options={
            "show_percentages": True,
            "labels_outside": True
        }
    )
    await doughnut.render(slide4, left=5.0, top=2.0, width=4.0, height=3.5)
    
    # ==========================================================================
    # SLIDE 5: SCATTER & BUBBLE CHARTS - Gradient Theme
    # ==========================================================================
    print("5. Creating scatter and bubble charts...")
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("sunset")
    theme.apply_to_slide(slide5)
    
    title_shape = slide5.shapes.title
    title_shape.text = "Scatter & Bubble Analysis"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Scatter plot
    scatter_chart = ScatterChart(
        series_data=[
            {
                "name": "Product Performance",
                "x_values": [12, 18, 25, 32, 38, 45, 52, 58, 65, 72],
                "y_values": [15, 22, 28, 35, 42, 48, 55, 62, 68, 75]
            },
            {
                "name": "Market Trend",
                "x_values": [10, 20, 30, 40, 50, 60, 70],
                "y_values": [25, 32, 38, 45, 52, 58, 65]
            }
        ],
        show_trendline=True,
        marker_size=10,
        title="Price vs Performance",
        theme=theme.__dict__
    )
    await scatter_chart.render(slide5, left=0.5, top=2.0, width=4.0, height=3.5)
    
    # Bubble chart for 3D data
    bubble_chart = BubbleChart(
        series_data=[
            {
                "name": "Products",
                "points": [
                    [25, 35, 15],   # [price, quality, market_share]
                    [45, 55, 25],
                    [35, 45, 20],
                    [55, 65, 30],
                    [65, 75, 35],
                    [30, 40, 12],
                    [50, 60, 28]
                ]
            }
        ],
        size_scale=2.0,
        transparency=30,
        title="Market Positioning",
        theme=theme.__dict__,
        options={"show_values": False}
    )
    await bubble_chart.render(slide5, left=5.0, top=2.0, width=4.0, height=3.5)
    
    # ==========================================================================
    # SLIDE 6: RADAR & GAUGE CHARTS - Dark Green Theme
    # ==========================================================================
    print("6. Creating radar and gauge charts...")
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("dark-green")
    theme.apply_to_slide(slide6)
    
    title_shape = slide6.shapes.title
    title_shape.text = "Radar & Performance Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Radar chart for product comparison
    radar_chart = RadarChart(
        categories=["Performance", "Reliability", "Usability", "Security", "Scalability", "Cost"],
        series={
            "Product A": [8.5, 9.2, 7.8, 9.0, 8.3, 6.5],
            "Product B": [7.2, 8.8, 9.1, 8.5, 7.9, 8.2],
            "Competitor": [6.8, 7.5, 8.2, 7.8, 8.8, 9.1]
        },
        variant="filled",
        max_value=10,
        title="Product Comparison",
        theme=theme.__dict__
    )
    await radar_chart.render(slide6, left=0.5, top=2.0, width=4.0, height=3.5)
    
    # Add a combo chart instead of gauges for KPIs
    combo_chart = ComboChart(
        categories=["Q1", "Q2", "Q3", "Q4"],
        column_series={"Sales": [120, 140, 165, 190]},
        line_series={"Growth %": [8, 12, 18, 15]},
        title="Performance Metrics",
        theme=theme.__dict__
    )
    await combo_chart.render(slide6, left=5.0, top=2.0, width=4.0, height=3.5)
    
    # ==========================================================================
    # SLIDE 7: COMBO & BAR CHARTS - Light Theme
    # ==========================================================================
    print("7. Creating combo and bar charts...")
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("light")
    theme.apply_to_slide(slide7)
    
    title_shape = slide7.shapes.title
    title_shape.text = "Combo & Bar Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Horizontal bar chart
    bar_chart = BarChart(
        categories=["Machine Learning", "Cloud Computing", "Cybersecurity", "IoT", "Blockchain", "AR/VR"],
        series={
            "2023": [95, 88, 82, 75, 68, 45],
            "2024 Target": [105, 95, 90, 85, 78, 65]
        },
        variant="clustered",
        title="Technology Adoption",
        theme=theme.__dict__,
        options={"gap_width": 120}
    )
    await bar_chart.render(slide7, left=0.5, top=2.0, width=4.0, height=4.0)
    
    # Combo chart (column + line)
    combo_chart = ComboChart(
        categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
        column_series={
            "Sales": [120, 135, 148, 162, 175, 190],
            "Costs": [80, 85, 92, 98, 105, 112]
        },
        line_series={
            "Profit Margin %": [33, 37, 38, 39, 40, 41]
        },
        secondary_axis=["Profit Margin %"],
        title="Sales Performance",
        theme=theme.__dict__
    )
    await combo_chart.render(slide7, left=5.0, top=2.0, width=4.0, height=4.0)
    
    # ==========================================================================
    # SLIDE 8: ADVANCED CHARTS - Dark Purple Theme
    # ==========================================================================
    print("8. Creating advanced chart types...")
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("dark-purple")
    theme.apply_to_slide(slide8)
    
    title_shape = slide8.shapes.title
    title_shape.text = "Advanced Chart Types"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Sunburst chart for hierarchical data
    sunburst_chart = SunburstChart(
        data={
            "name": "Total",
            "value": 1000,
            "children": [
                {
                    "name": "North America",
                    "value": 400,
                    "children": [
                        {"name": "USA", "value": 300},
                        {"name": "Canada", "value": 100}
                    ]
                },
                {
                    "name": "Europe",
                    "value": 350,
                    "children": [
                        {"name": "UK", "value": 150},
                        {"name": "Germany", "value": 120},
                        {"name": "France", "value": 80}
                    ]
                },
                {
                    "name": "Asia",
                    "value": 250,
                    "children": [
                        {"name": "China", "value": 150},
                        {"name": "Japan", "value": 100}
                    ]
                }
            ]
        },
        title="Revenue Distribution",
        theme=theme.__dict__
    )
    await sunburst_chart.render(slide8, left=0.5, top=2.0, width=4.0, height=3.5)
    
    # Additional area chart for trend visualization
    area_chart = AreaChart(
        categories=["2020", "2021", "2022", "2023", "2024"],
        series={
            "North America": [250, 280, 320, 380, 450],
            "Europe": [200, 230, 260, 300, 350],
            "Asia": [150, 180, 220, 270, 330]
        },
        variant="stacked",
        title="Regional Growth Trends",
        theme=theme.__dict__,
        options={"transparency": 20}
    )
    await area_chart.render(slide8, left=5.0, top=2.0, width=4.0, height=3.5)
    
    # ==========================================================================
    # SLIDE 9: SUMMARY & FEATURES
    # ==========================================================================
    print("9. Creating summary slide...")
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("dark")
    theme.apply_to_slide(slide9)
    
    title_shape = slide9.shapes.title
    title_shape.text = "Chart Component Features"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Add feature list as text
    from chuk_mcp_pptx.components import Card
    
    features_card = Card(
        title="🎨 Design System Features",
        description="Modern chart components with shadcn/ui-like architecture",
        variant="bordered",
        theme=theme.__dict__
    )
    features_card.render(slide9, left=1.0, top=2.0, width=3.5, height=4.0)
    
    capabilities_card = Card(
        title="📊 Chart Capabilities",
        description="Comprehensive chart types with validation and themes",
        variant="elevated",
        theme=theme.__dict__
    )
    capabilities_card.render(slide9, left=5.0, top=2.0, width=3.5, height=4.0)
    
    # Add text content to cards manually
    # (In a real implementation, you'd extend Card to support rich content)
    
    # ==========================================================================
    # SAVE PRESENTATION
    # ==========================================================================
    
    print("\n10. Saving presentation...")
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "beautiful_chart_gallery.pptx")
    prs.save(output_path)
    print(f"✅ Created {output_path}")
    
    print("\n📊 Beautiful Chart Gallery Features:")
    print("  • 15+ chart types with modern styling")
    print("  • 6 different themes showcased")
    print("  • Boundary detection and validation")
    print("  • Component-based architecture")
    print("  • Theme-aware color schemes")
    print("  • Beautiful default styling")
    
    print("\n🎨 Chart Types Demonstrated:")
    print("  • Column (clustered, stacked, waterfall)")
    print("  • Line (smooth, markers, sparklines)")
    print("  • Area (filled, stacked, transparent)")
    print("  • Pie & Doughnut (exploded, modern)")
    print("  • Scatter & Bubble (correlations, 3D data)")
    print("  • Radar (multi-criteria comparison)")
    print("  • Sunburst (hierarchical data)")
    print("  • Combo (mixed chart types)")
    
    print("\n💡 This is our 'shadcn for PowerPoint' - reusable,")
    print("    beautiful components with consistent theming!")


if __name__ == "__main__":
    asyncio.run(create_beautiful_chart_gallery())