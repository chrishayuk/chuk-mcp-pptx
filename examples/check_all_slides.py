#!/usr/bin/env python3
"""
Check ALL slides in both galleries for issues.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def check_slide(slide, slide_idx, title):
    """Check a single slide for issues."""
    issues = []
    shapes = 0
    connectors = 0
    placeholders = 0
    
    for shape in slide.shapes:
        shape_type = shape.shape_type if hasattr(shape, 'shape_type') else None
        
        if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shapes += 1
        elif shape_type == MSO_SHAPE_TYPE.LINE:
            connectors += 1
            # Check for arrow heads
            has_arrows = False
            if hasattr(shape._element, 'spPr') and hasattr(shape._element.spPr, 'ln'):
                ln = shape._element.spPr.ln
                head_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                tail_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                has_arrows = head_end is not None or tail_end is not None
            
            if not has_arrows and "arrow" in title.lower():
                issues.append("Connector missing arrow heads")
                
        elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, 'placeholder_format'):
                ph_type = shape.placeholder_format.type
                if ph_type not in [1, 3]:  # Not title
                    placeholders += 1
                    if "summary" not in title.lower():
                        issues.append("Content placeholder not removed")
    
    # Check specific slide expectations
    if "arrow" in title.lower() and connectors == 0:
        issues.append("No connectors on arrows slide")
    
    if "process flow" in title.lower() and connectors < 2:
        issues.append("Too few connectors for process flow")
        
    if "hierarchy" in title.lower() or "organizational" in title.lower():
        if connectors == 0 and shapes > 2:
            issues.append("No connectors in hierarchy diagram")
    
    if "cycle" in title.lower() and shapes > 2:
        if connectors < shapes - 1:
            issues.append("Incomplete cycle - missing connectors")
    
    return {
        'shapes': shapes,
        'connectors': connectors,
        'placeholders': placeholders,
        'issues': issues
    }

def check_presentation(file_path, name):
    """Check all slides in a presentation."""
    if not os.path.exists(file_path):
        print(f"❌ {name} not found")
        return
        
    print(f"\n{'='*70}")
    print(f"📊 {name}")
    print(f"{'='*70}")
    
    prs = Presentation(file_path)
    total_issues = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx}"
        
        result = check_slide(slide, slide_idx, title)
        
        # Report
        status = "✅" if not result['issues'] else "⚠️"
        print(f"\n{status} Slide {slide_idx}: {title}")
        print(f"   Shapes: {result['shapes']}, Connectors: {result['connectors']}")
        
        if result['issues']:
            for issue in result['issues']:
                print(f"   ❌ {issue}")
                total_issues += 1
    
    print(f"\n{'='*50}")
    if total_issues == 0:
        print(f"✅ {name}: All slides OK!")
    else:
        print(f"⚠️ {name}: {total_issues} issues found")
    
    return total_issues

# Main check
print("="*70)
print("COMPREHENSIVE SLIDE CHECK")
print("="*70)

total_issues = 0

# Check Shapes Gallery
issues = check_presentation("../outputs/shapes_gallery.pptx", "SHAPES GALLERY")
total_issues += issues if issues else 0

# Check SmartArt Gallery  
issues = check_presentation("../outputs/smartart_gallery.pptx", "SMARTART GALLERY")
total_issues += issues if issues else 0

# Final report
print("\n" + "="*70)
print("FINAL REPORT")
print("="*70)

if total_issues == 0:
    print("\n✅ ALL SLIDES PASS - Everything is working correctly!")
else:
    print(f"\n⚠️ TOTAL ISSUES: {total_issues}")
    print("\nIssues need to be addressed for full functionality.")

print("\n" + "="*70)