#!/usr/bin/env python3
"""
Theme Chart Galleries - Creates individual chart galleries for each theme.
Demonstrates how the same chart components look with different themes.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches

# Import chart components
from chuk_mcp_pptx.components.charts import (
    ColumnChart, BarChart, WaterfallChart,
    LineChart, AreaChart, SparklineChart,
    PieChart, DoughnutChart,
    ScatterChart, BubbleChart,
    RadarChart, GaugeChart
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager
from chuk_mcp_pptx.components import Card, MetricCard


async def create_chart_gallery_for_theme(theme_name: str, theme_obj):
    """Create a comprehensive chart gallery for a specific theme."""
    
    print(f"\n📊 Creating {theme_name} Chart Gallery")
    print("=" * 50)
    
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==========================================================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    theme_obj.apply_to_slide(slide1)
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = f"{theme_obj.name} Chart Gallery"
    subtitle.text = f"Beautiful charts with {theme_name} theme\\nComponent-based design system"
    
    # ==========================================================================
    # SLIDE 2: COLUMN & BAR CHARTS
    # ==========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide2)
    
    title_shape = slide2.shapes.title
    title_shape.text = "Column & Bar Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Clustered column chart
    column_chart = ColumnChart(
        categories=["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024"],
        series={
            "Revenue": [125, 148, 172, 195],
            "Profit": [28, 35, 42, 48],
            "Growth": [15, 22, 28, 35]
        },
        variant="clustered",
        title="Quarterly Performance",
        theme=theme_obj.__dict__,
        options={"show_values": True}
    )
    await column_chart.render(slide2, left=0.5, top=2.0, width=4.0, height=3.0)
    
    # Horizontal bar chart
    bar_chart = BarChart(
        categories=["AI/ML", "Cloud", "Security", "IoT", "Blockchain"],
        series={
            "Adoption Rate": [85, 92, 78, 65, 45],
            "Investment": [70, 88, 82, 58, 38]
        },
        variant="clustered",
        title="Technology Trends",
        theme=theme_obj.__dict__
    )
    await bar_chart.render(slide2, left=5.0, top=2.0, width=4.0, height=3.0)
    
    # Waterfall chart
    waterfall = WaterfallChart(
        categories=["Starting", "Revenue", "Costs", "Marketing", "R&D", "Net Profit"],
        values=[100, 65, -25, -12, -18, 110],
        title="Profit Analysis",
        theme=theme_obj.__dict__
    )
    await waterfall.render(slide2, left=2.5, top=5.5, width=4.0, height=1.8)
    
    # ==========================================================================
    # SLIDE 3: LINE & AREA CHARTS
    # ==========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide3)
    
    title_shape = slide3.shapes.title
    title_shape.text = "Line & Area Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Multi-series line chart
    line_chart = LineChart(
        categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul"],
        series={
            "Website Traffic": [1200, 1350, 1480, 1620, 1850, 2100, 2350],
            "Mobile App": [800, 920, 1050, 1180, 1350, 1520, 1750],
            "API Calls": [450, 520, 680, 820, 950, 1100, 1280]
        },
        smooth=True,
        markers=True,
        title="Platform Usage Growth",
        theme=theme_obj.__dict__
    )
    await line_chart.render(slide3, left=0.5, top=2.0, width=4.0, height=3.0)
    
    # Stacked area chart
    area_chart = AreaChart(
        categories=["2020", "2021", "2022", "2023", "2024", "2025"],
        series={
            "Enterprise": [45, 52, 68, 82, 95, 110],
            "SMB": [32, 38, 45, 58, 72, 85],
            "Startup": [18, 25, 32, 41, 55, 68]
        },
        variant="stacked",
        transparency=35,
        title="Customer Segment Growth",
        theme=theme_obj.__dict__
    )
    await area_chart.render(slide3, left=5.0, top=2.0, width=4.0, height=3.0)
    
    # Sparklines for KPIs
    sparkline_data = [42, 45, 48, 52, 49, 55, 58, 62, 60, 65, 68, 72, 75, 78]
    theme_colors = theme_obj.tokens.get("chart", ["#3b82f6"])
    
    sparkline1 = SparklineChart(
        values=sparkline_data,
        color=theme_colors[0] if theme_colors else "#3b82f6",
        theme=theme_obj.__dict__
    )
    await sparkline1.render(slide3, left=1.5, top=5.8, width=2.5, height=0.6)
    
    sparkline2 = SparklineChart(
        values=[v * 1.3 for v in sparkline_data],
        color=theme_colors[1] if len(theme_colors) > 1 else "#10b981",
        theme=theme_obj.__dict__
    )
    await sparkline2.render(slide3, left=5.5, top=5.8, width=2.5, height=0.6)
    
    # ==========================================================================
    # SLIDE 4: PIE & DOUGHNUT CHARTS
    # ==========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide4)
    
    title_shape = slide4.shapes.title
    title_shape.text = "Pie & Doughnut Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Pie chart with market share
    pie_chart = PieChart(
        categories=["SaaS", "Enterprise", "Mobile", "API", "Consulting"],
        values=[38, 24, 18, 12, 8],
        explode=0,  # Explode largest slice
        title="Revenue by Product Line",
        theme=theme_obj.__dict__,
        options={
            "show_percentages": True,
            "show_categories": True
        }
    )
    await pie_chart.render(slide4, left=0.5, top=2.0, width=4.0, height=3.5)
    
    # Doughnut chart
    doughnut = DoughnutChart(
        categories=["North America", "Europe", "Asia Pacific", "Latin America", "Other"],
        values=[42, 28, 22, 6, 2],
        hole_size=55,
        title="Global Market Distribution",
        theme=theme_obj.__dict__,
        options={
            "show_percentages": True,
            "labels_outside": True
        }
    )
    await doughnut.render(slide4, left=5.0, top=2.0, width=4.0, height=3.5)
    
    # ==========================================================================
    # SLIDE 5: SCATTER & BUBBLE ANALYSIS
    # ==========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide5)
    
    title_shape = slide5.shapes.title
    title_shape.text = "Scatter & Bubble Analysis"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Scatter plot for correlation
    scatter_chart = ScatterChart(
        series_data=[
            {
                "name": "Product Performance",
                "x_values": [15, 22, 28, 35, 42, 48, 55, 62, 68, 75, 82],
                "y_values": [18, 25, 32, 38, 45, 52, 58, 65, 72, 78, 85]
            },
            {
                "name": "Competitor Data",
                "x_values": [12, 18, 25, 32, 38, 45, 52, 58, 65],
                "y_values": [22, 28, 35, 42, 48, 55, 62, 68, 75]
            }
        ],
        marker_size=12,
        title="Performance vs Customer Satisfaction",
        theme=theme_obj.__dict__
    )
    await scatter_chart.render(slide5, left=0.5, top=2.0, width=4.0, height=3.5)
    
    # Bubble chart for market positioning
    bubble_chart = BubbleChart(
        series_data=[
            {
                "name": "Product Portfolio",
                "points": [
                    [25, 35, 18],   # [price, quality, market_share]
                    [45, 55, 28],
                    [35, 45, 22],
                    [55, 65, 35],
                    [65, 75, 42],
                    [30, 40, 15],
                    [50, 60, 32],
                    [40, 50, 25]
                ]
            }
        ],
        size_scale=1.8,
        transparency=25,
        title="Market Position Analysis",
        theme=theme_obj.__dict__
    )
    await bubble_chart.render(slide5, left=5.0, top=2.0, width=4.0, height=3.5)
    
    # ==========================================================================
    # SLIDE 6: RADAR & PERFORMANCE METRICS
    # ==========================================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide6)
    
    title_shape = slide6.shapes.title
    title_shape.text = "Performance & Comparison"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Radar chart for multi-criteria comparison
    radar_chart = RadarChart(
        categories=["Performance", "Reliability", "Security", "Usability", "Scalability", "Cost"],
        series={
            "Our Product": [8.5, 9.2, 9.0, 8.8, 8.3, 7.5],
            "Competitor A": [7.2, 8.8, 8.5, 9.1, 7.9, 8.8],
            "Competitor B": [6.8, 7.5, 8.8, 7.2, 9.1, 9.2]
        },
        variant="filled",
        max_value=10,
        title="Product Comparison Matrix",
        theme=theme_obj.__dict__
    )
    await radar_chart.render(slide6, left=0.5, top=2.0, width=4.0, height=3.5)
    
    # KPI Gauge cluster
    kpi_data = [
        {"value": 87, "title": "Customer Satisfaction", "pos": (5.0, 2.0)},
        {"value": 94, "title": "System Uptime", "pos": (7.2, 2.0)},
        {"value": 76, "title": "Performance Score", "pos": (5.0, 3.8)},
        {"value": 91, "title": "Security Rating", "pos": (7.2, 3.8)}
    ]
    
    for kpi in kpi_data:
        gauge = GaugeChart(
            value=kpi["value"],
            min_value=0,
            max_value=100,
            title=kpi["title"],
            theme=theme_obj.__dict__
        )
        await gauge.render(slide6, 
                          left=kpi["pos"][0], 
                          top=kpi["pos"][1], 
                          width=1.8, 
                          height=1.4)
    
    # ==========================================================================
    # SLIDE 7: THEME SHOWCASE & METRICS
    # ==========================================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide7)
    
    title_shape = slide7.shapes.title
    title_shape.text = f"{theme_obj.name} Theme Features"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Theme info card
    theme_card = Card(
        title=f"🎨 {theme_obj.name}",
        description=f"Mode: {theme_obj.mode.title()} • Primary: {theme_obj.primary_hue.title()}",
        variant="bordered",
        theme=theme_obj.__dict__
    )
    theme_card.render(slide7, left=1.0, top=2.0, width=3.5, height=2.0)
    
    # Metrics cards showing theme colors
    metrics = [
        {"label": "Chart Types", "value": "15+", "left": 5.0},
        {"label": "Color Variants", "value": f"{len(theme_obj.tokens.get('chart', []))}", "left": 6.5},
        {"label": "Themes", "value": "9", "left": 8.0}
    ]
    
    for metric in metrics:
        metric_card = MetricCard(
            label=metric["label"],
            value=metric["value"],
            theme=theme_obj.__dict__
        )
        metric_card.render(slide7, left=metric["left"], top=2.5, width=1.2, height=1.0)
    
    # Small sample charts to show theme colors
    sample_column = ColumnChart(
        categories=["A", "B", "C"],
        series={"Series": [10, 15, 12]},
        title="Sample",
        theme=theme_obj.__dict__
    )
    await sample_column.render(slide7, left=1.5, top=4.5, width=2.5, height=1.5)
    
    sample_pie = PieChart(
        categories=["X", "Y", "Z"],
        values=[40, 35, 25],
        title="Colors",
        theme=theme_obj.__dict__
    )
    await sample_pie.render(slide7, left=5.5, top=4.5, width=2.5, height=1.5)
    
    return prs


async def create_all_theme_galleries():
    """Create chart galleries for all themes."""
    
    print("\n🎨 Creating Individual Theme Chart Galleries")
    print("=" * 70)
    
    theme_manager = ThemeManager()
    
    # Define theme groups
    theme_groups = {
        "Dark Themes": ["dark", "dark-blue", "dark-violet", "dark-green", "dark-purple"],
        "Light Themes": ["light", "corporate", "light-warm"],
        "Special Themes": ["cyberpunk", "sunset", "ocean", "minimal"]
    }
    
    created_files = []
    
    for group_name, theme_names in theme_groups.items():
        print(f"\n📁 {group_name}")
        print("-" * 30)
        
        for theme_name in theme_names:
            theme_obj = theme_manager.get_theme(theme_name)
            if theme_obj:
                print(f"  Creating {theme_name} gallery...")
                
                try:
                    prs = await create_chart_gallery_for_theme(theme_name, theme_obj)
                    
                    # Save presentation
                    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs", "theme_galleries")
                    os.makedirs(output_dir, exist_ok=True)
                    
                    filename = f"chart_gallery_{theme_name.replace('-', '_')}.pptx"
                    output_path = os.path.join(output_dir, filename)
                    prs.save(output_path)
                    
                    created_files.append((theme_name, output_path))
                    print(f"    ✅ Created {filename}")
                    
                except Exception as e:
                    print(f"    ❌ Error creating {theme_name}: {e}")
            else:
                print(f"    ⚠️  Theme '{theme_name}' not found")
    
    # Summary
    print(f"\n🎉 Created {len(created_files)} theme galleries!")
    print("\n📊 Files created:")
    for theme_name, path in created_files:
        print(f"  • {theme_name}: {os.path.basename(path)}")
    
    print(f"\n📁 All files saved to: outputs/theme_galleries/")
    print("\n💡 Each gallery shows the same charts with different themes,")
    print("    demonstrating the power of the component-based design system!")
    
    return created_files


if __name__ == "__main__":
    asyncio.run(create_all_theme_galleries())