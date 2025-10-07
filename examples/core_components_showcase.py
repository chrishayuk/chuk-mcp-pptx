#!/usr/bin/env python3
"""
Core Components Showcase - Comprehensive demonstration of Button, Badge, Alert, and Card.
Shows all variants, sizes, composition patterns, and real-world usage examples.
"""

import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core.button import Button, IconButton, ButtonGroup
from chuk_mcp_pptx.components.core.badge import Badge, DotBadge, CountBadge
from chuk_mcp_pptx.components.core.alert import Alert
from chuk_mcp_pptx.components.core.card import Card, MetricCard
from chuk_mcp_pptx.components import ProgressBar, Icon, IconList, Timeline
from chuk_mcp_pptx.components.core.tile import Tile, IconTile, ValueTile
from chuk_mcp_pptx.components.core.avatar import Avatar, AvatarWithLabel, AvatarGroup
from chuk_mcp_pptx.layout import Container, Grid, Stack
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_button_showcase(prs, theme):
    """Showcase all button variants and types."""
    print("  • Creating Button Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Button Components"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Button variants - organized with better spacing
    button_variants = [
        ("Default", "default", 0.5),
        ("Secondary", "secondary", 2.3),
        ("Outline", "outline", 4.1),
        ("Ghost", "ghost", 5.7),
        ("Destructive", "destructive", 7.0),
    ]

    for text, variant, left in button_variants:
        btn = Button(text=text, variant=variant, size="md", theme=theme.__dict__)
        btn.render(slide, left=left, top=2.0, width=1.6)

    # Button sizes
    size_buttons = [
        ("Small", "sm", 0.5, 1.3),
        ("Medium", "md", 2.2, 2.0),
        ("Large", "lg", 4.6, 2.8),
    ]

    for text, size, left, width in size_buttons:
        btn = Button(text=text, variant="default", size=size, theme=theme.__dict__)
        btn.render(slide, left=left, top=3.2, width=width)

    # Icon buttons - compact row
    icon_buttons = [
        ("play", 0.5),
        ("pause", 1.2),
        ("settings", 1.9),
        ("star", 2.6),
        ("heart", 3.3),
        ("search", 4.0),
    ]

    for icon, left in icon_buttons:
        btn = IconButton(icon=icon, variant="ghost", size="md", theme=theme.__dict__)
        btn.render(slide, left=left, top=4.5)

    # Button groups
    buttons_config = [
        {"text": "Save", "variant": "default", "size": "md"},
        {"text": "Cancel", "variant": "ghost", "size": "md"}
    ]
    group = ButtonGroup(buttons=buttons_config, orientation="horizontal", spacing=0.2, theme=theme.__dict__)
    group.render(slide, left=0.5, top=5.8)


def create_badge_showcase(prs, theme):
    """Showcase all badge variants and types."""
    print("  • Creating Badge Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Badge Components"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Badge variants - organized row
    badge_variants = [
        ("Default", "default", 0.5),
        ("Secondary", "secondary", 1.8),
        ("Success", "success", 3.1),
        ("Warning", "warning", 4.4),
        ("Destructive", "destructive", 5.7),
        ("Outline", "outline", 7.2),
    ]

    for text, variant, left in badge_variants:
        badge = Badge(text=text, variant=variant, theme=theme.__dict__)
        badge.render(slide, left=left, top=2.0)

    # Dot badges - small indicators
    dot_badges = [
        ("default", 0.5),
        ("success", 1.5),
        ("warning", 2.5),
        ("destructive", 3.5),
    ]

    for variant, left in dot_badges:
        dot = DotBadge(variant=variant, theme=theme.__dict__)
        dot.render(slide, left=left, top=3.0)

    # Count badges - notification style
    count_badges = [
        (1, 5.0),
        (12, 5.8),
        (99, 6.6),
        (150, 7.4),  # Shows "99+"
    ]

    for count, left in count_badges:
        badge = CountBadge(count=count, variant="destructive", theme=theme.__dict__)
        badge.render(slide, left=left, top=3.0)

    # Badge use cases - combined with text
    use_cases = [
        ("New", "default", 0.5, 4.2),
        ("Beta", "warning", 2.0, 4.2),
        ("Active", "success", 3.5, 4.2),
        ("Deprecated", "destructive", 5.2, 4.2),
    ]

    for text, variant, left, top in use_cases:
        badge = Badge(text=text, variant=variant, theme=theme.__dict__)
        badge.render(slide, left=left, top=top)


def create_alert_showcase(prs, theme):
    """Showcase all alert variants."""
    print("  • Creating Alert Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Alert Components"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Alert variants - stacked vertically for readability
    alerts = [
        ("info", "Information", "This is an informational message.", 2.0),
        ("success", "Success!", "Your changes have been saved successfully.", 3.2),
        ("warning", "Warning", "Please review before proceeding.", 4.4),
        ("error", "Error", "An error occurred while processing.", 5.6),
    ]

    for variant, title, description, top in alerts:
        alert = Alert(variant=variant, title=title, description=description, theme=theme.__dict__)
        alert.render(slide, left=0.5, top=top, width=9.0, height=0.9)


def create_alert_composition_showcase(prs, theme):
    """Showcase alert composition patterns."""
    print("  • Creating Alert Composition examples...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Alert Composition Patterns"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Composed alerts - using children
    alert1 = Alert(variant="success", theme=theme.__dict__)
    alert1.add_child(Alert.Title("Deployment Complete"))
    alert1.add_child(Alert.Description("Your application has been deployed to production."))
    alert1.render(slide, left=0.5, top=2.0, width=9.0, height=1.0)

    alert2 = Alert(variant="warning", show_icon=True, theme=theme.__dict__)
    alert2.add_child(Alert.Title("Quota Warning"))
    alert2.add_child(Alert.Description("You've used 90% of your monthly quota."))
    alert2.render(slide, left=0.5, top=3.5, width=9.0, height=1.0)

    # Alert without icon
    alert3 = Alert(variant="info", show_icon=False, theme=theme.__dict__)
    alert3.add_child(Alert.Title("System Maintenance"))
    alert3.add_child(Alert.Description("Scheduled maintenance on Sunday 2:00 AM - 4:00 AM."))
    alert3.render(slide, left=0.5, top=5.0, width=9.0, height=1.0)


def create_card_showcase(prs, theme):
    """Showcase all card variants."""
    print("  • Creating Card Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Card Components"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Use Container → Grid pattern for card variants
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)

    # 12-column grid for card variants (4 cols each = 3 cards)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    card_variants = [
        ("default", "Default Card", 0),
        ("outlined", "Outlined Card", 4),
        ("elevated", "Elevated Card", 8),
    ]

    for variant, title, col_start in card_variants:
        pos = grid.get_cell(col_span=4, col_start=col_start)
        card = Card(variant=variant, theme=theme.__dict__)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description("Card with composition pattern"))
        card.render(slide, **pos)

    # Metric cards using 12-column grid (3 cols each = 4 cards)
    metrics = [
        ("Revenue", "$1.2M", "+12%", "up", 0),
        ("Users", "45.2K", "+8%", "up", 3),
        ("Retention", "92%", "-2%", "down", 6),
        ("NPS", "4.8", "0%", "neutral", 9),
    ]

    for label, value, change, trend, col_start in metrics:
        pos = grid.get_cell(col_span=3, col_start=col_start, row_start=1)
        metric = MetricCard(label=label, value=value, change=change, trend=trend, theme=theme.__dict__)
        metric.render(slide, **pos)


def create_progress_icon_timeline_showcase(prs, theme):
    """Showcase ProgressBar, Icon, and Timeline components."""
    print("  • Creating ProgressBar, Icon & Timeline showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Progress, Icons & Timeline"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Progress bars
    ProgressBar(value=75, label="Project Progress", show_percentage=True, variant="success", theme=theme.__dict__).render(
        slide, left=0.5, top=1.9, width=4.0
    )

    ProgressBar(value=60, segments=10, style="segmented", label="Milestones", theme=theme.__dict__).render(
        slide, left=0.5, top=2.8, width=4.0
    )

    # Icons row
    icons = [("check", "success"), ("star", "warning"), ("rocket", "primary"), ("target", "error")]
    left = 5.0
    for icon, variant in icons:
        Icon(icon, variant=variant, size="lg", theme=theme.__dict__).render(slide, left=left, top=2.0)
        left += 0.7

    # Icon list
    features = [
        ("check", "Fast & Reliable"),
        ("check", "Easy to Use"),
        ("rocket", "High Performance")
    ]
    IconList(features, variant="success", icon_size="sm", theme=theme.__dict__).render(
        slide, left=5.0, top=2.8, width=4.0
    )

    # Timeline
    events = [
        {"date": "Q1", "title": "Plan"},
        {"date": "Q2", "title": "Build"},
        {"date": "Q3", "title": "Launch", "highlight": True}
    ]
    Timeline(events, style="arrow", theme=theme.__dict__).render(
        slide, left=0.5, top=4.5, width=8.5
    )


def create_tile_avatar_showcase(prs, theme):
    """Showcase Tile and Avatar components."""
    print("  • Creating Tile & Avatar showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Tiles & Avatars"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Tiles - different variants in a row
    tiles = [
        (IconTile("rocket", label="Fast", variant="filled", color_variant="primary"), 0.5),
        (ValueTile("42", label="Tasks", variant="outlined"), 2.6),
        (IconTile("check", label="Done", variant="filled", color_variant="success"), 4.7),
        (ValueTile("98%", label="Score", variant="default"), 6.8),
    ]

    for tile, left in tiles:
        tile.theme = theme.__dict__
        tile.render(slide, left=left, top=2.0)

    # Avatars - different sizes and variants
    avatars = [
        (Avatar(text="JD", variant="filled", color_variant="primary", size="sm"), 0.5, 4.2),
        (Avatar(text="AS", variant="outlined", color_variant="success", size="md"), 1.5, 4.0),
        (Avatar(icon="user", variant="default", size="lg"), 3.0, 3.8),
        (Avatar(text="BM", variant="filled", color_variant="warning", size="md"), 5.0, 4.0),
    ]

    for avatar, left, top in avatars:
        avatar.theme = theme.__dict__
        avatar.render(slide, left=left, top=top)

    # Avatar with label - horizontal
    avatar_label = AvatarWithLabel(
        text="JD",
        label="John Doe",
        sublabel="Product Designer",
        variant="filled",
        color_variant="primary",
        orientation="horizontal",
        theme=theme.__dict__
    )
    avatar_label.render(slide, left=0.5, top=5.5, width=3.5)

    # Avatar group
    members = [
        {"text": "JD", "color_variant": "primary"},
        {"text": "AS", "color_variant": "success"},
        {"text": "BM", "color_variant": "warning"},
        {"text": "KL", "color_variant": "destructive"},
        {"text": "MN", "color_variant": "default"}
    ]
    group = AvatarGroup(members, max_display=3, overlap=True, size="sm", theme=theme.__dict__)
    group.render(slide, left=5.0, top=5.6)


def create_combined_dashboard(prs, theme):
    """Create a realistic dashboard combining all components."""
    print("  • Creating combined components dashboard...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Dashboard Example"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Status badges at top
    Badge(text="Live", variant="success", theme=theme.__dict__).render(slide, left=8.5, top=0.3)

    # Action buttons using Stack
    buttons = [
        Button(text="Refresh", variant="outline", size="sm", theme=theme.__dict__),
        Button(text="Export", variant="ghost", size="sm", theme=theme.__dict__)
    ]
    stack = Stack(direction="horizontal", gap="sm")
    stack.render_children(slide, buttons, left=0.5, top=1.8, item_width=1.2, item_height=0.4)

    # Key metrics using Grid
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=2.3)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    metrics = [
        ("Total Sales", "$245K", "+18%", "up", 0),
        ("Conversion", "3.2%", "+0.5%", "up", 4),
        ("Bounce Rate", "42%", "-5%", "down", 8),
    ]

    for label, value, change, trend, col_start in metrics:
        pos = grid.get_cell(col_span=4, col_start=col_start)
        MetricCard(label=label, value=value, change=change, trend=trend, theme=theme.__dict__).render(slide, **pos)

    # Alert notification
    alert = Alert(variant="info", theme=theme.__dict__)
    alert.add_child(Alert.Title("New Features Available"))
    alert.add_child(Alert.Description("Check out the latest updates in the changelog."))
    alert.render(slide, left=0.5, top=4.5, width=9.0, height=0.9)

    # Status indicators - simple horizontal placement
    DotBadge(variant="success", theme=theme.__dict__).render(slide, left=0.5, top=5.8)
    DotBadge(variant="warning", theme=theme.__dict__).render(slide, left=0.8, top=5.8)
    DotBadge(variant="destructive", theme=theme.__dict__).render(slide, left=1.1, top=5.8)


def main():
    """Generate comprehensive showcase presentation."""
    print("\n🎨 Creating Core Components Showcase")
    print("=" * 70)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-violet")

    # Create showcase slides
    create_button_showcase(prs, theme)
    create_badge_showcase(prs, theme)
    create_alert_showcase(prs, theme)
    create_alert_composition_showcase(prs, theme)
    create_card_showcase(prs, theme)
    create_progress_icon_timeline_showcase(prs, theme)
    create_tile_avatar_showcase(prs, theme)
    create_combined_dashboard(prs, theme)

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "core_components_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Theme: {theme.name}")
    print("\n🎨 Showcase Features:")
    print("  • Button Components (all variants, sizes, icons, groups)")
    print("  • Badge Components (all variants, dots, counts, tags)")
    print("  • Alert Components (all variants, composition patterns)")
    print("  • Card Components (variants, composition, metrics)")
    print("  • Progress, Icons & Timeline (PowerPoint-specific components)")
    print("  • Tiles & Avatars (dashboard elements)")
    print("  • Combined Dashboard (real-world usage)")
    print("\n💡 Demonstrates:")
    print("  • Component-based architecture")
    print("  • Theme-aware styling")
    print("  • Composition patterns (shadcn-style)")
    print("  • Variant system (cva-inspired)")
    print("  • Design tokens and semantic colors")
    print("  • PowerPoint-specific components (ProgressBar, Icon, Timeline, Tile, Avatar)")


if __name__ == "__main__":
    main()
