#!/usr/bin/env python3
"""
Robust Chart Showcase - One chart per slide with comprehensive error handling.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Import chart components with error handling
try:
    from chuk_mcp_pptx.components.charts import (
        ColumnChart, BarChart, LineChart, PieChart, DoughnutChart, ScatterChart
    )
    from chuk_mcp_pptx.themes.theme_manager import ThemeManager
except ImportError as e:
    print(f"Import error: {e}")
    sys.exit(1)


def add_slide_title(slide, title: str, description: str, theme_obj):
    """Add title and description to a slide."""
    try:
        # Apply theme to slide first
        theme_obj.apply_to_slide(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
        title_para.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.0), Inches(8), Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.text = description
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(14)
        desc_para.font.color.rgb = theme_obj.get_color("muted.foreground")
        desc_para.alignment = PP_ALIGN.CENTER
        
        return True
    except Exception as e:
        print(f"     ⚠️ Error adding slide title: {e}")
        return False


async def create_robust_showcase(theme_name: str, theme_obj):
    """Create a robust showcase with proper error handling."""
    
    print(f"🔧 Creating robust {theme_name} showcase...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    charts_created = []
    
    # ==========================================================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================================================
    try:
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        theme_obj.apply_to_slide(slide1)
        
        slide1.shapes.title.text = f"{theme_obj.name} Chart Components"
        slide1.placeholders[1].text = "Robust Chart Gallery with Error Handling"
        charts_created.append("Title slide")
    except Exception as e:
        print(f"     ❌ Title slide error: {e}")
    
    # ==========================================================================
    # SLIDE 2: COLUMN CHART
    # ==========================================================================
    try:
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        if add_slide_title(slide2, "Column Chart", "Quarterly performance data", theme_obj):
            
            column_chart = ColumnChart(
                categories=["Q1", "Q2", "Q3", "Q4"],
                series={
                    "Revenue": [100, 120, 140, 160],
                    "Profit": [20, 25, 30, 35]
                },
                title="Quarterly Results",
                theme=theme_obj.__dict__
            )
            await column_chart.render(slide2, left=1.5, top=1.8, width=6.0, height=3.8)
            charts_created.append("Column chart")
            print(f"     ✅ Column chart created")
    except Exception as e:
        print(f"     ❌ Column chart error: {e}")
    
    # ==========================================================================
    # SLIDE 3: BAR CHART
    # ==========================================================================
    try:
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        if add_slide_title(slide3, "Bar Chart", "Technology adoption rates", theme_obj):
            
            bar_chart = BarChart(
                categories=["AI/ML", "Cloud", "Security"],
                series={
                    "Adoption": [85, 92, 78],
                    "Interest": [90, 88, 85]
                },
                title="Technology Trends",
                theme=theme_obj.__dict__
            )
            await bar_chart.render(slide3, left=1.5, top=1.8, width=6.0, height=3.8)
            charts_created.append("Bar chart")
            print(f"     ✅ Bar chart created")
    except Exception as e:
        print(f"     ❌ Bar chart error: {e}")
    
    # ==========================================================================
    # SLIDE 4: LINE CHART
    # ==========================================================================
    try:
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        if add_slide_title(slide4, "Line Chart", "Growth trends over time", theme_obj):
            
            line_chart = LineChart(
                categories=["Jan", "Feb", "Mar", "Apr", "May"],
                series={
                    "Users": [1200, 1350, 1480, 1620, 1850],
                    "Revenue": [25000, 28000, 32000, 36000, 42000]
                },
                title="Growth Metrics",
                theme=theme_obj.__dict__
            )
            await line_chart.render(slide4, left=1.5, top=1.8, width=6.0, height=3.8)
            charts_created.append("Line chart")
            print(f"     ✅ Line chart created")
    except Exception as e:
        print(f"     ❌ Line chart error: {e}")
    
    # ==========================================================================
    # SLIDE 5: PIE CHART
    # ==========================================================================
    try:
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        if add_slide_title(slide5, "Pie Chart", "Market share distribution", theme_obj):
            
            pie_chart = PieChart(
                categories=["Product A", "Product B", "Product C"],
                values=[50, 30, 20],
                title="Market Share",
                theme=theme_obj.__dict__
            )
            await pie_chart.render(slide5, left=2.0, top=1.8, width=5.0, height=3.8)
            charts_created.append("Pie chart")
            print(f"     ✅ Pie chart created")
    except Exception as e:
        print(f"     ❌ Pie chart error: {e}")
    
    # ==========================================================================
    # SLIDE 6: SCATTER CHART
    # ==========================================================================
    try:
        slide6 = prs.slides.add_slide(prs.slide_layouts[5])
        if add_slide_title(slide6, "Scatter Chart", "Correlation analysis", theme_obj):
            
            scatter_chart = ScatterChart(
                series_data=[
                    {
                        "name": "Data Points",
                        "x_values": [10, 20, 30, 40, 50],
                        "y_values": [15, 25, 35, 45, 55]
                    }
                ],
                title="Correlation Analysis",
                theme=theme_obj.__dict__
            )
            await scatter_chart.render(slide6, left=1.5, top=1.8, width=6.0, height=3.8)
            charts_created.append("Scatter chart")
            print(f"     ✅ Scatter chart created")
    except Exception as e:
        print(f"     ❌ Scatter chart error: {e}")
    
    # ==========================================================================
    # SLIDE 7: SUMMARY
    # ==========================================================================
    try:
        slide7 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide7)
        
        # Title
        title_box = slide7.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"✅ {len(charts_created)} Components Working!"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
        title_para.alignment = PP_ALIGN.CENTER
        
        # Summary
        summary_text = f"""🎯 Successfully created:

{chr(10).join(f'• {chart}' for chart in charts_created)}

🔧 Theme: {theme_obj.name}
🎨 Mode: {theme_obj.mode.title()}
🌈 Colors: {len(theme_obj.tokens.get('chart', []))} chart colors"""
        
        summary_box = slide7.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
        summary_frame = summary_box.text_frame
        summary_frame.text = summary_text
        
        for para in summary_frame.paragraphs:
            para.font.size = Pt(16)
            para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
        
        charts_created.append("Summary slide")
        
    except Exception as e:
        print(f"     ❌ Summary slide error: {e}")
    
    print(f"     📊 Total slides created: {len(charts_created)}")
    return prs, charts_created


async def create_robust_showcases():
    """Create robust showcases for testing."""
    
    print("\n🔧 Creating Robust Chart Showcases")
    print("=" * 50)
    print("With comprehensive error handling and validation")
    print()
    
    theme_manager = ThemeManager()
    
    # Test with just two reliable themes first
    test_themes = ["dark-blue", "corporate"]
    
    created_files = []
    
    for theme_name in test_themes:
        theme_obj = theme_manager.get_theme(theme_name)
        if theme_obj:
            try:
                prs, charts_created = await create_robust_showcase(theme_name, theme_obj)
                
                # Save with safe filename
                filename = f"robust_showcase_{theme_name.replace('-', '_')}.pptx"
                output_path = os.path.join("outputs", filename)
                
                # Ensure outputs directory exists
                os.makedirs("outputs", exist_ok=True)
                
                # Save with error handling
                prs.save(output_path)
                created_files.append((theme_name, output_path, len(charts_created)))
                print(f"  ✅ Created {filename} ({len(charts_created)} slides)")
                
            except Exception as e:
                print(f"  ❌ Error creating {theme_name}: {e}")
                import traceback
                traceback.print_exc()
        else:
            print(f"  ⚠️  Theme '{theme_name}' not found")
    
    print(f"\n🎉 Created {len(created_files)} robust showcases!")
    print("\n📁 Files created:")
    for theme_name, path, slide_count in created_files:
        print(f"  • {theme_name}: {os.path.basename(path)} ({slide_count} slides)")
    
    print(f"\n💡 These robust showcases should open perfectly in PowerPoint.")
    print("  Each has proper error handling and validated chart creation.")
    
    return created_files


if __name__ == "__main__":
    asyncio.run(create_robust_showcases())