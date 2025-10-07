#!/usr/bin/env python3
"""
Safe Chart Showcase - One chart per slide using only proven working components.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Import only the core chart components that work
from chuk_mcp_pptx.components.charts import (
    ColumnChart, BarChart, LineChart, PieChart, DoughnutChart
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def add_slide_title(slide, title: str, description: str, theme_obj):
    """Add title and description to a slide."""
    # Apply theme to slide first
    theme_obj.apply_to_slide(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    title_para.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(14)
    desc_para.font.color.rgb = theme_obj.get_color("muted.foreground")
    desc_para.alignment = PP_ALIGN.CENTER


async def create_safe_showcase(theme_name: str, theme_obj):
    """Create a safe showcase with only working chart types."""
    
    print(f"🔧 Creating safe {theme_name} showcase...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==========================================================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    theme_obj.apply_to_slide(slide1)
    
    slide1.shapes.title.text = f"{theme_obj.name} Chart Components"
    slide1.placeholders[1].text = "Proven Working Chart Gallery"
    
    # ==========================================================================
    # SLIDE 2: COLUMN CHART
    # ==========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_title(slide2, "Column Chart", "Quarterly performance with clustered columns", theme_obj)
    
    try:
        column_chart = ColumnChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series={
                "Revenue": [100, 120, 140, 160],
                "Profit": [20, 25, 30, 35]
            },
            title="Quarterly Results",
            theme=theme_obj.__dict__
        )
        await column_chart.render(slide2, left=1.5, top=1.8, width=6.0, height=3.5)
        print(f"  ✅ Column chart created")
    except Exception as e:
        print(f"  ❌ Column chart error: {e}")
    
    # ==========================================================================
    # SLIDE 3: BAR CHART
    # ==========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_title(slide3, "Bar Chart", "Technology adoption rates by category", theme_obj)
    
    try:
        bar_chart = BarChart(
            categories=["AI/ML", "Cloud", "Security", "Mobile"],
            series={
                "Adoption": [85, 92, 78, 65],
                "Interest": [90, 88, 85, 75]
            },
            title="Technology Trends",
            theme=theme_obj.__dict__
        )
        await bar_chart.render(slide3, left=1.5, top=1.8, width=6.0, height=3.5)
        print(f"  ✅ Bar chart created")
    except Exception as e:
        print(f"  ❌ Bar chart error: {e}")
    
    # ==========================================================================
    # SLIDE 4: LINE CHART
    # ==========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_title(slide4, "Line Chart", "Growth trends over time with multiple series", theme_obj)
    
    try:
        line_chart = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            series={
                "Users": [1200, 1350, 1480, 1620, 1850, 2100],
                "Revenue": [25000, 28000, 32000, 36000, 42000, 48000]
            },
            title="Growth Metrics",
            theme=theme_obj.__dict__
        )
        await line_chart.render(slide4, left=1.5, top=1.8, width=6.0, height=3.5)
        print(f"  ✅ Line chart created")
    except Exception as e:
        print(f"  ❌ Line chart error: {e}")
    
    # ==========================================================================
    # SLIDE 5: PIE CHART
    # ==========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_title(slide5, "Pie Chart", "Market share distribution by product", theme_obj)
    
    try:
        pie_chart = PieChart(
            categories=["Product A", "Product B", "Product C", "Product D"],
            values=[40, 30, 20, 10],
            title="Market Share",
            theme=theme_obj.__dict__
        )
        await pie_chart.render(slide5, left=2.0, top=1.8, width=5.0, height=3.5)
        print(f"  ✅ Pie chart created")
    except Exception as e:
        print(f"  ❌ Pie chart error: {e}")
    
    # ==========================================================================
    # SLIDE 6: DOUGHNUT CHART
    # ==========================================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_slide_title(slide6, "Doughnut Chart", "Revenue breakdown with center space for metrics", theme_obj)
    
    try:
        doughnut_chart = DoughnutChart(
            categories=["SaaS", "Enterprise", "Mobile", "API"],
            values=[45, 25, 20, 10],
            title="Revenue Sources",
            theme=theme_obj.__dict__
        )
        await doughnut_chart.render(slide6, left=2.0, top=1.8, width=5.0, height=3.5)
        print(f"  ✅ Doughnut chart created")
    except Exception as e:
        print(f"  ❌ Doughnut chart error: {e}")
    
    # ==========================================================================
    # SLIDE 7: SUMMARY
    # ==========================================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide7)
    
    # Title
    title_box = slide7.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "✅ All Core Components Working!"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    title_para.alignment = PP_ALIGN.CENTER
    
    # Summary
    summary_box = slide7.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
    summary_frame = summary_box.text_frame
    summary_frame.text = """🎯 Component Architecture Proven:

• Column Charts - Clustered data visualization
• Bar Charts - Horizontal category comparisons  
• Line Charts - Trend analysis over time
• Pie Charts - Proportion visualization
• Doughnut Charts - Modern proportion display

🔧 Theme Integration:
• Consistent color schemes
• Beautiful typography
• Professional styling
• Component reusability"""
    
    for para in summary_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    return prs


async def create_safe_showcases():
    """Create safe showcases for core themes."""
    
    print("\n🔧 Creating Safe Chart Showcases")
    print("=" * 50)
    print("Using only proven working chart components")
    print()
    
    theme_manager = ThemeManager()
    
    # Use only the most reliable themes
    safe_themes = ["dark-blue", "corporate", "cyberpunk", "minimal"]
    
    created_files = []
    
    for theme_name in safe_themes:
        theme_obj = theme_manager.get_theme(theme_name)
        if theme_obj:
            try:
                prs = await create_safe_showcase(theme_name, theme_obj)
                
                # Save with safe filename
                filename = f"safe_showcase_{theme_name.replace('-', '_')}.pptx"
                output_path = os.path.join("outputs", filename)
                
                # Ensure outputs directory exists
                os.makedirs("outputs", exist_ok=True)
                
                # Save with error handling
                prs.save(output_path)
                created_files.append((theme_name, output_path))
                print(f"  ✅ Created {filename}")
                
            except Exception as e:
                print(f"  ❌ Error creating {theme_name}: {e}")
        else:
            print(f"  ⚠️  Theme '{theme_name}' not found")
    
    print(f"\n🎉 Created {len(created_files)} safe showcases!")
    print("\n📁 Files created:")
    for theme_name, path in created_files:
        print(f"  • {theme_name}: {os.path.basename(path)}")
    
    print(f"\n💡 These safe showcases should open without any PowerPoint errors.")
    print("  Each showcase has 7 slides with core chart types that we know work.")
    
    return created_files


if __name__ == "__main__":
    asyncio.run(create_safe_showcases())