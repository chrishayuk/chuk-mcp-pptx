#!/usr/bin/env python3
"""Direct inspection of gallery files."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def inspect_presentation(file_path, name):
    """Directly inspect a presentation file."""
    if not os.path.exists(file_path):
        print(f"❌ File not found: {file_path}")
        return
        
    print(f"\n{'='*70}")
    print(f"🔍 INSPECTING {name}")
    print(f"{'='*70}")
    
    prs = Presentation(file_path)
    print(f"Total slides: {len(prs.slides)}")
    
    for slide_idx, slide in enumerate(prs.slides):
        # Skip title slide
        if slide_idx == 0:
            continue
            
        print(f"\n{'='*50}")
        print(f"📌 Slide {slide_idx}: {slide.shapes.title.text if slide.shapes.title else 'No title'}")
        print(f"{'='*50}")
        
        # Count different shape types
        autoshapes = []
        connectors = []
        text_boxes = []
        placeholders = []
        
        for shape in slide.shapes:
            shape_type = shape.shape_type if hasattr(shape, 'shape_type') else None
            
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Check if it has text
                text = ""
                if shape.has_text_frame and shape.text_frame.text:
                    text = shape.text_frame.text[:20]
                    # Check text positioning
                    tf = shape.text_frame
                    margin_info = f"margins(L:{tf.margin_left/914400:.1f} R:{tf.margin_right/914400:.1f} T:{tf.margin_top/914400:.1f} B:{tf.margin_bottom/914400:.1f})"
                else:
                    margin_info = ""
                    
                autoshapes.append(f"  • AutoShape: {text if text else '(no text)'} {margin_info}")
                
            elif shape_type == MSO_SHAPE_TYPE.LINE:
                connectors.append(f"  • Connector/Line")
                
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text = shape.text_frame.text[:30] if shape.has_text_frame else ""
                text_boxes.append(f"  • TextBox: '{text}'")
                
            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if not (hasattr(shape, 'placeholder_format') and 
                       shape.placeholder_format.type in [1, 3]):  # Not title placeholders
                    placeholders.append(f"  • Placeholder (should be removed)")
        
        # Report findings
        if autoshapes:
            print(f"\n🔷 AutoShapes ({len(autoshapes)}):")
            for s in autoshapes:
                print(s)
                
        if connectors:
            print(f"\n🔗 Connectors ({len(connectors)}):")
            for c in connectors:
                print(c)
                
        if text_boxes:
            print(f"\n📝 Text Boxes ({len(text_boxes)}):")
            for t in text_boxes:
                print(t)
                
        if placeholders:
            print(f"\n⚠️  ISSUE: Content placeholders not removed ({len(placeholders)}):")
            for p in placeholders:
                print(p)
        
        # Check for potential text issues
        issues = []
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                # Check if margins are too small
                if tf.margin_left < 72000:  # Less than 0.08 inches
                    issues.append(f"  • Shape has very small left margin")
                if tf.margin_top < 36000:  # Less than 0.04 inches  
                    issues.append(f"  • Shape has very small top margin")
                    
        if issues:
            print(f"\n⚠️  POTENTIAL TEXT ISSUES:")
            for issue in issues:
                print(issue)

# Inspect both galleries
inspect_presentation("../outputs/shapes_gallery.pptx", "SHAPES GALLERY")
inspect_presentation("../outputs/smartart_gallery.pptx", "SMARTART GALLERY")

print("\n" + "="*70)
print("✅ Direct inspection complete!")