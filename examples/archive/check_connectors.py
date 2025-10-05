#!/usr/bin/env python3
"""
Detailed inspection of connectors in gallery presentations.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def inspect_connectors_detailed(file_path, gallery_name):
    """Inspect connectors in detail."""
    if not os.path.exists(file_path):
        print(f"❌ File not found: {file_path}")
        return
    
    print(f"\n{'='*80}")
    print(f"🔍 CONNECTOR INSPECTION: {gallery_name}")
    print(f"{'='*80}")
    
    prs = Presentation(file_path)
    
    connector_slides = []
    
    for slide_idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx}"
        
        # Count connectors
        connectors = []
        shapes = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type if hasattr(shape, 'shape_type') else None
            
            if shape_type == MSO_SHAPE_TYPE.LINE:
                # Found a connector
                connector_info = {
                    'index': shape_idx,
                    'type': 'LINE/CONNECTOR'
                }
                
                # Get positions
                if hasattr(shape, 'begin_x'):
                    connector_info['start'] = (shape.begin_x / 914400, shape.begin_y / 914400)
                    connector_info['end'] = (shape.end_x / 914400, shape.end_y / 914400)
                elif hasattr(shape, 'left'):
                    left = shape.left / 914400
                    top = shape.top / 914400
                    width = shape.width / 914400
                    height = shape.height / 914400
                    connector_info['position'] = (left, top, width, height)
                
                # Check line properties
                if hasattr(shape, 'line'):
                    line = shape.line
                    connector_info['line_width'] = line.width.pt if line.width else 0
                    
                    # Try to detect arrow heads (this may not work directly)
                    connector_info['has_line'] = True
                    
                    # Check XML for arrow heads
                    if hasattr(shape._element, 'spPr') and hasattr(shape._element.spPr, 'ln'):
                        ln = shape._element.spPr.ln
                        # Check for headEnd and tailEnd elements
                        head_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                        tail_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                        connector_info['arrow_end'] = head_end is not None
                        connector_info['arrow_start'] = tail_end is not None
                
                connectors.append(connector_info)
            
            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if shape.has_text_frame and shape.text_frame.text:
                    left = shape.left / 914400
                    top = shape.top / 914400
                    width = shape.width / 914400
                    height = shape.height / 914400
                    shapes.append({
                        'text': shape.text_frame.text[:20],
                        'bounds': (left, top, left + width, top + height)
                    })
        
        if connectors:
            connector_slides.append((slide_idx, title, connectors, shapes))
    
    # Report findings
    if not connector_slides:
        print("\n⚠️ No connectors found in any slides!")
        return
    
    for slide_idx, title, connectors, shapes in connector_slides:
        print(f"\n{'='*60}")
        print(f"📌 Slide {slide_idx}: {title}")
        print(f"{'='*60}")
        
        print(f"\n📊 Summary:")
        print(f"  • Shapes with text: {len(shapes)}")
        print(f"  • Connectors: {len(connectors)}")
        
        if shapes:
            print(f"\n🔷 Shapes:")
            for shape in shapes:
                print(f"  • '{shape['text']}' at ({shape['bounds'][0]:.1f}, {shape['bounds'][1]:.1f})")
        
        if connectors:
            print(f"\n🔗 Connectors:")
            for conn in connectors:
                print(f"\n  [{conn['index']}] Connector:")
                
                if 'start' in conn:
                    print(f"    Start: ({conn['start'][0]:.2f}, {conn['start'][1]:.2f})")
                    print(f"    End: ({conn['end'][0]:.2f}, {conn['end'][1]:.2f})")
                    
                    # Calculate length
                    import math
                    dx = conn['end'][0] - conn['start'][0]
                    dy = conn['end'][1] - conn['start'][1]
                    length = math.sqrt(dx*dx + dy*dy)
                    print(f"    Length: {length:.2f} inches")
                elif 'position' in conn:
                    pos = conn['position']
                    print(f"    Position: ({pos[0]:.2f}, {pos[1]:.2f})")
                    print(f"    Size: {pos[2]:.2f} x {pos[3]:.2f}")
                
                if 'line_width' in conn:
                    print(f"    Line width: {conn['line_width']:.1f} pt")
                
                # Arrow information
                arrow_info = []
                if conn.get('arrow_start'):
                    arrow_info.append("start")
                if conn.get('arrow_end'):
                    arrow_info.append("end")
                
                if arrow_info:
                    print(f"    ✓ Arrows at: {', '.join(arrow_info)}")
                else:
                    print(f"    ⚠️ No arrows detected")
                
                # Check if connector connects shapes
                if 'start' in conn and shapes:
                    connections = []
                    for shape in shapes:
                        # Check if start point is near shape
                        if (shape['bounds'][0] - 0.2 <= conn['start'][0] <= shape['bounds'][2] + 0.2 and
                            shape['bounds'][1] - 0.2 <= conn['start'][1] <= shape['bounds'][3] + 0.2):
                            connections.append(f"FROM '{shape['text']}'")
                        # Check if end point is near shape
                        if (shape['bounds'][0] - 0.2 <= conn['end'][0] <= shape['bounds'][2] + 0.2 and
                            shape['bounds'][1] - 0.2 <= conn['end'][1] <= shape['bounds'][3] + 0.2):
                            connections.append(f"TO '{shape['text']}'")
                    
                    if connections:
                        print(f"    🔗 Connects: {' '.join(connections)}")

# Main inspection
print("="*80)
print("COMPREHENSIVE CONNECTOR INSPECTION")
print("="*80)

# Check Shapes Gallery
print("\n" + "="*80)
print("1. SHAPES GALLERY")
print("="*80)
inspect_connectors_detailed("../outputs/shapes_gallery.pptx", "SHAPES GALLERY")

# Check SmartArt Gallery
print("\n" + "="*80)
print("2. SMARTART GALLERY")
print("="*80)
inspect_connectors_detailed("../outputs/smartart_gallery.pptx", "SMARTART GALLERY")

# Check connector verification
print("\n" + "="*80)
print("3. CONNECTOR VERIFICATION TEST")
print("="*80)
inspect_connectors_detailed("../outputs/connector_verify.pptx", "CONNECTOR VERIFY")

print("\n" + "="*80)
print("✅ INSPECTION COMPLETE")
print("="*80)