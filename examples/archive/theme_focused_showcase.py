#!/usr/bin/env python3
"""
Theme-Focused Showcase
Creates one comprehensive presentation per theme showing ALL components and chart types.
Perfect for seeing how a consistent theme looks across all elements.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches
from typing import Dict, List, Any

# Import all components
from chuk_mcp_pptx.components.charts import (
    ColumnChart, BarChart, WaterfallChart,
    LineChart, AreaChart, SparklineChart,
    PieChart, DoughnutChart, SunburstChart,
    ScatterChart, BubbleChart,
    RadarChart, ComboChart, FunnelChart
)
from chuk_mcp_pptx.components import Card, Button
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


class ThemeFocusedShowcase:
    """Create a comprehensive showcase for each theme."""
    
    def __init__(self):
        self.theme_manager = ThemeManager()
        self.themes = [
            "dark", "dark-blue", "dark-violet", "dark-green", "dark-purple",
            "light", "corporate", "light-warm"
        ]
    
    async def create_theme_showcase(self, theme_name: str) -> Presentation:
        """Create a complete showcase for a single theme."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        theme = self.theme_manager.get_theme(theme_name)
        
        # ======================================================================
        # SLIDE 1: Title Slide
        # ======================================================================
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        theme.apply_to_slide(title_slide)
        
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None
        
        title.text = f"{theme_name.replace('-', ' ').title()} Theme"
        if subtitle:
            subtitle.text = "Complete Component & Chart Showcase\nShadcn-inspired Design System for PowerPoint"
        
        # ======================================================================
        # SLIDE 2: Column & Bar Charts
        # ======================================================================
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide2)
        
        slide2.shapes.title.text = "Column & Bar Chart Variations"
        
        # Clustered Column
        column = ColumnChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series={
                "Revenue": [120, 135, 155, 180],
                "Profit": [20, 28, 35, 42],
                "Costs": [100, 107, 120, 138]
            },
            variant="clustered",
            title="Quarterly Performance",
            theme=theme.__dict__
        )
        await column.render(slide2, left=0.5, top=1.5, width=3.0, height=2.5)
        
        # Stacked Bar
        bar = BarChart(
            categories=["Product A", "Product B", "Product C", "Product D"],
            series={
                "Online": [45, 52, 38, 42],
                "Retail": [35, 28, 42, 38],
                "Wholesale": [20, 20, 20, 20]
            },
            variant="stacked",
            title="Sales by Channel",
            theme=theme.__dict__
        )
        await bar.render(slide2, left=3.75, top=1.5, width=3.0, height=2.5)
        
        # Waterfall
        waterfall = WaterfallChart(
            categories=["Start", "Q1", "Q2", "Q3", "Q4", "Total"],
            values=[100, 20, 15, -10, 25, 150],
            title="Revenue Bridge",
            theme=theme.__dict__
        )
        await waterfall.render(slide2, left=7.0, top=1.5, width=2.75, height=2.5)
        
        # 100% Stacked Column
        stacked100 = ColumnChart(
            categories=["2020", "2021", "2022", "2023"],
            series={
                "Americas": [30, 35, 40, 45],
                "Europe": [40, 35, 35, 30],
                "Asia": [30, 30, 25, 25]
            },
            variant="stacked",
            title="Market Share by Region (%)",
            theme=theme.__dict__
        )
        await stacked100.render(slide2, left=0.5, top=4.5, width=4.5, height=2.5)
        
        # Funnel Chart
        funnel = FunnelChart(
            stages=["Leads", "Qualified", "Proposal", "Negotiation", "Closed"],
            values=[10000, 7500, 3500, 1500, 500],
            show_percentages=True,
            show_values=True,
            title="Sales Funnel",
            theme=theme.__dict__
        )
        await funnel.render(slide2, left=5.25, top=4.5, width=4.5, height=2.5)
        
        # ======================================================================
        # SLIDE 3: Line & Area Charts
        # ======================================================================
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide3)
        
        slide3.shapes.title.text = "Line & Area Chart Variations"
        
        # Smooth Line
        line = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            series={
                "2023": [65, 68, 72, 78, 85, 92],
                "2024": [72, 75, 80, 88, 95, 105],
                "Target": [70, 75, 80, 85, 90, 95]
            },
            variant="smooth",
            title="Performance Trends",
            theme=theme.__dict__
        )
        await line.render(slide3, left=0.5, top=1.5, width=3.0, height=2.5)
        
        # Line with Markers
        line_markers = LineChart(
            categories=["Mon", "Tue", "Wed", "Thu", "Fri"],
            series={
                "Week 1": [82, 85, 88, 86, 90],
                "Week 2": [78, 82, 85, 88, 92],
                "Week 3": [85, 87, 90, 92, 95]
            },
            variant="markers",
            title="Daily Metrics",
            theme=theme.__dict__
        )
        await line_markers.render(slide3, left=3.75, top=1.5, width=3.0, height=2.5)
        
        # Area Chart
        area = AreaChart(
            categories=["2019", "2020", "2021", "2022", "2023"],
            series={
                "Revenue": [200, 250, 320, 380, 450],
                "Costs": [150, 180, 210, 240, 280],
                "Profit": [50, 70, 110, 140, 170]
            },
            variant="stacked",
            title="Financial Growth",
            theme=theme.__dict__
        )
        await area.render(slide3, left=7.0, top=1.5, width=2.75, height=2.5)
        
        # Sparklines
        sparkline1 = SparklineChart(
            values=[10, 12, 8, 14, 16, 12, 18, 20],
            title="CPU Usage",
            theme=theme.__dict__
        )
        await sparkline1.render(slide3, left=0.5, top=4.5, width=2.0, height=1.2)
        
        sparkline2 = SparklineChart(
            values=[45, 48, 42, 50, 55, 52, 58, 60],
            title="Memory",
            theme=theme.__dict__
        )
        await sparkline2.render(slide3, left=2.75, top=4.5, width=2.0, height=1.2)
        
        sparkline3 = SparklineChart(
            values=[100, 95, 98, 88, 85, 90, 92, 95],
            title="Network",
            theme=theme.__dict__
        )
        await sparkline3.render(slide3, left=5.0, top=4.5, width=2.0, height=1.2)
        
        # Combined Area
        combined_area = AreaChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series={
                "Baseline": [100, 100, 100, 100],
                "Growth": [20, 35, 45, 60],
                "Bonus": [5, 10, 15, 20]
            },
            variant="filled",
            title="Cumulative Performance",
            theme=theme.__dict__
        )
        await combined_area.render(slide3, left=7.25, top=4.5, width=2.5, height=2.5)
        
        # ======================================================================
        # SLIDE 4: Pie, Doughnut & Circular Charts
        # ======================================================================
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide4)
        
        slide4.shapes.title.text = "Circular Chart Variations"
        
        # Basic Pie
        pie = PieChart(
            categories=["Product A", "Product B", "Product C", "Product D", "Other"],
            values=[35, 25, 20, 15, 5],
            variant="standard",
            title="Market Share",
            theme=theme.__dict__
        )
        await pie.render(slide4, left=0.5, top=1.5, width=3.0, height=3.0)
        
        # Doughnut
        doughnut = DoughnutChart(
            categories=["Completed", "In Progress", "Pending", "Blocked"],
            values=[65, 20, 10, 5],
            variant="modern",
            title="Project Status",
            theme=theme.__dict__
        )
        await doughnut.render(slide4, left=3.75, top=1.5, width=3.0, height=3.0)
        
        # Sunburst
        sunburst_data = {
            "name": "Total",
            "value": 100,
            "children": [
                {
                    "name": "Division A",
                    "value": 45,
                    "children": [
                        {"name": "Team 1", "value": 20},
                        {"name": "Team 2", "value": 25}
                    ]
                },
                {
                    "name": "Division B",
                    "value": 35,
                    "children": [
                        {"name": "Team 3", "value": 20},
                        {"name": "Team 4", "value": 15}
                    ]
                },
                {
                    "name": "Division C",
                    "value": 20
                }
            ]
        }
        sunburst = SunburstChart(
            data=sunburst_data,
            title="Organization Structure",
            theme=theme.__dict__
        )
        await sunburst.render(slide4, left=7.0, top=1.5, width=2.75, height=3.0)
        
        # Exploded Pie
        exploded_pie = PieChart(
            categories=["Critical", "High", "Medium", "Low"],
            values=[5, 15, 35, 45],
            variant="exploded",
            title="Risk Distribution",
            theme=theme.__dict__
        )
        await exploded_pie.render(slide4, left=2.0, top=5.0, width=3.0, height=2.0)
        
        # Nested Doughnut
        nested_doughnut = DoughnutChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            values=[22, 28, 25, 25],
            variant="nested",
            title="Quarterly Split",
            theme=theme.__dict__
        )
        await nested_doughnut.render(slide4, left=5.5, top=5.0, width=3.0, height=2.0)
        
        # ======================================================================
        # SLIDE 5: Scatter, Bubble & Advanced Charts
        # ======================================================================
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide5)
        
        slide5.shapes.title.text = "Advanced Analytics Charts"
        
        # Scatter Plot
        scatter_data = [
            {
                "name": "Product Line A",
                "x_values": [2, 3, 4, 5, 6],
                "y_values": [15, 18, 22, 28, 35]
            },
            {
                "name": "Product Line B",
                "x_values": [2, 3, 4, 5, 6],
                "y_values": [12, 14, 17, 20, 24]
            },
            {
                "name": "Product Line C",
                "x_values": [2, 3, 4, 5, 6],
                "y_values": [8, 10, 11, 13, 15]
            }
        ]
        scatter = ScatterChart(
            series_data=scatter_data,
            show_trendline=True,
            marker_size=10,
            title="Price vs Performance",
            theme=theme.__dict__
        )
        await scatter.render(slide5, left=0.5, top=1.5, width=3.0, height=2.5)
        
        # Bubble Chart
        bubble_data = [
            {
                "name": "Emerging Markets",
                "points": [[25, 45, 15], [30, 55, 25], [35, 65, 35]]  # [x, y, size]
            },
            {
                "name": "Developed Markets",
                "points": [[45, 75, 45], [50, 80, 50], [55, 85, 55]]
            },
            {
                "name": "Frontier Markets",
                "points": [[15, 25, 10], [20, 30, 15], [25, 35, 20]]
            }
        ]
        bubble = BubbleChart(
            series_data=bubble_data,
            size_scale=2.0,
            transparency=30,
            title="Market Analysis",
            theme=theme.__dict__
        )
        await bubble.render(slide5, left=3.75, top=1.5, width=3.0, height=2.5)
        
        # Radar Chart
        radar = RadarChart(
            categories=["Speed", "Reliability", "Cost", "Features", "Support", "Security"],
            series={
                "Product A": [8, 9, 6, 8, 7, 9],
                "Product B": [7, 8, 8, 9, 8, 7],
                "Competitor": [6, 7, 9, 7, 6, 8]
            },
            variant="filled",
            max_value=10,
            title="Competitive Analysis",
            theme=theme.__dict__
        )
        await radar.render(slide5, left=7.0, top=1.5, width=2.75, height=2.5)
        
        # Combo Chart
        combo = ComboChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            column_series={
                "Sales": [120, 135, 125, 145, 160, 175],
                "Costs": [80, 85, 82, 90, 95, 100]
            },
            line_series={
                "Margin %": [33, 37, 34, 38, 41, 43],
                "Target %": [35, 35, 35, 35, 35, 35]
            },
            title="Sales & Margins",
            theme=theme.__dict__
        )
        await combo.render(slide5, left=0.5, top=4.5, width=4.5, height=2.5)
        
        # Advanced Bubble
        advanced_bubble = BubbleChart(
            series_data=[
                {
                    "name": "High Growth",
                    "points": [[70, 85, 40], [75, 90, 45], [80, 95, 50]]
                },
                {
                    "name": "Stable",
                    "points": [[50, 60, 30], [55, 65, 35], [60, 70, 38]]
                }
            ],
            size_scale=1.5,
            title="Portfolio Distribution",
            theme=theme.__dict__
        )
        await advanced_bubble.render(slide5, left=5.25, top=4.5, width=4.5, height=2.5)
        
        # ======================================================================
        # SLIDE 6: UI Components
        # ======================================================================
        slide6 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide6)
        
        slide6.shapes.title.text = "UI Components & Elements"
        
        # Cards
        card1 = Card(
            title="Revenue",
            description="Total revenue for Q4 2024: $12.5M",
            variant="elevated",
            theme=theme.__dict__
        )
        card1.render(slide6, left=0.5, top=1.5, width=2.5, height=2.0)
        
        card2 = Card(
            title="Growth Rate",
            description="Year over year growth: +28%",
            variant="bordered",
            theme=theme.__dict__
        )
        card2.render(slide6, left=3.25, top=1.5, width=2.5, height=2.0)
        
        card3 = Card(
            title="Customer Score",
            description="Average satisfaction rating: 4.8/5.0",
            variant="default",
            theme=theme.__dict__
        )
        card3.render(slide6, left=6.0, top=1.5, width=2.5, height=2.0)
        
        # Additional Cards as badges replacement
        card4 = Card(
            title="Status",
            description="✅ Active",
            variant="elevated",
            theme=theme.__dict__
        )
        card4.render(slide6, left=0.5, top=4.0, width=2.0, height=1.0)
        
        card5 = Card(
            title="Alert",
            description="⚠️ Warning",
            variant="bordered",
            theme=theme.__dict__
        )
        card5.render(slide6, left=3.0, top=4.0, width=2.0, height=1.0)
        
        # Buttons
        button_y = 5.0
        for i, (text, variant) in enumerate([
            ("Primary", "primary"),
            ("Secondary", "secondary"),
            ("Outline", "outline"),
            ("Ghost", "ghost")
        ]):
            button = Button(
                text=text,
                variant=variant,
                theme=theme.__dict__
            )
            button.render(slide6, left=0.5 + (i * 2.2), top=button_y, width=2.0, height=0.6)
        
        # More status cards
        card6 = Card(
            title="Progress",
            description="75% Complete",
            variant="default",
            theme=theme.__dict__
        )
        card6.render(slide6, left=5.5, top=4.0, width=2.0, height=1.0)
        
        card7 = Card(
            title="Score",
            description="Grade: A+",
            variant="elevated",
            theme=theme.__dict__
        )
        card7.render(slide6, left=8.0, top=4.0, width=1.75, height=1.0)
        
        # ======================================================================
        # SLIDE 7: Data Visualization Summary
        # ======================================================================
        slide7 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide7)
        
        slide7.shapes.title.text = "Data Visualization Grid"
        
        # Create a grid of small charts
        charts_grid = [
            ("Revenue", ColumnChart(
                categories=["Q1", "Q2", "Q3", "Q4"],
                series={"Revenue": [100, 120, 140, 160]},
                variant="clustered",
                theme=theme.__dict__
            )),
            ("Expenses", BarChart(
                categories=["Salaries", "Marketing", "R&D", "Other"],
                series={"Amount": [45, 20, 25, 10]},
                variant="standard",
                theme=theme.__dict__
            )),
            ("Trends", LineChart(
                categories=["Jan", "Feb", "Mar", "Apr"],
                series={"Sales": [80, 85, 90, 95]},
                variant="smooth",
                theme=theme.__dict__
            )),
            ("Distribution", PieChart(
                categories=["A", "B", "C", "D"],
                values=[30, 25, 25, 20],
                variant="standard",
                theme=theme.__dict__
            )),
            ("Performance", RadarChart(
                categories=["Q1", "Q2", "Q3", "Q4"],
                series={"Score": [7, 8, 9, 8]},
                max_value=10,
                variant="filled",
                theme=theme.__dict__
            )),
            ("Pipeline", FunnelChart(
                stages=["Lead", "Qualified", "Closed"],
                values=[1000, 500, 100],
                show_percentages=True,
                theme=theme.__dict__
            ))
        ]
        
        # Arrange in 3x2 grid
        for i, (title, chart) in enumerate(charts_grid):
            row = i // 3
            col = i % 3
            x = 0.5 + (col * 3.2)
            y = 1.5 + (row * 3.0)
            
            chart.title = title
            await chart.render(slide7, left=x, top=y, width=3.0, height=2.5)
        
        # ======================================================================
        # SLIDE 8: Theme Summary
        # ======================================================================
        slide8 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide8)
        
        slide8.shapes.title.text = f"{theme_name.replace('-', ' ').title()} Theme Summary"
        
        # Summary cards showing theme characteristics
        summary_card1 = Card(
            title="🎨 Design Philosophy",
            description=self._get_theme_philosophy(theme_name),
            variant="elevated",
            theme=theme.__dict__
        )
        summary_card1.render(slide8, left=0.5, top=1.5, width=4.5, height=2.5)
        
        summary_card2 = Card(
            title="📊 Best Use Cases",
            description=self._get_theme_use_cases(theme_name),
            variant="bordered",
            theme=theme.__dict__
        )
        summary_card2.render(slide8, left=5.25, top=1.5, width=4.5, height=2.5)
        
        summary_card3 = Card(
            title="✨ Key Features",
            description=self._get_theme_features(theme_name),
            variant="default",
            theme=theme.__dict__
        )
        summary_card3.render(slide8, left=2.5, top=4.5, width=5.0, height=2.5)
        
        return prs
    
    def _get_theme_philosophy(self, theme_name: str) -> str:
        """Get design philosophy for theme."""
        philosophies = {
            "dark": "Modern, sleek design with high contrast for maximum readability",
            "dark-blue": "Professional dark theme with calming blue accents",
            "dark-violet": "Creative and bold with purple highlights",
            "dark-green": "Nature-inspired dark theme with growth-oriented green tones",
            "dark-purple": "Sophisticated and elegant with royal purple accents",
            "light": "Clean, minimal design perfect for formal presentations",
            "corporate": "Traditional business styling with professional appeal",
            "light-warm": "Friendly and approachable with warm, inviting colors"
        }
        return philosophies.get(theme_name, "Professional design system")
    
    def _get_theme_use_cases(self, theme_name: str) -> str:
        """Get best use cases for theme."""
        use_cases = {
            "dark": "Tech demos, developer presentations, evening events",
            "dark-blue": "Financial reports, corporate updates, analytics dashboards",
            "dark-violet": "Creative pitches, design reviews, innovation workshops",
            "dark-green": "Sustainability reports, growth metrics, environmental data",
            "dark-purple": "Executive briefings, luxury brands, premium products",
            "light": "Board meetings, academic presentations, formal reports",
            "corporate": "Quarterly reviews, investor relations, business proposals",
            "light-warm": "Team meetings, training sessions, internal communications"
        }
        return use_cases.get(theme_name, "Business presentations")
    
    def _get_theme_features(self, theme_name: str) -> str:
        """Get key features for theme."""
        if "dark" in theme_name:
            return "• High contrast ratios\n• Reduced eye strain in low light\n• Modern, sophisticated appearance\n• Excellent for data visualization"
        else:
            return "• Maximum readability\n• Professional appearance\n• Print-friendly colors\n• Traditional business styling"
    
    async def generate_all_theme_showcases(self):
        """Generate comprehensive showcase for each theme."""
        output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs", "theme_showcases")
        os.makedirs(output_dir, exist_ok=True)
        
        print("\n🎨 Generating Theme-Focused Showcases")
        print("=" * 70)
        print(f"Creating {len(self.themes)} comprehensive theme presentations")
        print("Each presentation shows ALL components and charts in that theme")
        print()
        
        for theme_name in self.themes:
            print(f"📊 Creating {theme_name} showcase...")
            
            # Create presentation for this theme
            prs = await self.create_theme_showcase(theme_name)
            
            # Save presentation
            filename = f"theme_showcase_{theme_name.replace('-', '_')}.pptx"
            filepath = os.path.join(output_dir, filename)
            prs.save(filepath)
            
            print(f"   ✅ Saved: {filename}")
            print(f"      • 8 slides of comprehensive components")
            print(f"      • 30+ chart variations")
            print(f"      • UI components and badges")
            print(f"      • Complete design system demonstration")
        
        print("\n" + "=" * 70)
        print("✨ All theme showcases generated successfully!")
        print(f"📁 Output directory: {output_dir}")
        
        return output_dir


async def main():
    """Generate all theme-focused showcases."""
    showcase = ThemeFocusedShowcase()
    output_dir = await showcase.generate_all_theme_showcases()
    
    print("\n🎯 Theme Showcases Complete!")
    print("Each file demonstrates the complete design system in a single theme:")
    print("  • All chart types (Column, Bar, Line, Area, Pie, etc.)")
    print("  • UI Components (Cards, Badges, Buttons)")
    print("  • Data visualization grids")
    print("  • Theme-specific styling and colors")
    print("\n💡 Perfect for:")
    print("  • Choosing a theme for your presentation")
    print("  • Seeing consistency across all components")
    print("  • Understanding the design system capabilities")


if __name__ == "__main__":
    asyncio.run(main())