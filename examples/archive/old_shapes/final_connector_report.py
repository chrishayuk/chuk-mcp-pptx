#!/usr/bin/env python3
"""
Final comprehensive connector report for all galleries.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def analyze_connectors(file_path):
    """Analyze connector quality in a presentation."""
    if not os.path.exists(file_path):
        return None
    
    prs = Presentation(file_path)
    
    stats = {
        'total_connectors': 0,
        'with_arrows': 0,
        'bidirectional': 0,
        'no_arrows': 0,
        'properly_connected': 0,
        'connector_types': set()
    }
    
    for slide in prs.slides:
        shapes_positions = {}
        
        # First pass: collect shape positions
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if shape.has_text_frame and shape.text_frame.text:
                    left = shape.left / 914400
                    top = shape.top / 914400
                    width = shape.width / 914400
                    height = shape.height / 914400
                    shapes_positions[shape.text_frame.text[:20]] = (left, top, left + width, top + height)
        
        # Second pass: analyze connectors
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.LINE:
                stats['total_connectors'] += 1
                
                # Check for arrows
                has_start = False
                has_end = False
                
                if hasattr(shape._element, 'spPr') and hasattr(shape._element.spPr, 'ln'):
                    ln = shape._element.spPr.ln
                    head_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                    tail_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                    
                    has_end = head_end is not None
                    has_start = tail_end is not None
                
                if has_start and has_end:
                    stats['bidirectional'] += 1
                elif has_start or has_end:
                    stats['with_arrows'] += 1
                else:
                    stats['no_arrows'] += 1
                
                # Check if properly connected to shapes
                if hasattr(shape, 'begin_x') and shapes_positions:
                    start_x = shape.begin_x / 914400
                    start_y = shape.begin_y / 914400
                    end_x = shape.end_x / 914400
                    end_y = shape.end_y / 914400
                    
                    connected = False
                    for text, bounds in shapes_positions.items():
                        # Check if connector touches shape boundaries
                        start_near = (bounds[0] - 0.3 <= start_x <= bounds[2] + 0.3 and 
                                     bounds[1] - 0.3 <= start_y <= bounds[3] + 0.3)
                        end_near = (bounds[0] - 0.3 <= end_x <= bounds[2] + 0.3 and 
                                   bounds[1] - 0.3 <= end_y <= bounds[3] + 0.3)
                        
                        if start_near and end_near:
                            connected = True
                            break
                    
                    if connected:
                        stats['properly_connected'] += 1
    
    return stats

# Generate final report
print("="*80)
print("FINAL CONNECTOR QUALITY REPORT")
print("="*80)

galleries = [
    ("../outputs/shapes_gallery.pptx", "Shapes Gallery"),
    ("../outputs/smartart_gallery.pptx", "SmartArt Gallery"),
    ("../outputs/connector_verify.pptx", "Connector Verification"),
    ("../outputs/chart_gallery.pptx", "Chart Gallery"),
    ("../outputs/table_gallery.pptx", "Table Gallery")
]

total_stats = {
    'total_connectors': 0,
    'with_arrows': 0,
    'bidirectional': 0,
    'no_arrows': 0,
    'properly_connected': 0
}

for file_path, name in galleries:
    stats = analyze_connectors(file_path)
    if stats:
        print(f"\n{'='*60}")
        print(f"📊 {name}")
        print(f"{'='*60}")
        
        print(f"\n  Total Connectors: {stats['total_connectors']}")
        
        if stats['total_connectors'] > 0:
            arrow_pct = ((stats['with_arrows'] + stats['bidirectional']) / stats['total_connectors']) * 100
            connected_pct = (stats['properly_connected'] / stats['total_connectors']) * 100
            
            print(f"  ✓ With arrows: {stats['with_arrows']}")
            print(f"  ↔ Bidirectional: {stats['bidirectional']}")
            print(f"  ✗ No arrows: {stats['no_arrows']}")
            print(f"  🔗 Properly connected: {stats['properly_connected']}")
            
            print(f"\n  Quality Metrics:")
            print(f"    • Arrow coverage: {arrow_pct:.1f}%")
            print(f"    • Connection accuracy: {connected_pct:.1f}%")
            
            # Quality assessment
            if arrow_pct >= 90 and connected_pct >= 80:
                print(f"    • Overall: ✅ EXCELLENT")
            elif arrow_pct >= 70 and connected_pct >= 60:
                print(f"    • Overall: ✓ GOOD")
            else:
                print(f"    • Overall: ⚠️ NEEDS IMPROVEMENT")
        
        # Update totals
        for key in total_stats:
            total_stats[key] += stats.get(key, 0)

# Overall summary
print(f"\n{'='*80}")
print(f"📈 OVERALL CONNECTOR QUALITY")
print(f"{'='*80}")

if total_stats['total_connectors'] > 0:
    overall_arrow_pct = ((total_stats['with_arrows'] + total_stats['bidirectional']) / total_stats['total_connectors']) * 100
    overall_connected_pct = (total_stats['properly_connected'] / total_stats['total_connectors']) * 100
    
    print(f"\n  Total Connectors Across All Galleries: {total_stats['total_connectors']}")
    print(f"  Connectors with Arrows: {total_stats['with_arrows'] + total_stats['bidirectional']}/{total_stats['total_connectors']} ({overall_arrow_pct:.1f}%)")
    print(f"  Properly Connected: {total_stats['properly_connected']}/{total_stats['total_connectors']} ({overall_connected_pct:.1f}%)")
    
    print(f"\n  🏆 FINAL ASSESSMENT:")
    if overall_arrow_pct >= 90 and overall_connected_pct >= 80:
        print(f"     ✅ EXCELLENT - Connectors are well-implemented!")
    elif overall_arrow_pct >= 70:
        print(f"     ✓ GOOD - Most connectors working properly")
    else:
        print(f"     ⚠️ NEEDS WORK - Some connector issues remain")

print("\n" + "="*80)
print("✅ CONNECTOR QUALITY REPORT COMPLETE")
print("="*80)