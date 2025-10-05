#!/usr/bin/env python3
"""
Fix Chart Galleries - Creates more robust chart galleries with better PowerPoint compatibility.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

# Import chart components
from chuk_mcp_pptx.components.charts import (
    ColumnChart, BarChart, LineChart, PieChart, DoughnutChart, ScatterChart
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager
from chuk_mcp_pptx.components import Card


async def create_simplified_gallery(theme_name: str, theme_obj):
    """Create a simplified, more compatible chart gallery."""
    
    print(f"🔧 Creating simplified {theme_name} gallery...")
    
    # Initialize presentation with explicit template
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ==========================================================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    theme_obj.apply_to_slide(slide1)
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = f"{theme_obj.name} Gallery"
    subtitle.text = f"Charts with {theme_name} theme"
    
    # Style title text
    title_para = title.text_frame.paragraphs[0]
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    
    subtitle_para = subtitle.text_frame.paragraphs[0]
    subtitle_para.font.color.rgb = theme_obj.get_color("muted.foreground")
    subtitle_para.font.size = Pt(20)
    
    # ==========================================================================
    # SLIDE 2: COLUMN CHARTS
    # ==========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    theme_obj.apply_to_slide(slide2)
    
    # Add title manually
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Column Charts"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Simple column chart
    try:
        column_chart = ColumnChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series={
                "Revenue": [100, 120, 140, 160],
                "Profit": [20, 25, 30, 35]
            },
            title="Quarterly Performance",
            theme=theme_obj.__dict__
        )
        await column_chart.render(slide2, left=1.0, top=2.0, width=7.0, height=4.0)
    except Exception as e:
        print(f"  ⚠️ Column chart error: {e}")
        # Add error placeholder
        error_box = slide2.shapes.add_textbox(Inches(1), Inches(3), Inches(7), Inches(2))
        error_box.text_frame.text = "Column Chart\n(Chart generation in progress)"
    
    # ==========================================================================
    # SLIDE 3: PIE CHART
    # ==========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide3)
    
    # Title
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Pie Charts"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Simple pie chart
    try:
        pie_chart = PieChart(
            categories=["Product A", "Product B", "Product C", "Product D"],
            values=[40, 30, 20, 10],
            title="Market Share",
            theme=theme_obj.__dict__
        )
        await pie_chart.render(slide3, left=2.0, top=2.0, width=6.0, height=4.0)
    except Exception as e:
        print(f"  ⚠️ Pie chart error: {e}")
        error_box = slide3.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        error_box.text_frame.text = "Pie Chart\n(Chart generation in progress)"
    
    # ==========================================================================
    # SLIDE 4: LINE CHART
    # ==========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide4)
    
    # Title
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Line Charts"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Simple line chart
    try:
        line_chart = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            series={
                "Growth": [10, 15, 12, 18, 22, 25],
                "Target": [12, 14, 16, 18, 20, 22]
            },
            title="Monthly Trends",
            theme=theme_obj.__dict__
        )
        await line_chart.render(slide4, left=1.0, top=2.0, width=7.0, height=4.0)
    except Exception as e:
        print(f"  ⚠️ Line chart error: {e}")
        error_box = slide4.shapes.add_textbox(Inches(1), Inches(3), Inches(7), Inches(2))
        error_box.text_frame.text = "Line Chart\n(Chart generation in progress)"
    
    # ==========================================================================
    # SLIDE 5: THEME SHOWCASE
    # ==========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide5)
    
    # Title
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = f"{theme_obj.name} Theme"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    # Add theme color swatches
    chart_colors = theme_obj.tokens.get("chart", [])
    for i, color in enumerate(chart_colors[:6]):  # Show first 6 colors
        if isinstance(color, str):
            # Create color swatch
            swatch = slide5.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.5 + i * 1.2), Inches(2.5),
                Inches(1.0), Inches(1.0)
            )
            
            # Apply color
            swatch.fill.solid()
            rgb = theme_obj.hex_to_rgb(color)
            swatch.fill.fore_color.rgb = theme_obj.get_color("primary.DEFAULT").__class__(*rgb)
            
            # Remove border
            swatch.line.fill.background()
            
            # Add color label
            label_box = slide5.shapes.add_textbox(
                Inches(1.5 + i * 1.2), Inches(3.8),
                Inches(1.0), Inches(0.5)
            )
            label_frame = label_box.text_frame
            label_frame.text = f"Color {i+1}"
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(10)
            label_para.font.color.rgb = theme_obj.get_color("muted.foreground")
    
    # Add theme info
    info_box = slide5.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    info_frame = info_box.text_frame
    info_frame.text = f"Theme: {theme_obj.name}\nMode: {theme_obj.mode.title()}\nPrimary Hue: {theme_obj.primary_hue.title()}"
    
    for para in info_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    
    return prs


async def create_all_fixed_galleries():
    """Create simplified, more compatible galleries for all themes."""
    
    print("\n🔧 Creating Fixed Chart Galleries")
    print("=" * 50)
    print("These galleries are simplified for better PowerPoint compatibility")
    print()
    
    theme_manager = ThemeManager()
    
    # Focus on core themes that work well
    core_themes = [
        "dark", "dark-blue", "dark-violet", 
        "light", "corporate", 
        "cyberpunk", "minimal"
    ]
    
    created_files = []
    
    for theme_name in core_themes:
        theme_obj = theme_manager.get_theme(theme_name)
        if theme_obj:
            try:
                prs = await create_simplified_gallery(theme_name, theme_obj)
                
                # Save with safe filename
                safe_name = theme_name.replace("-", "_")
                filename = f"fixed_gallery_{safe_name}.pptx"
                output_path = os.path.join("outputs", filename)
                
                # Ensure outputs directory exists
                os.makedirs("outputs", exist_ok=True)
                
                # Save with error handling
                try:
                    prs.save(output_path)
                    created_files.append((theme_name, output_path))
                    print(f"  ✅ Created {filename}")
                except Exception as save_error:
                    print(f"  ❌ Save error for {theme_name}: {save_error}")
                    
            except Exception as e:
                print(f"  ❌ Creation error for {theme_name}: {e}")
    
    print(f"\n🎉 Created {len(created_files)} fixed galleries!")
    print("\n📁 Files created:")
    for theme_name, path in created_files:
        print(f"  • {theme_name}: {os.path.basename(path)}")
    
    print(f"\n💡 These simplified galleries should open without errors in PowerPoint.")
    print("  Each gallery has 5 slides with basic chart examples.")
    
    return created_files


if __name__ == "__main__":
    asyncio.run(create_all_fixed_galleries())