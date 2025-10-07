#!/usr/bin/env python3
"""
Single Theme Demo - Creates a simple, guaranteed-to-work chart demo.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt

from chuk_mcp_pptx.components.charts import ColumnChart, PieChart, LineChart
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


async def create_simple_demo():
    """Create a simple, robust chart demo."""
    
    print("\n🎯 Creating Simple Chart Demo")
    print("=" * 40)
    
    # Use dark-blue theme (very reliable)
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    theme.apply_to_slide(slide1)
    
    slide1.shapes.title.text = "Chart Component Demo"
    slide1.placeholders[1].text = "Beautiful, Theme-Aware Charts"
    
    # SLIDE 2: Column Chart
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide2)
    
    # Add title
    title_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Column Chart Example"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Simple column chart
    column_chart = ColumnChart(
        categories=["Q1", "Q2", "Q3", "Q4"],
        series={
            "Sales": [100, 120, 140, 160],
            "Profit": [20, 25, 30, 35]
        },
        title="Quarterly Results",
        theme=theme.__dict__
    )
    await column_chart.render(slide2, left=1.5, top=2.0, width=6.0, height=4.0)
    
    # SLIDE 3: Success message
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide3)
    
    # Success message
    msg_box = slide3.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(3))
    msg_frame = msg_box.text_frame
    msg_frame.text = "✅ Success!\n\nYour chart component system\nis working perfectly.\n\nThis demonstrates:\n• Theme integration\n• Component architecture\n• Beautiful defaults"
    
    for para in msg_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Save
    output_path = "outputs/simple_demo.pptx"
    os.makedirs("outputs", exist_ok=True)
    prs.save(output_path)
    
    print(f"✅ Created {output_path}")
    print("🎯 This simple demo should open perfectly in PowerPoint!")
    
    return output_path


if __name__ == "__main__":
    asyncio.run(create_simple_demo())