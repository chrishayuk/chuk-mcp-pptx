#!/usr/bin/env python3
"""Detailed inspection of connectors in test file."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_CONNECTOR
import os

def inspect_connectors(file_path):
    """Inspect connector details."""
    prs = Presentation(file_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx}"
        print(f"\n{'='*60}")
        print(f"Slide {slide_idx}: {title}")
        print(f"{'='*60}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type if hasattr(shape, 'shape_type') else None
            
            if shape_type == MSO_SHAPE_TYPE.LINE:
                print(f"\n[{shape_idx}] CONNECTOR/LINE:")
                
                # Get connector type if available
                if hasattr(shape, 'connector_type'):
                    print(f"   Type: {shape.connector_type}")
                
                # Get position info
                if hasattr(shape, 'begin_x'):
                    begin_x = shape.begin_x / 914400
                    begin_y = shape.begin_y / 914400
                    end_x = shape.end_x / 914400
                    end_y = shape.end_y / 914400
                    print(f"   Start: ({begin_x:.2f}, {begin_y:.2f})")
                    print(f"   End: ({end_x:.2f}, {end_y:.2f})")
                elif hasattr(shape, 'left'):
                    left = shape.left / 914400
                    top = shape.top / 914400
                    width = shape.width / 914400
                    height = shape.height / 914400
                    print(f"   Position: ({left:.2f}, {top:.2f})")
                    print(f"   Size: {width:.2f} x {height:.2f}")
                
                # Check line properties
                if hasattr(shape, 'line'):
                    line = shape.line
                    if line.width:
                        print(f"   Line width: {line.width.pt:.1f} pt")
                    if hasattr(line, 'color') and line.color.rgb:
                        color = line.color.rgb
                        print(f"   Line color: RGB({color[0]}, {color[1]}, {color[2]})")
                    
                    # Check for arrowheads
                    if hasattr(line, 'begin_arrow_type'):
                        print(f"   Begin arrow: {line.begin_arrow_type}")
                    if hasattr(line, 'end_arrow_type'):
                        print(f"   End arrow: {line.end_arrow_type}")
            
            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if shape.has_text_frame and shape.text_frame.text:
                    left = shape.left / 914400
                    top = shape.top / 914400
                    width = shape.width / 914400
                    height = shape.height / 914400
                    print(f"\n[{shape_idx}] SHAPE: {shape.text_frame.text}")
                    print(f"   Position: ({left:.2f}, {top:.2f})")
                    print(f"   Size: {width:.2f} x {height:.2f}")
                    print(f"   Edges: L={left:.2f}, R={left+width:.2f}, T={top:.2f}, B={top+height:.2f}")

if __name__ == "__main__":
    print("CONNECTOR INSPECTION REPORT")
    print("="*60)
    
    test_file = "../outputs/connector_test.pptx"
    if os.path.exists(test_file):
        inspect_connectors(test_file)
    else:
        print(f"❌ File not found: {test_file}")
    
    print("\n" + "="*60)
    print("CHECKING GALLERIES")
    print("="*60)
    
    # Check shapes gallery for connectors
    shapes_file = "../outputs/shapes_gallery.pptx"
    if os.path.exists(shapes_file):
        print("\n📐 SHAPES GALLERY - Slide 2 (Arrows and Connectors):")
        prs = Presentation(shapes_file)
        if len(prs.slides) > 2:
            slide = prs.slides[2]
            connector_count = 0
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    connector_count += 1
                    if hasattr(shape, 'begin_x'):
                        begin_x = shape.begin_x / 914400
                        begin_y = shape.begin_y / 914400
                        end_x = shape.end_x / 914400
                        end_y = shape.end_y / 914400
                        print(f"   Connector: ({begin_x:.1f}, {begin_y:.1f}) → ({end_x:.1f}, {end_y:.1f})")
            print(f"   Total connectors found: {connector_count}")