#!/usr/bin/env python3
"""
Comprehensive inspection of all shapes and connectors in galleries.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def inspect_shapes_and_connectors(file_path, name):
    """Comprehensively inspect all shapes and connectors."""
    if not os.path.exists(file_path):
        print(f"❌ File not found: {file_path}")
        return
        
    print(f"\n{'='*80}")
    print(f"🔍 COMPREHENSIVE INSPECTION: {name}")
    print(f"{'='*80}")
    
    prs = Presentation(file_path)
    
    total_shapes = 0
    total_connectors = 0
    total_placeholders = 0
    issues = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_idx}"
        
        print(f"\n{'='*60}")
        print(f"📌 Slide {slide_idx}: {slide_title}")
        print(f"{'='*60}")
        
        # Categorize shapes
        autoshapes = []
        connectors = []
        text_boxes = []
        placeholders = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type if hasattr(shape, 'shape_type') else None
            
            # Get position and size
            if hasattr(shape, 'left'):
                left = shape.left / 914400  # EMU to inches
                top = shape.top / 914400
                width = shape.width / 914400 if hasattr(shape, 'width') else 0
                height = shape.height / 914400 if hasattr(shape, 'height') else 0
                
                # Check for out-of-bounds
                if left < 0 or top < 0:
                    issues.append(f"Slide {slide_idx}: Shape at negative position ({left:.1f}, {top:.1f})")
                if left + width > 10:
                    issues.append(f"Slide {slide_idx}: Shape extends beyond slide width")
                if top + height > 7.5:
                    issues.append(f"Slide {slide_idx}: Shape extends beyond slide height")
            else:
                left = top = width = height = 0
            
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Get text and shape details
                text = ""
                if shape.has_text_frame and shape.text_frame.text:
                    text = shape.text_frame.text[:30]
                    # Check text alignment
                    tf = shape.text_frame
                    if tf.margin_left < 0 or tf.margin_top < 0:
                        issues.append(f"Slide {slide_idx}: Shape '{text}' has negative text margins")
                
                # Get auto shape type if available
                auto_type = shape.auto_shape_type if hasattr(shape, 'auto_shape_type') else "Unknown"
                
                autoshapes.append(f"  [{shape_idx}] AutoShape ({auto_type}): '{text}' at ({left:.1f}, {top:.1f}) size {width:.1f}x{height:.1f}")
                total_shapes += 1
                
            elif shape_type == MSO_SHAPE_TYPE.LINE:
                # Connector details
                begin_x = shape.begin_x / 914400 if hasattr(shape, 'begin_x') else left
                begin_y = shape.begin_y / 914400 if hasattr(shape, 'begin_y') else top
                end_x = shape.end_x / 914400 if hasattr(shape, 'end_x') else left + width
                end_y = shape.end_y / 914400 if hasattr(shape, 'end_y') else top + height
                
                connectors.append(f"  [{shape_idx}] Connector: ({begin_x:.1f}, {begin_y:.1f}) → ({end_x:.1f}, {end_y:.1f})")
                total_connectors += 1
                
                # Check if connector is too short
                import math
                length = math.sqrt((end_x - begin_x)**2 + (end_y - begin_y)**2)
                if length < 0.1:
                    issues.append(f"Slide {slide_idx}: Very short connector (length {length:.2f})")
                
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text = shape.text_frame.text[:30] if shape.has_text_frame else "(empty)"
                text_boxes.append(f"  [{shape_idx}] TextBox: '{text}' at ({left:.1f}, {top:.1f})")
                
            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                ph_type = "unknown"
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    # Check for content placeholders that should be removed
                    if ph_type not in [1, 3]:  # Not title placeholders
                        placeholders.append(f"  [{shape_idx}] ⚠️ Content Placeholder (type {ph_type}) - should be removed")
                        total_placeholders += 1
                        issues.append(f"Slide {slide_idx}: Content placeholder not removed")
                    else:
                        placeholders.append(f"  [{shape_idx}] Title Placeholder")
        
        # Report findings
        if autoshapes:
            print(f"\n🔷 SHAPES ({len(autoshapes)}):")
            for s in autoshapes:
                print(s)
        
        if connectors:
            print(f"\n🔗 CONNECTORS ({len(connectors)}):")
            for c in connectors:
                print(c)
        
        if text_boxes:
            print(f"\n📝 TEXT BOXES ({len(text_boxes)}):")
            for t in text_boxes:
                print(t)
        
        if placeholders:
            print(f"\n📍 PLACEHOLDERS ({len(placeholders)}):")
            for p in placeholders:
                print(p)
        
        # Check for overlapping shapes
        shapes_with_bounds = []
        for shape in slide.shapes:
            if hasattr(shape, 'left') and shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX]:
                shapes_with_bounds.append({
                    'left': shape.left / 914400,
                    'top': shape.top / 914400,
                    'right': (shape.left + shape.width) / 914400,
                    'bottom': (shape.top + shape.height) / 914400,
                    'text': shape.text_frame.text[:20] if shape.has_text_frame else ""
                })
        
        # Check for overlaps
        overlaps = []
        for i in range(len(shapes_with_bounds)):
            for j in range(i + 1, len(shapes_with_bounds)):
                s1 = shapes_with_bounds[i]
                s2 = shapes_with_bounds[j]
                
                # Check if rectangles overlap
                if not (s1['right'] < s2['left'] or s2['right'] < s1['left'] or
                       s1['bottom'] < s2['top'] or s2['bottom'] < s1['top']):
                    overlaps.append(f"  ⚠️ '{s1['text']}' overlaps with '{s2['text']}'")
        
        if overlaps:
            print(f"\n⚠️ OVERLAPPING SHAPES:")
            for o in overlaps:
                print(o)
                issues.append(f"Slide {slide_idx}: {o}")
    
    # Summary
    print(f"\n{'='*60}")
    print(f"📊 SUMMARY")
    print(f"{'='*60}")
    print(f"Total Shapes: {total_shapes}")
    print(f"Total Connectors: {total_connectors}")
    print(f"Total Unwanted Placeholders: {total_placeholders}")
    
    if issues:
        print(f"\n⚠️ ISSUES FOUND ({len(issues)}):")
        for issue in issues:
            print(f"  • {issue}")
    else:
        print("\n✅ No issues found!")
    
    return issues

# Run comprehensive inspection
print("="*80)
print("COMPREHENSIVE SHAPE AND CONNECTOR INSPECTION")
print("="*80)

shapes_issues = inspect_shapes_and_connectors("../outputs/shapes_gallery.pptx", "SHAPES GALLERY")
smartart_issues = inspect_shapes_and_connectors("../outputs/smartart_gallery.pptx", "SMARTART GALLERY")

print("\n" + "="*80)
print("FINAL REPORT")
print("="*80)

if shapes_issues or smartart_issues:
    print("\n⚠️ Issues requiring attention:")
    if shapes_issues:
        print(f"\nShapes Gallery: {len(shapes_issues)} issues")
    if smartart_issues:
        print(f"\nSmartArt Gallery: {len(smartart_issues)} issues")
else:
    print("\n✅ All shapes and connectors are properly positioned and rendered!")