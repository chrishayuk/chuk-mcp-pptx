"""
Demonstration of enhanced component system with variants, composition, and registry.
Shows the shadcn-inspired improvements to the PowerPoint design system.
"""

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.card_v2 import Card, MetricCard
from chuk_mcp_pptx.composition import (
    CardHeader, CardContent, CardFooter, CardTitle, CardDescription,
    CompositionBuilder, Badge, Separator, compose
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager
from chuk_mcp_pptx.registry import registry
from chuk_mcp_pptx.variants import create_variants, CARD_VARIANTS

import json


def demo_variants():
    """Demonstrate the variant system."""
    print("=== Variant System Demo ===\n")

    # Show card variants
    print("Card Variants Schema:")
    schema = CARD_VARIANTS.get_schema()
    print(json.dumps(schema, indent=2))
    print()

    # Build different variant combinations
    print("Building variant combinations:")

    default_card = CARD_VARIANTS.build(variant="default", padding="md")
    print(f"Default + Medium: {default_card}")

    outlined_large = CARD_VARIANTS.build(variant="outlined", padding="lg")
    print(f"Outlined + Large: {outlined_large}")

    elevated_small = CARD_VARIANTS.build(variant="elevated", padding="sm")
    print(f"Elevated + Small: {elevated_small}")
    print()

    # Create custom variant
    custom_variants = create_variants(
        base={"font_family": "Inter"},
        variants={
            "style": {
                "modern": {"bg_color": "primary.DEFAULT", "rounded": True},
                "classic": {"bg_color": "secondary.DEFAULT", "rounded": False},
            },
            "emphasis": {
                "high": {"font_weight": "bold", "opacity": 1.0},
                "low": {"font_weight": "normal", "opacity": 0.7},
            }
        },
        default_variants={"style": "modern", "emphasis": "high"},
        compound_variants=[
            {
                "conditions": {"style": "modern", "emphasis": "high"},
                "props": {"border_glow": True}
            }
        ]
    )

    print("Custom variant combination:")
    custom_props = custom_variants.build(style="modern", emphasis="high")
    print(json.dumps(custom_props, indent=2))
    print("\n" + "="*50 + "\n")


def demo_composition():
    """Demonstrate composition patterns."""
    print("=== Composition Patterns Demo ===\n")

    # Method 1: Manual composition
    print("Method 1: Manual Composition")
    print("---")
    print("""
card = Card(variant="outlined")
card.add_child(CardHeader("Dashboard", "Overview of key metrics"))
card.add_child(CardContent("View your performance data"))
card.add_child(CardFooter("Updated 5 minutes ago"))
    """)
    print()

    # Method 2: Using class attributes (shadcn style)
    print("Method 2: Class Attribute Style (shadcn-inspired)")
    print("---")
    print("""
card = Card(variant="elevated")
card.add_child(Card.Header("Product Features", "What we offer"))
card.add_child(Card.Content("Feature description here"))
card.add_child(Card.Footer("Learn more →", align="right"))
    """)
    print()

    # Method 3: Composition builder
    print("Method 3: Fluent Builder Pattern")
    print("---")
    print("""
builder = CompositionBuilder(theme)
children = (builder
    .header("Analytics", "Real-time insights")
    .separator()
    .content("Your metrics are trending upward")
    .badge("New", "success")
    .footer("View details")
    .build())

card = Card(variant="default")
for child in children:
    card.add_child(child)
    """)
    print()

    # Method 4: Compose helper
    print("Method 4: Compose Helper Function")
    print("---")
    print("""
from chuk_mcp_pptx.composition import compose, with_separator

card = Card(variant="outlined")
card._children = compose(
    CardTitle("Welcome"),
    CardDescription("Get started with our platform"),
    CardContent("Follow the steps below")
)

# With automatic separators
card._children = with_separator(
    CardTitle("Section 1"),
    CardContent("Content 1"),
    CardTitle("Section 2"),
    CardContent("Content 2")
)
    """)
    print("\n" + "="*50 + "\n")


def demo_registry():
    """Demonstrate component registry."""
    print("=== Component Registry Demo ===\n")

    # List all components
    print(f"Registered Components: {registry.list_components()}\n")

    # Get component schema
    if "Card" in registry.list_components():
        print("Card Component Schema:")
        schema = registry.get_schema("Card")
        print(json.dumps(schema, indent=2))
        print()

        # Show component signature
        sig = registry.get_component_signature("Card")
        print(f"Component Signature: {sig}\n")

        # Show variants
        variants = registry.list_variants("Card")
        print(f"Available Variants: {json.dumps(variants, indent=2)}\n")

        # Show examples
        examples = registry.get_examples("Card")
        print("Usage Examples:")
        for i, ex in enumerate(examples, 1):
            print(f"\nExample {i}: {ex.get('description', 'No description')}")
            print(ex.get('code', 'No code'))

    # Search functionality
    print("\n" + "="*50)
    print("\nSearch Results for 'metric':")
    results = registry.search("metric")
    for r in results:
        print(f"  - {r.name}: {r.description}")

    # Export for LLM
    print("\n" + "="*50)
    print("\nLLM Documentation Export (first 500 chars):")
    llm_docs = registry.export_for_llm()
    print(llm_docs[:500] + "...")
    print("\n" + "="*50 + "\n")


def create_showcase_presentation():
    """Create a presentation showcasing the new features."""
    print("=== Creating Showcase Presentation ===\n")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Setup theme
    theme_mgr = ThemeManager()
    theme = theme_mgr.get_theme("dark-violet")
    theme_dict = {
        "name": theme.name,
        "mode": theme.mode,
        "primary_hue": theme.primary_hue,
        "font_family": theme.font_family
    }

    # Slide 1: Variant Showcase
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    theme.apply_to_slide(slide1)

    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Variant System Showcase"
    title_frame.paragraphs[0].font.size = Inches(0.4)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Different card variants
    variants_to_show = ["default", "outlined", "elevated", "ghost"]
    x_positions = [0.5, 3.0, 5.5, 8.0]

    for variant, x in zip(variants_to_show, x_positions):
        card = Card(variant=variant, padding="sm", theme=theme_dict)
        card.add_child(CardTitle(variant.capitalize()))
        card.add_child(CardDescription(f"{variant} variant"))
        card.render(slide1, left=x, top=1.5, width=2.0, height=1.5)

    # Slide 2: Composition Showcase
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    theme.apply_to_slide(slide2)

    # Title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Composition Patterns"
    title_frame.paragraphs[0].font.size = Inches(0.4)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Card with full composition
    card1 = Card(variant="outlined", padding="md", theme=theme_dict)
    card1.add_child(Card.Header("Dashboard", "Real-time analytics"))
    card1.add_child(Separator(theme_dict))
    card1.add_child(Card.Content("View your performance metrics and insights"))
    card1.add_child(Badge("New", "success", theme_dict))
    card1.add_child(Card.Footer("Updated 5 min ago", align="right"))
    card1.render(slide2, left=0.5, top=1.5, width=4.5, height=3.0)

    # Card built with builder pattern
    builder = CompositionBuilder(theme_dict)
    children = (builder
                .title("Quick Stats")
                .description("Key metrics at a glance")
                .separator()
                .content("↑ 25% increase")
                .badge("Trending", "success")
                .build())

    card2 = Card(variant="elevated", padding="lg", theme=theme_dict)
    for child in children:
        card2.add_child(child)
    card2.render(slide2, left=5.5, top=1.5, width=4.0, height=3.0)

    # Slide 3: Metric Cards
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    theme.apply_to_slide(slide3)

    # Title
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Metric Cards with Variants"
    title_frame.paragraphs[0].font.size = Inches(0.4)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Different metric cards
    metrics = [
        {"label": "Revenue", "value": "$1.2M", "change": "+12%", "trend": "up"},
        {"label": "Users", "value": "45.2K", "change": "+8%", "trend": "up"},
        {"label": "Churn", "value": "3.2%", "change": "-1.5%", "trend": "down"},
        {"label": "NPS", "value": "72", "change": "→ 0%", "trend": "neutral"},
    ]

    variants_for_metrics = ["outlined", "elevated", "default", "outlined"]

    for i, (metric_data, var) in enumerate(zip(metrics, variants_for_metrics)):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 4.75
        y = 1.5 + row * 2.0

        metric = MetricCard(
            label=metric_data["label"],
            value=metric_data["value"],
            change=metric_data["change"],
            trend=metric_data["trend"],
            variant=var,
            theme=theme_dict
        )
        metric.render(slide3, left=x, top=y, width=4.25, height=1.5)

    # Save presentation
    output_path = "outputs/enhanced_components_showcase.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}\n")
    return output_path


def main():
    """Run all demonstrations."""
    print("\n" + "="*70)
    print(" Enhanced Component System Demo".center(70))
    print(" shadcn-inspired improvements for PowerPoint".center(70))
    print("="*70 + "\n")

    # Run demos
    demo_variants()
    demo_composition()
    demo_registry()

    # Create presentation
    output = create_showcase_presentation()

    print("="*70)
    print("\n✨ Demo Complete!\n")
    print("Key Features Demonstrated:")
    print("  1. Variant System - Composable variants like cva")
    print("  2. Composition Patterns - shadcn-style component composition")
    print("  3. Component Registry - LLM-friendly schemas and discovery")
    print(f"\nOpen the presentation to see it in action:")
    print(f"  {output}\n")
    print("="*70 + "\n")


if __name__ == "__main__":
    main()
