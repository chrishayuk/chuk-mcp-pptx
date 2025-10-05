#!/usr/bin/env python3
"""
Domain-Focused Showcase
Creates comprehensive presentations for each business domain using a consistent theme.
Shows the full range of charts and visualizations appropriate for each domain.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches
from typing import Dict, List, Any

# Import all chart components
from chuk_mcp_pptx.components.charts import (
    ColumnChart, BarChart, WaterfallChart,
    LineChart, AreaChart, SparklineChart,
    PieChart, DoughnutChart, SunburstChart,
    ScatterChart, BubbleChart,
    RadarChart, ComboChart, FunnelChart
)
from chuk_mcp_pptx.components import Card
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


class DomainFocusedShowcase:
    """Create comprehensive domain showcases with consistent theming."""
    
    def __init__(self, theme_name: str = "dark-blue"):
        """
        Initialize with a specific theme for consistency.
        
        Args:
            theme_name: Theme to use across all domain presentations
        """
        self.theme_manager = ThemeManager()
        self.theme_name = theme_name
        self.theme = self.theme_manager.get_theme(theme_name)
    
    async def create_general_business_showcase(self) -> Presentation:
        """Create comprehensive general business/strategy presentation."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        self.theme.apply_to_slide(title_slide)
        title_slide.shapes.title.text = "Business Strategy Dashboard"
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "Comprehensive Analytics Suite"
        
        # Slide 1: Executive Summary
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        self.theme.apply_to_slide(slide1)
        slide1.shapes.title.text = "Executive Summary"
        
        # KPI Cards
        for i, (title, value, change) in enumerate([
            ("Revenue", "$48.2M", "+12%"),
            ("Profit Margin", "23.5%", "+2.3pp"),
            ("Market Share", "18.5%", "+1.2pp"),
            ("Customer NPS", "72", "+5")
        ]):
            card = Card(
                title=title,
                description=f"{value} ({change} YoY)",
                variant="elevated" if i % 2 == 0 else "bordered",
                theme=self.theme.__dict__
            )
            card.render(slide1, left=0.5 + (i * 2.4), top=1.5, width=2.2, height=1.5)
        
        # Revenue Waterfall
        waterfall = WaterfallChart(
            categories=["Q1", "+Sales", "+Services", "-Costs", "-Tax", "Q2"],
            values=[10.5, 3.2, 1.8, -2.1, -0.9, 12.5],
            title="Quarterly Revenue Bridge ($M)",
            theme=self.theme.__dict__
        )
        await waterfall.render(slide1, left=0.5, top=3.5, width=4.5, height=3.5)
        
        # Market Share Pie
        pie = PieChart(
            categories=["Us", "Competitor A", "Competitor B", "Others"],
            values=[18.5, 24.2, 19.8, 37.5],
            variant="exploded",
            title="Market Share Distribution",
            theme=self.theme.__dict__
        )
        await pie.render(slide1, left=5.25, top=3.5, width=4.5, height=3.5)
        
        # Slide 2: Revenue Analysis
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        self.theme.apply_to_slide(slide2)
        slide2.shapes.title.text = "Revenue Analysis"
        
        # Regional Revenue Comparison
        regional = ColumnChart(
            categories=["North America", "Europe", "Asia Pacific", "Latin America", "Middle East"],
            series={
                "2022": [18.5, 14.2, 12.8, 6.4, 3.8],
                "2023": [21.2, 16.5, 15.9, 7.8, 4.6],
                "2024 Target": [24.0, 18.5, 19.0, 9.0, 5.5]
            },
            variant="clustered",
            title="Revenue by Region ($M)",
            theme=self.theme.__dict__
        )
        await regional.render(slide2, left=0.5, top=1.5, width=4.5, height=2.8)
        
        # Product Mix Sunburst
        sunburst_data = {
            "name": "Total Revenue",
            "value": 48.2,
            "children": [
                {
                    "name": "Software",
                    "value": 28.5,
                    "children": [
                        {"name": "SaaS", "value": 18.2},
                        {"name": "Licenses", "value": 10.3}
                    ]
                },
                {
                    "name": "Services",
                    "value": 15.3,
                    "children": [
                        {"name": "Consulting", "value": 9.1},
                        {"name": "Support", "value": 6.2}
                    ]
                },
                {
                    "name": "Hardware",
                    "value": 4.4
                }
            ]
        }
        sunburst = SunburstChart(
            data=sunburst_data,
            title="Revenue Breakdown",
            theme=self.theme.__dict__
        )
        await sunburst.render(slide2, left=5.25, top=1.5, width=4.5, height=2.8)
        
        # Growth Trends
        growth = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"],
            series={
                "Revenue": [3.8, 3.9, 4.1, 3.9, 4.2, 4.3, 4.1, 4.4, 4.5, 4.3, 4.6, 4.8],
                "Target": [3.5, 3.6, 3.8, 3.9, 4.0, 4.1, 4.2, 4.3, 4.4, 4.5, 4.6, 4.7],
                "Forecast": [3.8, 3.9, 4.1, 3.9, 4.2, 4.3, 4.1, 4.4, 4.5, 4.5, 4.7, 5.0]
            },
            variant="smooth",
            title="Monthly Revenue Trend ($M)",
            theme=self.theme.__dict__
        )
        await growth.render(slide2, left=0.5, top=4.8, width=9.25, height=2.5)
        
        # Slide 3: Sales Pipeline
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        self.theme.apply_to_slide(slide3)
        slide3.shapes.title.text = "Sales Pipeline & Conversion"
        
        # Sales Funnel
        funnel = FunnelChart(
            stages=["Prospects", "Qualified Leads", "Proposals", "Negotiation", "Closed Won"],
            values=[15000, 8500, 3200, 1800, 750],
            show_percentages=True,
            show_values=True,
            title="Q4 Sales Pipeline",
            theme=self.theme.__dict__
        )
        await funnel.render(slide3, left=0.5, top=1.5, width=4.5, height=5.5)
        
        # Conversion Metrics
        conversion = BarChart(
            categories=["Lead to Qualified", "Qualified to Proposal", "Proposal to Negotiation", "Negotiation to Close"],
            series={
                "Current": [56.7, 37.6, 56.3, 41.7],
                "Target": [60.0, 45.0, 60.0, 50.0],
                "Industry Avg": [50.0, 35.0, 50.0, 40.0]
            },
            variant="clustered",
            title="Conversion Rates (%)",
            theme=self.theme.__dict__
        )
        await conversion.render(slide3, left=5.25, top=1.5, width=4.5, height=2.5)
        
        # Win/Loss Analysis
        winloss = DoughnutChart(
            categories=["Won", "Lost - Price", "Lost - Features", "Lost - Competition", "No Decision"],
            values=[750, 280, 180, 320, 270],
            variant="modern",
            title="Win/Loss Analysis",
            theme=self.theme.__dict__
        )
        await winloss.render(slide3, left=5.25, top=4.5, width=4.5, height=2.5)
        
        # Slide 4: Customer Analytics
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        self.theme.apply_to_slide(slide4)
        slide4.shapes.title.text = "Customer Analytics"
        
        # Customer Segments
        segments = ColumnChart(
            categories=["Enterprise", "Mid-Market", "SMB", "Startup"],
            series={
                "Customer Count": [145, 320, 680, 455],
                "Revenue ($M)": [22.5, 15.8, 7.2, 2.7]
            },
            variant="clustered",
            title="Customer Segmentation",
            theme=self.theme.__dict__
        )
        await segments.render(slide4, left=0.5, top=1.5, width=4.5, height=2.5)
        
        # Churn Analysis
        churn = LineChart(
            categories=["Q1-22", "Q2-22", "Q3-22", "Q4-22", "Q1-23", "Q2-23", "Q3-23", "Q4-23"],
            series={
                "Gross Churn": [3.2, 3.5, 3.1, 2.9, 2.8, 2.6, 2.4, 2.2],
                "Net Churn": [1.8, 2.1, 1.7, 1.5, 1.3, 1.1, 0.9, 0.7],
                "Target": [2.0, 2.0, 2.0, 2.0, 2.0, 2.0, 2.0, 2.0]
            },
            variant="markers",
            title="Churn Rate Trends (%)",
            theme=self.theme.__dict__
        )
        await churn.render(slide4, left=5.25, top=1.5, width=4.5, height=2.5)
        
        # Customer Satisfaction
        satisfaction = RadarChart(
            categories=["Product Quality", "Customer Support", "Value for Money", "Features", "Reliability", "Ease of Use"],
            series={
                "Current": [8.2, 8.5, 7.8, 8.0, 8.7, 8.3],
                "Previous": [7.9, 8.2, 7.5, 7.8, 8.5, 8.0],
                "Target": [8.5, 8.5, 8.5, 8.5, 8.5, 8.5]
            },
            variant="filled",
            max_value=10,
            title="Customer Satisfaction Scores",
            theme=self.theme.__dict__
        )
        await satisfaction.render(slide4, left=0.5, top=4.5, width=4.5, height=2.5)
        
        # Lifetime Value
        ltv = BubbleChart(
            series_data=[
                {
                    "name": "Enterprise",
                    "points": [[24, 85, 45], [36, 92, 52], [48, 95, 58]]  # [months, retention%, value]
                },
                {
                    "name": "Mid-Market",
                    "points": [[24, 78, 28], [36, 82, 32], [48, 85, 35]]
                },
                {
                    "name": "SMB",
                    "points": [[24, 65, 12], [36, 68, 14], [48, 70, 15]]
                }
            ],
            size_scale=2.0,
            transparency=30,
            title="Customer LTV Analysis",
            theme=self.theme.__dict__
        )
        await ltv.render(slide4, left=5.25, top=4.5, width=4.5, height=2.5)
        
        # Slide 5: Competitive Analysis
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        self.theme.apply_to_slide(slide5)
        slide5.shapes.title.text = "Competitive Landscape"
        
        # Competitive Positioning
        positioning = ScatterChart(
            series_data=[
                {
                    "name": "Us",
                    "x_values": [7.5],
                    "y_values": [8.2]
                },
                {
                    "name": "Competitors",
                    "x_values": [6.2, 8.1, 5.8, 7.2, 6.8],
                    "y_values": [7.5, 7.8, 6.2, 8.5, 7.0]
                }
            ],
            show_trendline=False,
            marker_size=15,
            title="Market Positioning (Price vs Quality)",
            theme=self.theme.__dict__
        )
        await positioning.render(slide5, left=0.5, top=1.5, width=4.5, height=3.0)
        
        # Feature Comparison
        features = RadarChart(
            categories=["Features", "Price", "Support", "Performance", "Security", "Integration"],
            series={
                "Us": [8.5, 7.0, 9.0, 8.8, 9.2, 8.0],
                "Competitor A": [8.0, 8.5, 7.5, 8.2, 8.5, 8.5],
                "Competitor B": [7.5, 6.5, 8.0, 7.8, 8.8, 7.0],
                "Market Leader": [9.0, 6.0, 8.5, 9.0, 9.0, 9.5]
            },
            variant="filled",
            max_value=10,
            title="Competitive Feature Analysis",
            theme=self.theme.__dict__
        )
        await features.render(slide5, left=5.25, top=1.5, width=4.5, height=3.0)
        
        # Market Share Trends
        share_trends = AreaChart(
            categories=["2020", "2021", "2022", "2023", "2024E"],
            series={
                "Us": [12, 14, 16, 18.5, 21],
                "Competitor A": [28, 27, 26, 24.2, 23],
                "Competitor B": [22, 21, 20.5, 19.8, 19],
                "Others": [38, 38, 37.5, 37.5, 37]
            },
            variant="stacked",
            title="Market Share Evolution (%)",
            theme=self.theme.__dict__
        )
        await share_trends.render(slide5, left=0.5, top=5.0, width=9.25, height=2.0)
        
        return prs
    
    async def create_tech_teams_showcase(self) -> Presentation:
        """Create comprehensive tech team/engineering presentation."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        self.theme.apply_to_slide(title_slide)
        title_slide.shapes.title.text = "Engineering Team Dashboard"
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "Development Metrics & Analytics"
        
        # Slide 1: Sprint Overview
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        self.theme.apply_to_slide(slide1)
        slide1.shapes.title.text = "Current Sprint Overview"
        
        # Sprint Burndown
        burndown = LineChart(
            categories=["Day 1", "Day 2", "Day 3", "Day 4", "Day 5", "Day 6", "Day 7", "Day 8", "Day 9", "Day 10"],
            series={
                "Ideal": [100, 90, 80, 70, 60, 50, 40, 30, 20, 10],
                "Actual": [100, 95, 88, 78, 72, 65, 58, 48, 35, 22],
                "Forecast": [100, 95, 88, 78, 72, 65, 58, 48, 38, 25]
            },
            variant="markers",
            title="Sprint Burndown (Story Points)",
            theme=self.theme.__dict__
        )
        await burndown.render(slide1, left=0.5, top=1.5, width=4.5, height=2.5)
        
        # Team Velocity
        velocity = ColumnChart(
            categories=["Sprint 1", "Sprint 2", "Sprint 3", "Sprint 4", "Sprint 5", "Sprint 6"],
            series={
                "Committed": [45, 48, 52, 50, 55, 58],
                "Delivered": [42, 47, 48, 52, 53, 60],
                "Bugs": [5, 3, 6, 4, 3, 2]
            },
            variant="clustered",
            title="Team Velocity Trend",
            theme=self.theme.__dict__
        )
        await velocity.render(slide1, left=5.25, top=1.5, width=4.5, height=2.5)
        
        # Story Status
        story_status = DoughnutChart(
            categories=["Completed", "In Progress", "In Review", "Blocked", "Not Started"],
            values=[35, 22, 15, 5, 23],
            variant="modern",
            title="Story Status Distribution",
            theme=self.theme.__dict__
        )
        await story_status.render(slide1, left=0.5, top=4.5, width=3.0, height=2.5)
        
        # Team Allocation
        allocation = BarChart(
            categories=["Frontend", "Backend", "DevOps", "QA", "Mobile"],
            series={
                "Allocated": [18, 22, 8, 12, 10],
                "Available": [15, 20, 10, 10, 8]
            },
            variant="clustered",
            title="Team Allocation (Hours/Day)",
            theme=self.theme.__dict__
        )
        await allocation.render(slide1, left=3.75, top=4.5, width=3.0, height=2.5)
        
        # Sprint Health Metrics
        health = RadarChart(
            categories=["Velocity", "Quality", "Predictability", "Morale", "Innovation", "Collaboration"],
            series={
                "Current": [8, 7.5, 7, 8.5, 7, 9],
                "Previous": [7.5, 7, 6.5, 8, 6.5, 8.5],
                "Target": [8, 8, 8, 8, 8, 8]
            },
            variant="filled",
            max_value=10,
            title="Sprint Health Metrics",
            theme=self.theme.__dict__
        )
        await health.render(slide1, left=7.0, top=4.5, width=2.75, height=2.5)
        
        # Continue with more slides for tech teams...
        # (Code Quality, CI/CD Metrics, Bug Tracking, Performance Monitoring, etc.)
        
        return prs
    
    async def create_finance_showcase(self) -> Presentation:
        """Create comprehensive finance/CFO presentation."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        self.theme.apply_to_slide(title_slide)
        title_slide.shapes.title.text = "Financial Dashboard"
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = "Q4 2024 Financial Review"
        
        # Slide 1: Financial Summary
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        self.theme.apply_to_slide(slide1)
        slide1.shapes.title.text = "Financial Summary"
        
        # EBITDA Waterfall
        ebitda = WaterfallChart(
            categories=["Revenue", "COGS", "SG&A", "R&D", "Other", "EBITDA"],
            values=[125.5, -48.2, -32.8, -12.5, -4.5, 27.5],
            title="EBITDA Bridge ($M)",
            theme=self.theme.__dict__
        )
        await ebitda.render(slide1, left=0.5, top=1.5, width=4.5, height=2.8)
        
        # Revenue & Margins Combo
        margins = ComboChart(
            categories=["Q1-23", "Q2-23", "Q3-23", "Q4-23", "Q1-24", "Q2-24", "Q3-24", "Q4-24E"],
            column_series={
                "Revenue": [28.2, 29.5, 31.2, 32.8, 30.5, 31.8, 33.2, 35.0]
            },
            line_series={
                "Gross Margin %": [42, 43, 44, 44.5, 45, 45.5, 46, 46.5],
                "EBITDA Margin %": [18, 19, 20, 21, 21.5, 22, 22.5, 23]
            },
            title="Revenue & Margin Trends",
            theme=self.theme.__dict__
        )
        await margins.render(slide1, left=5.25, top=1.5, width=4.5, height=2.8)
        
        # Cash Flow
        cashflow = AreaChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"],
            series={
                "Operating": [8.2, 7.5, 9.1, 8.8, 9.5, 10.2, 9.8, 10.5, 11.2, 10.8, 11.5, 12.0],
                "Investing": [-2.5, -1.8, -3.2, -2.0, -2.2, -1.5, -2.8, -1.2, -1.8, -2.0, -1.5, -1.0],
                "Financing": [-1.0, -1.0, -1.0, -1.0, -1.0, -1.0, -1.0, -1.0, -1.0, -1.0, -1.0, -1.0]
            },
            variant="stacked",
            title="Cash Flow Analysis ($M)",
            theme=self.theme.__dict__
        )
        await cashflow.render(slide1, left=0.5, top=4.8, width=9.25, height=2.2)
        
        return prs
    
    async def generate_all_domain_showcases(self):
        """Generate comprehensive showcase for each domain."""
        output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs", "domain_showcases")
        os.makedirs(output_dir, exist_ok=True)
        
        print(f"\n📊 Generating Domain-Focused Showcases")
        print("=" * 70)
        print(f"Theme: {self.theme_name}")
        print(f"Creating comprehensive presentations for each business domain")
        print()
        
        domains = {
            "general_business": ("General Business & Strategy", self.create_general_business_showcase),
            "tech_teams": ("Engineering & Tech Teams", self.create_tech_teams_showcase),
            "finance": ("Finance & CFO", self.create_finance_showcase),
            # Add more domains as needed
        }
        
        for domain_key, (domain_name, create_func) in domains.items():
            print(f"📈 Creating {domain_name} showcase...")
            
            # Create presentation
            prs = await create_func()
            
            # Save presentation
            filename = f"domain_{domain_key}_{self.theme_name.replace('-', '_')}.pptx"
            filepath = os.path.join(output_dir, filename)
            prs.save(filepath)
            
            print(f"   ✅ Saved: {filename}")
            print(f"      • Comprehensive {domain_name.lower()} analytics")
            print(f"      • Multiple dashboard views")
            print(f"      • Industry-specific metrics")
            print(f"      • Consistent {self.theme_name} theme throughout")
        
        print("\n" + "=" * 70)
        print("✨ All domain showcases generated successfully!")
        print(f"📁 Output directory: {output_dir}")
        
        return output_dir


async def main():
    """Generate all domain-focused showcases."""
    # Use dark-blue as the default professional theme
    showcase = DomainFocusedShowcase(theme_name="dark-blue")
    output_dir = await showcase.generate_all_domain_showcases()
    
    print("\n🎯 Domain Showcases Complete!")
    print(f"Theme used: {showcase.theme_name}")
    print("\nEach presentation demonstrates:")
    print("  • Complete analytics suite for the domain")
    print("  • Industry-appropriate charts and metrics")
    print("  • Professional dashboard layouts")
    print("  • Consistent theming throughout")
    print("\n💡 Perfect for:")
    print("  • Domain-specific presentations")
    print("  • Industry best practices")
    print("  • Comprehensive analytics examples")


if __name__ == "__main__":
    asyncio.run(main())