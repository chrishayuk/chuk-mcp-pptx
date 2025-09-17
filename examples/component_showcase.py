#!/usr/bin/env python3
"""
Component Showcase - Demonstrating the new component-based architecture.
Shows how to build beautiful, theme-driven presentations using components.
"""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from pptx import Presentation
from pptx.util import Inches

# Import components
from chuk_mcp_pptx.components import (
    Button, IconButton,
    Card, MetricCard, FeatureCard,
    CodeBlock, Terminal
)
from chuk_mcp_pptx.components.chart import BarChart, LineChart, PieChart
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


async def create_component_showcase():
    """Create a showcase presentation using the component system."""
    
    print("\n🎨 Creating Component Showcase Presentation")
    print("=" * 70)
    
    # Initialize presentation and theme manager
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    theme_manager = ThemeManager()
    
    # Slide 1: Dark Theme with Buttons
    print("\n1. Creating Dark Theme Button Showcase...")
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    theme = theme_manager.get_theme("dark-violet")
    theme.apply_to_slide(slide1)
    
    # Add title
    title_shape = slide1.shapes.title
    title_shape.text = "Button Components"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Add button variants
    button_variants = [
        ("Primary", "primary", 1.5),
        ("Secondary", "secondary", 3.0),
        ("Outline", "outline", 4.5),
        ("Ghost", "ghost", 6.0),
        ("Destructive", "destructive", 7.5),
    ]
    
    for text, variant, left in button_variants:
        btn = Button(text, variant=variant, size="md", theme=theme.__dict__)
        btn.render(slide1, left=left, top=2.0)
    
    # Add icon buttons
    icon_buttons = [
        ("play", 1.5),
        ("pause", 2.5),
        ("settings", 3.5),
        ("star", 4.5),
        ("heart", 5.5),
    ]
    
    for icon, left in icon_buttons:
        icon_btn = IconButton(icon, variant="ghost", theme=theme.__dict__)
        icon_btn.render(slide1, left=left, top=3.5)
    
    # Slide 2: Card Components with Metrics
    print("2. Creating Card Components...")
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("dark-blue")
    theme.apply_to_slide(slide2)
    
    # Add title
    title_shape = slide2.shapes.title
    title_shape.text = "Card Components & Metrics"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Add metric cards
    metrics = [
        ("Revenue", "$1.2M", "+12%", "up", 1.0),
        ("Users", "45.2K", "+8%", "up", 3.5),
        ("Retention", "92%", "-2%", "down", 6.0),
        ("NPS", "4.8", "0%", "neutral", 8.5),
    ]
    
    for label, value, change, trend, left in metrics:
        metric = MetricCard(label, value, change, trend, theme=theme.__dict__)
        metric.render(slide2, left=left, top=2.0)
    
    # Add feature card
    features = FeatureCard(
        icon="🚀",
        title="Key Features",
        features=["Component-based", "Theme-aware", "Design tokens", "Beautiful defaults"],
        theme=theme.__dict__
    )
    features.render(slide2, left=3.5, top=4.0, width=3.0, height=2.5)
    
    # Slide 3: Code Components
    print("3. Creating Code Component Examples...")
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("cyberpunk")
    theme.apply_to_slide(slide3)
    
    # Add title
    title_shape = slide3.shapes.title
    title_shape.text = "Code Components"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Add Python code block
    python_code = """async def create_presentation(theme):
    prs = Presentation()
    manager = ThemeManager()
    theme = manager.get_theme(theme)
    
    slide = prs.slides.add_slide(layout)
    theme.apply_to_slide(slide)
    
    return prs"""
    
    code_block = CodeBlock(python_code, "python", show_line_numbers=True, theme=theme.__dict__)
    code_block.render(slide3, left=1.0, top=2.0, width=4.0, height=2.5)
    
    # Add terminal output
    terminal_output = """npm run build
Building presentation components...
✓ Tokens compiled
✓ Components built
✓ Themes generated
Build completed in 2.3s"""
    
    terminal = Terminal(terminal_output, theme=theme.__dict__)
    terminal.render(slide3, left=5.5, top=2.0, width=4.0, height=2.5)
    
    # Slide 4: Chart Components
    print("4. Creating Chart Components...")
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("sunset")
    theme.apply_to_slide(slide4)
    
    # Add title
    title_shape = slide4.shapes.title
    title_shape.text = "Data Visualization Components"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Bar chart
    bar_chart = BarChart(
        title="Q4 Performance",
        categories=["Oct", "Nov", "Dec"],
        series={
            "Revenue": [120, 145, 165],
            "Profit": [45, 52, 61]
        },
        theme=theme.__dict__
    )
    await bar_chart.render(slide4, left=0.5, top=2.0, width=4.5, height=3.0)
    
    # Pie chart
    pie_chart = PieChart(
        title="Market Share",
        categories=["Product A", "Product B", "Product C", "Others"],
        values=[35, 28, 22, 15],
        explode=0,
        theme=theme.__dict__
    )
    await pie_chart.render(slide4, left=5.0, top=2.0, width=4.5, height=3.0)
    
    # Slide 5: Light Theme Showcase
    print("5. Creating Light Theme Examples...")
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("light")
    theme.apply_to_slide(slide5)
    
    # Add title
    title_shape = slide5.shapes.title
    title_shape.text = "Light Theme Components"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Add various cards
    card1 = Card(
        title="Clean Design",
        description="Minimal and elegant components with light backgrounds.",
        variant="bordered",
        theme=theme.__dict__
    )
    card1.render(slide5, left=1.0, top=2.0, width=3.0, height=2.0)
    
    card2 = Card(
        title="Professional",
        description="Perfect for corporate presentations and reports.",
        variant="elevated",
        theme=theme.__dict__
    )
    card2.render(slide5, left=4.5, top=2.0, width=3.0, height=2.0)
    
    # Add buttons in light theme
    for i, variant in enumerate(["primary", "secondary", "outline"]):
        btn = Button(f"{variant.title()}", variant=variant, theme=theme.__dict__)
        btn.render(slide5, left=2.5 + i * 2, top=5.0)
    
    # Slide 6: Corporate Theme
    print("6. Creating Corporate Theme Example...")
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    theme = theme_manager.get_theme("corporate")
    theme.apply_to_slide(slide6)
    
    # Add title
    title_shape = slide6.shapes.title
    title_shape.text = "Corporate Theme"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")
    
    # Add business metrics
    corp_metrics = [
        ("Q4 Revenue", "$4.2M", "+18%", "up", 1.0),
        ("Employees", "1,250", "+5%", "up", 3.5),
        ("Offices", "12", "+2", "up", 6.0),
        ("Projects", "48", "+15", "up", 8.5),
    ]
    
    for label, value, change, trend, left in corp_metrics:
        metric = MetricCard(label, value, change, trend, theme=theme.__dict__)
        metric.render(slide6, left=left, top=2.0, width=1.8)
    
    # Add line chart for corporate data
    line_chart = LineChart(
        title="Annual Growth Trend",
        categories=["Q1", "Q2", "Q3", "Q4"],
        series={
            "Revenue": [3.2, 3.5, 3.8, 4.2],
            "Profit": [0.8, 0.9, 1.0, 1.2]
        },
        smooth=True,
        markers=True,
        theme=theme.__dict__
    )
    await line_chart.render(slide6, left=2.5, top=4.0, width=5.0, height=3.0)
    
    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "component_showcase.pptx")
    prs.save(output_path)
    print(f"\n✅ Created {output_path}")
    
    print("\n🎨 Component Showcase Features:")
    print("  • Multiple theme variations (dark, light, cyberpunk, gradient)")
    print("  • Reusable components (buttons, cards, charts, code blocks)")
    print("  • Design token system for consistent styling")
    print("  • Theme-aware color schemes")
    print("  • Professional, modern design")
    
    print("\n💡 This demonstrates:")
    print("  • Component-based architecture")
    print("  • Theme flexibility")
    print("  • Design system approach")
    print("  • Similar to shadcn/ui but for PowerPoint")


if __name__ == "__main__":
    asyncio.run(create_component_showcase())