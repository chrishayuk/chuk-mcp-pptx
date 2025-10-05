"""
Pytest configuration and shared fixtures for PowerPoint MCP tests.
"""

import pytest
import sys
import os
from pathlib import Path
from typing import Dict, Any
from unittest.mock import MagicMock, Mock
from pptx import Presentation
from pptx.slide import Slide

# Add src to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from chuk_mcp_pptx.themes.theme_manager import ThemeManager


@pytest.fixture
def mock_slide():
    """Create a mock slide for testing."""
    slide = MagicMock(spec=Slide)
    
    # Mock shapes collection
    slide.shapes = MagicMock()
    slide.shapes.add_chart = MagicMock(return_value=MagicMock())
    slide.shapes.add_shape = MagicMock(return_value=MagicMock())
    slide.shapes.add_textbox = MagicMock(return_value=MagicMock())
    slide.shapes.title = MagicMock()
    slide.shapes.title.text = ""
    
    # Mock placeholders
    slide.placeholders = {}
    
    return slide


@pytest.fixture
def mock_presentation():
    """Create a mock presentation for testing."""
    prs = MagicMock(spec=Presentation)
    prs.slides = MagicMock()
    prs.slide_layouts = [MagicMock() for _ in range(11)]
    prs.slide_width = 9144000  # 10 inches in EMU
    prs.slide_height = 5143500  # 5.625 inches in EMU
    
    # Mock add_slide method
    def add_slide(layout):
        return mock_slide()
    
    prs.slides.add_slide = MagicMock(side_effect=add_slide)
    
    return prs


@pytest.fixture
def sample_theme():
    """Get a sample theme for testing."""
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark")
    return theme.__dict__ if hasattr(theme, '__dict__') else {}


@pytest.fixture
def dark_theme():
    """Get dark theme for testing."""
    return {
        "name": "dark",
        "mode": "dark",
        "primary_hue": "blue",
        "background": {"DEFAULT": "#0f0f0f"},
        "foreground": {"DEFAULT": "#f5f5f5"},
        "primary": {"DEFAULT": "#3b82f6"},
        "secondary": {"DEFAULT": "#64748b"},
        "accent": {"DEFAULT": "#8b5cf6"},
        "chart": ["#3b82f6", "#10b981", "#f59e0b", "#ef4444", "#8b5cf6", "#ec4899"]
    }


@pytest.fixture
def light_theme():
    """Get light theme for testing."""
    return {
        "name": "light",
        "mode": "light",
        "primary_hue": "blue",
        "background": {"DEFAULT": "#ffffff"},
        "foreground": {"DEFAULT": "#0f0f0f"},
        "primary": {"DEFAULT": "#2563eb"},
        "secondary": {"DEFAULT": "#64748b"},
        "accent": {"DEFAULT": "#7c3aed"},
        "chart": ["#2563eb", "#059669", "#d97706", "#dc2626", "#7c3aed", "#db2777"]
    }


@pytest.fixture
def sample_chart_data():
    """Sample data for chart testing."""
    return {
        "column_data": {
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": {
                "Revenue": [100, 120, 140, 160],
                "Profit": [20, 25, 30, 35]
            }
        },
        "line_data": {
            "categories": ["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            "series": {
                "Sales": [65, 68, 72, 78, 85, 92],
                "Target": [70, 75, 80, 85, 90, 95]
            }
        },
        "pie_data": {
            "categories": ["Product A", "Product B", "Product C", "Product D"],
            "values": [35, 25, 20, 20]
        },
        "scatter_data": [
            {
                "name": "Series 1",
                "x_values": [1, 2, 3, 4, 5],
                "y_values": [10, 20, 15, 25, 30]
            },
            {
                "name": "Series 2",
                "x_values": [1, 2, 3, 4, 5],
                "y_values": [5, 15, 20, 18, 22]
            }
        ],
        "bubble_data": [
            {
                "name": "Group A",
                "points": [[10, 20, 5], [15, 25, 8], [20, 30, 12]]
            },
            {
                "name": "Group B",
                "points": [[5, 15, 3], [8, 18, 5], [12, 22, 7]]
            }
        ],
        "funnel_data": {
            "stages": ["Leads", "Qualified", "Proposal", "Negotiation", "Closed"],
            "values": [1000, 750, 300, 150, 50]
        },
        "waterfall_data": {
            "categories": ["Start", "Q1", "Q2", "Q3", "Q4", "End"],
            "values": [100, 20, -10, 15, 25, 150]
        }
    }


@pytest.fixture
def sample_component_data():
    """Sample data for UI component testing."""
    return {
        "card": {
            "title": "Test Card",
            "description": "This is a test card description"
        },
        "button": {
            "text": "Click Me",
            "variant": "primary"
        },
        "code_block": {
            "code": "print('Hello, World!')",
            "language": "python"
        }
    }


class MockChartData:
    """Mock chart data for testing."""
    def __init__(self):
        self.series = []
    
    def add_series(self, name, values=None):
        """Add a series to the chart data."""
        series = MagicMock()
        series.name = name
        series.values = values or []
        series.add_data_point = MagicMock()
        self.series.append(series)
        return series


@pytest.fixture
def mock_chart_data():
    """Create mock chart data."""
    return MockChartData()


@pytest.fixture
def temp_output_dir(tmp_path):
    """Create a temporary output directory for test files."""
    output_dir = tmp_path / "test_outputs"
    output_dir.mkdir()
    return output_dir


# Test helper functions
def assert_color_valid(color):
    """Assert that a color value is valid hex format."""
    if isinstance(color, str):
        assert color.startswith('#'), f"Color {color} should start with #"
        assert len(color) in [4, 7], f"Color {color} should be 3 or 6 hex digits"
        # Check if valid hex
        try:
            int(color[1:], 16)
        except ValueError:
            pytest.fail(f"Color {color} is not valid hex")


def assert_chart_renders(chart, slide):
    """Assert that a chart renders without errors."""
    try:
        # For async charts
        import asyncio
        if asyncio.iscoroutinefunction(chart.render):
            asyncio.run(chart.render(slide, left=1, top=1, width=4, height=3))
        else:
            chart.render(slide, left=1, top=1, width=4, height=3)
    except Exception as e:
        pytest.fail(f"Chart failed to render: {e}")


def assert_component_renders(component, slide):
    """Assert that a component renders without errors."""
    try:
        component.render(slide, left=1, top=1, width=3, height=2)
    except Exception as e:
        pytest.fail(f"Component failed to render: {e}")


# Async test helper
def async_test(coro):
    """Decorator to run async tests."""
    import asyncio
    
    def wrapper(*args, **kwargs):
        loop = asyncio.get_event_loop()
        return loop.run_until_complete(coro(*args, **kwargs))
    
    return wrapper