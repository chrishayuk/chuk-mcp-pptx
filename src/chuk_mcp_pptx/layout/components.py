# src/chuk_mcp_pptx/layout/components.py
"""
Layout components for PowerPoint presentations.
Grid system, containers, stacks - inspired by modern UI frameworks.
"""

from typing import Optional, Dict, Any, List, Literal
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ..components.base import Component
from ..tokens.spacing import SPACING, MARGINS, PADDING, GAPS, GRID, CONTAINERS
from .helpers import (
    SLIDE_WIDTH, SLIDE_HEIGHT, CONTENT_WIDTH, CONTENT_HEIGHT,
    validate_position, calculate_grid_layout
)
from ..registry import component, ComponentCategory, prop, example


@component(
    name="Container",
    category=ComponentCategory.LAYOUT,
    description="Responsive container for centering and constraining content width",
    props=[
        prop("size", "string", "Container size",
             options=["sm", "md", "lg", "xl", "2xl", "full"],
             default="lg"),
        prop("padding", "string", "Internal padding",
             options=["none", "sm", "md", "lg", "xl"],
             default="md"),
        prop("center", "boolean", "Center horizontally", default=True),
    ],
    examples=[
        example(
            "Centered content container",
            """
container = Container(size="lg", padding="md")
container.render(slide, top=1.5)
            """,
            size="lg",
            padding="md"
        )
    ],
    tags=["layout", "container", "responsive"]
)
class Container(Component):
    """
    Container component for centering and constraining content.

    Usage:
        # Standard container
        container = Container(size="lg")
        container.render(slide, top=1.5)

        # Full width
        container = Container(size="full", padding="none")
        container.render(slide, top=0)
    """

    def __init__(self,
                 size: Literal["sm", "md", "lg", "xl", "2xl", "full"] = "lg",
                 padding: Literal["none", "sm", "md", "lg", "xl"] = "md",
                 center: bool = True,
                 theme: Optional[Dict[str, Any]] = None):
        super().__init__(theme)
        self.size = size
        self.padding = padding
        self.center = center

    def render(self, slide, top: float = 0, height: Optional[float] = None):
        """
        Render container and return its bounds.

        Returns:
            Dict with container dimensions for child rendering
        """
        container_width = CONTAINERS.get(self.size, CONTAINERS["lg"])
        padding_value = PADDING.get(self.padding, PADDING["md"])

        # Center horizontally if requested
        if self.center:
            left = (SLIDE_WIDTH - container_width) / 2
        else:
            left = 0.5

        # Container doesn't render a visual element, just returns bounds
        return {
            'left': left + padding_value,
            'top': top + padding_value,
            'width': container_width - (2 * padding_value),
            'height': (height or SLIDE_HEIGHT - top) - (2 * padding_value)
        }


@component(
    name="Grid",
    category=ComponentCategory.LAYOUT,
    description="12-column grid system for flexible layouts",
    props=[
        prop("columns", "number", "Number of columns", default=12),
        prop("gap", "string", "Gap between columns",
             options=["none", "xs", "sm", "md", "lg", "xl"],
             default="md"),
        prop("rows", "number", "Number of rows", default=1),
    ],
    examples=[
        example(
            "Three column grid",
            """
grid = Grid(columns=3, gap="md")
positions = grid.get_cell_positions(slide, top=2.0, height=3.0)
            """,
            columns=3,
            gap="md"
        )
    ],
    tags=["layout", "grid", "columns"]
)
class Grid(Component):
    """
    12-column grid system (like Bootstrap/Tailwind).

    Usage:
        # 3-column grid
        grid = Grid(columns=3, gap="md")
        cells = grid.get_cell_positions(slide, top=2.0, height=3.0)

        for i, cell in enumerate(cells):
            # Render content in each cell
            card.render(slide, **cell)
    """

    def __init__(self,
                 columns: int = 12,
                 gap: Literal["none", "xs", "sm", "md", "lg", "xl"] = "md",
                 rows: int = 1,
                 bounds: Optional[Dict[str, float]] = None,
                 theme: Optional[Dict[str, Any]] = None):
        """
        Initialize grid.

        Args:
            columns: Number of columns
            gap: Gap between cells
            rows: Number of rows
            bounds: Optional container bounds dict with 'left', 'top', 'width', 'height'
                   If provided, grid will use these as defaults for all cell calculations
            theme: Optional theme
        """
        super().__init__(theme)
        self.columns = columns
        self.gap = GAPS.get(gap, GAPS["md"])
        self.rows = rows
        self.bounds = bounds or {}

    def get_cell_positions(self,
                          slide,
                          left: float = 0.5,
                          top: float = 1.5,
                          width: Optional[float] = None,
                          height: Optional[float] = None) -> List[Dict[str, float]]:
        """
        Calculate positions for all grid cells.

        Returns:
            List of dicts with 'left', 'top', 'width', 'height' for each cell
        """
        grid_width = width or CONTENT_WIDTH
        grid_height = height or CONTENT_HEIGHT

        # Calculate cell dimensions
        total_h_gap = self.gap * (self.columns - 1)
        total_v_gap = self.gap * (self.rows - 1)

        cell_width = (grid_width - total_h_gap) / self.columns
        cell_height = (grid_height - total_v_gap) / self.rows

        positions = []
        for row in range(self.rows):
            for col in range(self.columns):
                cell_left = left + col * (cell_width + self.gap)
                cell_top = top + row * (cell_height + self.gap)

                positions.append({
                    'left': cell_left,
                    'top': cell_top,
                    'width': cell_width,
                    'height': cell_height,
                    'row': row,
                    'col': col
                })

        return positions

    def get_span(self,
                 col_span: int = 1,
                 row_span: int = 1,
                 col_start: int = 0,
                 row_start: int = 0,
                 left: float = 0.5,
                 top: float = 1.5,
                 width: Optional[float] = None,
                 height: Optional[float] = None) -> Dict[str, float]:
        """
        Get position for a cell that spans multiple columns/rows.

        Args:
            col_span: Number of columns to span
            row_span: Number of rows to span
            col_start: Starting column (0-indexed)
            row_start: Starting row (0-indexed)

        Returns:
            Dict with position and dimensions
        """
        grid_width = width or CONTENT_WIDTH
        grid_height = height or CONTENT_HEIGHT

        # Calculate cell dimensions
        total_h_gap = self.gap * (self.columns - 1)
        total_v_gap = self.gap * (self.rows - 1)

        cell_width = (grid_width - total_h_gap) / self.columns
        cell_height = (grid_height - total_v_gap) / self.rows

        # Calculate span dimensions
        span_width = cell_width * col_span + self.gap * (col_span - 1)
        span_height = cell_height * row_span + self.gap * (row_span - 1)

        # Calculate position
        cell_left = left + col_start * (cell_width + self.gap)
        cell_top = top + row_start * (cell_height + self.gap)

        return {
            'left': cell_left,
            'top': cell_top,
            'width': span_width,
            'height': span_height
        }

    def get_cell(self,
                 col_span: int = 1,
                 row_span: int = 1,
                 col_start: int = 0,
                 row_start: int = 0,
                 left: Optional[float] = None,
                 top: Optional[float] = None,
                 width: Optional[float] = None,
                 height: Optional[float] = None,
                 auto_height: bool = True) -> Dict[str, float]:
        """
        Get position for a cell with optional auto-height for components.

        This is the preferred method for laying out components in a grid.
        Use auto_height=True (default) for components that calculate their own height.
        Use auto_height=False when you need fixed-height cells.

        Args:
            col_span: Number of columns to span
            row_span: Number of rows to span
            col_start: Starting column (0-indexed)
            row_start: Starting row (0-indexed)
            left: Grid left position (uses bounds if not provided)
            top: Grid top position (uses bounds if not provided)
            width: Grid total width (uses bounds if not provided)
            height: Grid total height (uses bounds if not provided)
            auto_height: If True, omit height from result (for auto-sizing components)

        Returns:
            Dict with 'left', 'top', 'width' and optionally 'height'

        Example:
            # With bounds set on grid
            grid = Grid(columns=12, rows=2, bounds=container_bounds)
            pos = grid.get_cell(col_span=8, col_start=0, row_start=1)
            card.render(slide, **pos)

            # Without bounds (manual positioning)
            grid = Grid(columns=12, rows=2)
            pos = grid.get_cell(col_span=8, col_start=0, row_start=1,
                               left=0.5, top=2.0, width=9.0, height=4.0)
            card.render(slide, **pos)
        """
        # Use bounds as defaults if available
        grid_left = left if left is not None else self.bounds.get('left', 0.5)
        grid_top = top if top is not None else self.bounds.get('top', 1.5)
        grid_width = width if width is not None else self.bounds.get('width', CONTENT_WIDTH)
        grid_height = height if height is not None else self.bounds.get('height', CONTENT_HEIGHT)

        # Get full position from get_span
        position = self.get_span(
            col_span=col_span,
            row_span=row_span,
            col_start=col_start,
            row_start=row_start,
            left=grid_left,
            top=grid_top,
            width=grid_width,
            height=grid_height
        )

        # Remove height if auto-sizing
        if auto_height:
            return {k: v for k, v in position.items() if k != 'height'}

        return position


@component(
    name="Stack",
    category=ComponentCategory.LAYOUT,
    description="Stack elements vertically or horizontally with consistent spacing",
    props=[
        prop("direction", "string", "Stack direction",
             options=["vertical", "horizontal"],
             default="vertical"),
        prop("gap", "string", "Gap between items",
             options=["none", "xs", "sm", "md", "lg", "xl"],
             default="md"),
        prop("align", "string", "Alignment",
             options=["start", "center", "end", "stretch"],
             default="start"),
    ],
    examples=[
        example(
            "Vertical stack",
            """
stack = Stack(direction="vertical", gap="md")
positions = stack.distribute(3, item_height=1.0, top=2.0)
            """,
            direction="vertical",
            gap="md"
        )
    ],
    tags=["layout", "stack", "flexbox"]
)
class Stack(Component):
    """
    Stack component for arranging items (like CSS Flexbox).

    Usage:
        # Vertical stack
        stack = Stack(direction="vertical", gap="lg")
        positions = stack.distribute(
            num_items=3,
            item_height=1.0,
            top=2.0
        )

        for pos in positions:
            component.render(slide, **pos)
    """

    def __init__(self,
                 direction: Literal["vertical", "horizontal"] = "vertical",
                 gap: Literal["none", "xs", "sm", "md", "lg", "xl"] = "md",
                 align: Literal["start", "center", "end", "stretch"] = "start",
                 theme: Optional[Dict[str, Any]] = None):
        super().__init__(theme)
        self.direction = direction
        self.gap = GAPS.get(gap, GAPS["md"])
        self.align = align

    def distribute(self,
                   num_items: int,
                   item_width: Optional[float] = None,
                   item_height: Optional[float] = None,
                   left: float = 0.5,
                   top: float = 1.5,
                   container_width: Optional[float] = None,
                   container_height: Optional[float] = None) -> List[Dict[str, float]]:
        """
        Distribute items in stack.

        Returns:
            List of position dicts for each item
        """
        positions = []

        if self.direction == "vertical":
            # Vertical stack
            current_top = top
            width = item_width or container_width or CONTENT_WIDTH

            for i in range(num_items):
                height = item_height or 1.0

                # Handle alignment
                if self.align == "center":
                    item_left = (SLIDE_WIDTH - width) / 2
                elif self.align == "end":
                    item_left = SLIDE_WIDTH - width - 0.5
                else:  # start or stretch
                    item_left = left

                positions.append({
                    'left': item_left,
                    'top': current_top,
                    'width': width,
                    'height': height
                })

                current_top += height + self.gap

        else:  # horizontal
            # Horizontal stack
            current_left = left
            height = item_height or container_height or 1.0

            for i in range(num_items):
                width = item_width or 2.0

                # Handle alignment
                if self.align == "center":
                    item_top = (SLIDE_HEIGHT - height) / 2
                elif self.align == "end":
                    item_top = SLIDE_HEIGHT - height - 0.5
                else:  # start or stretch
                    item_top = top

                positions.append({
                    'left': current_left,
                    'top': item_top,
                    'width': width,
                    'height': height
                })

                current_left += width + self.gap

        return positions

    def render_children(self,
                       slide,
                       children: List[Any],
                       left: float = 0.5,
                       top: float = 1.5,
                       item_width: Optional[float] = None,
                       item_height: Optional[float] = None) -> List[Any]:
        """
        Render a list of components in a stack layout.

        This is a convenience method that handles positioning automatically.

        Args:
            slide: PowerPoint slide
            children: List of component instances to render
            left: Starting left position
            top: Starting top position
            item_width: Width for all items (optional)
            item_height: Height for all items (optional)

        Returns:
            List of rendered shapes

        Example:
            stack = Stack(direction="vertical", gap="md")
            buttons = [
                Button("Save", "default"),
                Button("Cancel", "ghost")
            ]
            stack.render_children(slide, buttons, left=1, top=2)
        """
        positions = self.distribute(
            num_items=len(children),
            item_width=item_width,
            item_height=item_height,
            left=left,
            top=top
        )

        shapes = []
        for child, pos in zip(children, positions):
            # Render each child at its calculated position
            if hasattr(child, 'render'):
                shape = child.render(slide, **pos)
                shapes.append(shape)

        return shapes


@component(
    name="Spacer",
    category=ComponentCategory.LAYOUT,
    description="Invisible spacer for adding spacing between elements",
    props=[
        prop("size", "string", "Spacer size",
             options=["xs", "sm", "md", "lg", "xl", "2xl"],
             default="md"),
        prop("direction", "string", "Spacer direction",
             options=["vertical", "horizontal"],
             default="vertical"),
    ],
    tags=["layout", "spacer", "margin"]
)
class Spacer(Component):
    """
    Spacer component for adding space (like SwiftUI Spacer).

    Usage:
        # Get spacer height
        spacer = Spacer(size="lg", direction="vertical")
        spacing = spacer.get_size()
    """

    def __init__(self,
                 size: Literal["xs", "sm", "md", "lg", "xl", "2xl"] = "md",
                 direction: Literal["vertical", "horizontal"] = "vertical",
                 theme: Optional[Dict[str, Any]] = None):
        super().__init__(theme)
        self.size = size
        self.direction = direction

    def get_size(self) -> float:
        """Get the spacer size in inches."""
        size_map = {
            "xs": SPACING["4"],
            "sm": SPACING["6"],
            "md": SPACING["8"],
            "lg": SPACING["12"],
            "xl": SPACING["16"],
            "2xl": SPACING["24"],
        }
        return size_map.get(self.size, SPACING["8"])

    def render(self, slide, left: float = 0, top: float = 0):
        """Spacer doesn't render anything, just returns size."""
        size = self.get_size()
        if self.direction == "vertical":
            return {'height': size, 'width': 0}
        else:
            return {'width': size, 'height': 0}


@component(
    name="Divider",
    category=ComponentCategory.LAYOUT,
    description="Visual divider line for separating content",
    props=[
        prop("orientation", "string", "Divider orientation",
             options=["horizontal", "vertical"],
             default="horizontal"),
        prop("thickness", "number", "Line thickness in points", default=1),
        prop("color", "string", "Line color token", default="border.DEFAULT"),
    ],
    examples=[
        example(
            "Horizontal divider",
            """
divider = Divider(orientation="horizontal", thickness=1)
divider.render(slide, left=0.5, top=3.0, width=9.0)
            """,
            orientation="horizontal"
        )
    ],
    tags=["layout", "divider", "separator"]
)
class Divider(Component):
    """
    Divider line for visual separation.

    Usage:
        # Horizontal divider
        divider = Divider(orientation="horizontal")
        divider.render(slide, left=0.5, top=3.0, width=9.0)

        # Vertical divider
        divider = Divider(orientation="vertical")
        divider.render(slide, left=5.0, top=1.5, height=4.0)
    """

    def __init__(self,
                 orientation: Literal["horizontal", "vertical"] = "horizontal",
                 thickness: float = 1,
                 color: str = "border.DEFAULT",
                 theme: Optional[Dict[str, Any]] = None):
        super().__init__(theme)
        self.orientation = orientation
        self.thickness = thickness
        self.color = color

    def render(self, slide,
               left: float,
               top: float,
               width: Optional[float] = None,
               height: Optional[float] = None):
        """Render divider line."""
        if self.orientation == "horizontal":
            # Horizontal line
            line_width = width or CONTENT_WIDTH
            line_height = Pt(self.thickness)

            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(line_width),
                line_height
            )
        else:
            # Vertical line
            line_width = Pt(self.thickness)
            line_height = height or CONTENT_HEIGHT

            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left),
                Inches(top),
                line_width,
                Inches(line_height)
            )

        # Style line
        line.fill.solid()
        line.fill.fore_color.rgb = self.get_color(self.color)
        line.line.fill.background()

        return line
