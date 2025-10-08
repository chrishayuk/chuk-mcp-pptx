#!/usr/bin/env python3
"""
Async PowerPoint MCP Server using chuk-mcp-server

This server provides async MCP tools for creating and managing PowerPoint presentations
using the python-pptx library. It supports multiple presentations with optional
virtual filesystem integration for persistence and multi-server access.
"""
import asyncio
import json
import logging
import os
from pathlib import Path
from typing import List, Optional
import base64
import io

from chuk_mcp_server import ChukMCPServer
from pptx.util import Inches
from .presentation_manager import PresentationManager
from .slide_templates import (
    create_title_slide, create_content_slide, create_comparison_slide,
    create_key_metrics_slide, list_templates, list_color_schemes
)
# Text utilities now handled by tools/text.py via register_text_tools()
from .utilities.chart_utils import (
    add_chart, add_pie_chart, add_scatter_chart, add_data_table
)
# Shape utilities now available as components in components.core

# Import modular tools modules
from .chart_tools import register_chart_tools
from .tools.image import register_image_tools
from .tools.text import register_text_tools
from .inspection_tools import register_inspection_tools
from .tools.table import register_table_tools
from .tools.layout import register_layout_tools

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Create the MCP server instance
mcp = ChukMCPServer("chuk-mcp-pptx-async")

# Initialize presentation manager
# Check environment variable to enable VFS mode
USE_VFS = os.getenv("PPTX_USE_VFS", "false").lower() == "true"
VFS_BASE_PATH = os.getenv("PPTX_VFS_PATH", "vfs://presentations")

# Create a single manager instance for the server
manager = PresentationManager(use_vfs=USE_VFS, vfs_base_path=VFS_BASE_PATH)

# Register all modular tools
chart_tools = register_chart_tools(mcp, manager)
image_tools = register_image_tools(mcp, manager)
text_tools = register_text_tools(mcp, manager)
inspection_tools = register_inspection_tools(mcp, manager)
table_tools = register_table_tools(mcp, manager)
layout_tools = register_layout_tools(mcp, manager)
from .tools.shape import register_shape_tools
shape_tools = register_shape_tools(mcp, manager)

# Make tools available at module level for easier imports
if chart_tools:
    pptx_add_chart = chart_tools['pptx_add_chart']

if image_tools:
    pptx_add_image_slide = image_tools['pptx_add_image_slide']
    pptx_add_image = image_tools['pptx_add_image']
    pptx_add_background_image = image_tools['pptx_add_background_image']
    pptx_add_image_gallery = image_tools['pptx_add_image_gallery']
    pptx_add_image_with_caption = image_tools['pptx_add_image_with_caption']
    pptx_add_logo = image_tools['pptx_add_logo']
    pptx_replace_image = image_tools['pptx_replace_image']

if inspection_tools:
    pptx_inspect_slide = inspection_tools['pptx_inspect_slide']
    pptx_fix_slide_layout = inspection_tools['pptx_fix_slide_layout']
    pptx_analyze_presentation_layout = inspection_tools['pptx_analyze_presentation_layout']

if table_tools:
    pptx_add_data_table = table_tools['pptx_add_data_table']
    pptx_add_comparison_table = table_tools['pptx_add_comparison_table']
    pptx_update_table_cell = table_tools['pptx_update_table_cell']
    pptx_format_table = table_tools['pptx_format_table']

if layout_tools:
    pptx_list_layouts = layout_tools['pptx_list_layouts']
    pptx_add_slide_with_layout = layout_tools['pptx_add_slide_with_layout']
    pptx_customize_layout = layout_tools['pptx_customize_layout']
    pptx_apply_master_layout = layout_tools['pptx_apply_master_layout']
    pptx_duplicate_slide = layout_tools['pptx_duplicate_slide']
    pptx_reorder_slides = layout_tools['pptx_reorder_slides']

if shape_tools:
    pptx_add_shape = shape_tools['pptx_add_shape']
    pptx_add_arrow = shape_tools['pptx_add_arrow']
    pptx_add_smart_art = shape_tools['pptx_add_smart_art']
    pptx_add_code_block = shape_tools['pptx_add_code_block']
    pptx_apply_theme = shape_tools['pptx_apply_theme']
    pptx_list_themes = shape_tools['pptx_list_themes']

# Create backward compatibility references for chart and image tools
def _create_function_references():
    """Create references to chart and image functions for backward compatibility."""
    # Create a temporary MCP instance to get the function references
    temp_mcp = type('TempMCP', (), {})()
    temp_tools = {}
    
    def tool_decorator(func):
        temp_tools[func.__name__] = func
        return func
    
    temp_mcp.tool = tool_decorator
    
    # Register both chart and image tools to get all function references
    register_chart_tools(temp_mcp, manager)
    register_image_tools(temp_mcp, manager)
    
    # Make the functions available in this module's namespace
    for name, func in temp_tools.items():
        globals()[name] = func

_create_function_references()


@mcp.tool
async def pptx_create(name: str) -> str:
    """
    Create a new PowerPoint presentation.
    
    Creates a new blank presentation and sets it as the current active presentation.
    If VFS mode is enabled, the presentation is automatically saved to the virtual filesystem.
    
    Args:
        name: Unique name for the presentation (used for reference in other commands)
        
    Returns:
        Success message confirming presentation creation
        
    Example:
        await pptx_create(name="quarterly_report")
        # Returns: "Created presentation 'quarterly_report'"
    """
    # Run synchronous operation in thread pool
    return await asyncio.get_event_loop().run_in_executor(None, manager.create, name)


@mcp.tool
async def pptx_add_title_slide(
    title: str,
    subtitle: str = "",
    presentation: Optional[str] = None
) -> str:
    """
    Add a title slide to the current presentation.
    
    Creates a standard title slide with a main title and optional subtitle.
    This is typically used as the first slide in a presentation.
    
    Args:
        title: Main title text for the slide
        subtitle: Optional subtitle text (appears below the title)
        presentation: Name of presentation to add slide to (uses current if not specified)
        
    Returns:
        Success message confirming slide addition
        
    Example:
        await pptx_add_title_slide(
            title="Annual Report 2024",
            subtitle="Financial Results and Strategic Outlook"
        )
    """
    def _add_title_slide():
        prs = manager.get(presentation)
        if not prs:
            return "Error: No presentation found. Create one first with pptx_create()"
        
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        
        # Update in VFS if enabled
        manager.update(presentation)
        
        pres_name = presentation or manager.get_current_name()
        return f"Added title slide to '{pres_name}'"
    
    return await asyncio.get_event_loop().run_in_executor(None, _add_title_slide)


@mcp.tool
async def pptx_add_slide(
    title: str,
    content: List[str],
    presentation: Optional[str] = None
) -> str:
    """
    Add a content slide with title and bullet points.
    
    Creates a slide with a title and bulleted list. Perfect for agendas,
    key points, or any structured list content.
    
    Args:
        title: Title text for the slide
        content: List of strings, each becoming a bullet point
        presentation: Name of presentation to add slide to (uses current if not specified)
        
    Returns:
        Success message confirming slide addition
        
    Example:
        await pptx_add_slide(
            title="Project Milestones",
            content=[
                "Phase 1: Research completed",
                "Phase 2: Development in progress",
                "Phase 3: Testing scheduled for Q2"
            ]
        )
    """
    def _add_slide():
        prs = manager.get(presentation)
        if not prs:
            return "Error: No presentation found. Create one first with pptx_create()"
        
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = title
        
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            for idx, bullet in enumerate(content):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0  # First level bullet
        
        # Update in VFS if enabled
        manager.update(presentation)
        
        pres_name = presentation or manager.get_current_name()
        return f"Added content slide to '{pres_name}'"
    
    return await asyncio.get_event_loop().run_in_executor(None, _add_slide)


# Note: pptx_add_text_slide is now provided by text_tools.py
# The function is registered via register_text_tools()


@mcp.tool
async def pptx_save(
    path: str,
    presentation: Optional[str] = None
) -> str:
    """
    Save the presentation to a PowerPoint file.
    
    Saves the current or specified presentation to a .pptx file on disk.
    
    Args:
        path: File path where to save the .pptx file
        presentation: Name of presentation to save (uses current if not specified)
        
    Returns:
        Success message with the save path, or error message if save fails
        
    Example:
        await pptx_save(path="reports/quarterly_report.pptx")
        # Returns: "Saved presentation to: reports/quarterly_report.pptx"
    """
    def _save():
        prs = manager.get(presentation)
        if not prs:
            return "Error: No presentation found. Create one first with pptx_create()"
        
        try:
            prs.save(path)
            return f"Saved presentation to: {path}"
        except Exception as e:
            return f"Error: Failed to save presentation: {str(e)}"
    
    return await asyncio.get_event_loop().run_in_executor(None, _save)


@mcp.tool
async def pptx_export_base64(presentation: Optional[str] = None) -> str:
    """
    Export the presentation as a base64-encoded string.
    
    Exports the current or specified presentation as a base64 string that can be
    saved, transmitted, or imported later. Useful for transferring presentations
    between systems or storing in databases.
    
    Args:
        presentation: Name of presentation to export (uses current if not specified)
        
    Returns:
        JSON string containing presentation name, base64 data, and MIME type
        
    Example:
        result = await pptx_export_base64()
        # Returns JSON with structure:
        # {
        #   "presentation": "my_pres",
        #   "data": "UEsDBAoAAA...", 
        #   "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        # }
    """
    def _export():
        data = manager.export_base64(presentation)
        if not data:
            return json.dumps({"error": "No presentation found"})
        return json.dumps(data)
    
    return await asyncio.get_event_loop().run_in_executor(None, _export)


@mcp.tool
async def pptx_import_base64(data: str, name: str) -> str:
    """
    Import a presentation from a base64-encoded string.
    
    Imports a presentation from a base64 string and creates it with the given name.
    The imported presentation becomes the current active presentation.
    
    Args:
        data: Base64-encoded string of the .pptx file
        name: Name to give to the imported presentation
        
    Returns:
        Success message confirming import
        
    Example:
        await pptx_import_base64(
            data="UEsDBBQABgAIAAAAIQA...",
            name="imported_presentation"
        )
        # Returns: "Imported presentation 'imported_presentation'"
    """
    def _import():
        if manager.import_base64(data, name):
            return f"Imported presentation '{name}'"
        else:
            return "Error: Failed to import presentation"
    
    return await asyncio.get_event_loop().run_in_executor(None, _import)


@mcp.tool
async def pptx_list() -> str:
    """
    List all presentations currently in memory.
    
    Returns a JSON array of presentation names that are currently loaded.
    Useful for managing multiple presentations in a session.
    
    Returns:
        Formatted list of presentations with slide counts and status
        
    Example:
        presentations = await pptx_list()
        # Returns formatted list of presentations
    """
    def _list():
        if not manager._presentations:
            return "No presentations currently in memory"
        
        lines = ["Current presentations:"]
        for name in manager._presentations:
            prs = manager.get(name)
            if prs:
                slide_count = len(prs.slides)
                is_current = name == manager.get_current_name()
                status = " (current)" if is_current else ""
                lines.append(f"  - {name}: {slide_count} slides{status}")
        
        return "\n".join(lines)
    
    return await asyncio.get_event_loop().run_in_executor(None, _list)


@mcp.tool
async def pptx_switch(name: str) -> str:
    """
    Switch to a different presentation.
    
    Changes the current active presentation to the specified one.
    All subsequent operations will affect this presentation unless
    explicitly specified otherwise.
    
    Args:
        name: Name of the presentation to switch to
        
    Returns:
        Success message confirming the switch
        
    Example:
        await pptx_switch(name="sales_pitch")
        # Returns: "Switched to presentation 'sales_pitch'"
    """
    def _switch():
        if name not in manager._presentations:
            return f"Error: Presentation '{name}' not found"
        manager._current_presentation = name
        return f"Switched to presentation '{name}'"
    
    return await asyncio.get_event_loop().run_in_executor(None, _switch)


@mcp.tool
async def pptx_delete(name: str) -> str:
    """
    Delete a presentation from memory.
    
    Removes the specified presentation from memory. If it's the current
    presentation, you'll need to switch to or create another one.
    
    Args:
        name: Name of the presentation to delete
        
    Returns:
        Success message confirming deletion
        
    Example:
        await pptx_delete(name="old_presentation")
        # Returns: "Deleted presentation 'old_presentation'"
    """
    def _delete():
        if name not in manager._presentations:
            return f"Error: Presentation '{name}' not found"
        
        del manager._presentations[name]
        
        # Clear current if it was deleted
        if manager._current_presentation == name:
            manager._current_presentation = None
        
        # Delete from VFS if enabled
        if manager.use_vfs:
            vfs_path = f"{manager.vfs_base_path}/{name}.pptx"
            # In real implementation, would delete from VFS here
        
        return f"Deleted presentation '{name}'"
    
    return await asyncio.get_event_loop().run_in_executor(None, _delete)


@mcp.tool
async def pptx_get_info(presentation: Optional[str] = None) -> str:
    """
    Get information about a presentation.
    
    Returns detailed information about the specified presentation including
    number of slides, slide titles, and slide content summaries.
    
    Args:
        presentation: Name of presentation to get info for (uses current if not specified)
        
    Returns:
        JSON object with presentation information
        
    Example:
        info = await pptx_get_info()
        # Returns JSON with structure:
        # {
        #   "name": "quarterly_report",
        #   "slides": 10,
        #   "slide_details": [
        #     {"index": 0, "title": "Q4 Report", "shapes": 2},
        #     ...
        #   ]
        # }
    """
    def _get_info():
        prs = manager.get(presentation)
        if not prs:
            return json.dumps({"error": "No presentation found"})
        
        info = {
            "name": presentation or manager.get_current_name(),
            "slides": len(prs.slides),
            "slide_details": []
        }
        
        for idx, slide in enumerate(prs.slides):
            slide_info = {
                "index": idx,
                "title": slide.shapes.title.text if slide.shapes.title else "No title",
                "shapes": len(slide.shapes)
            }
            info["slide_details"].append(slide_info)
        
        return json.dumps(info, indent=2)
    
    return await asyncio.get_event_loop().run_in_executor(None, _get_info)


# Enhanced tools using template utilities

@mcp.tool
async def pptx_create_title_slide(
    title: str,
    subtitle: str = "",
    author: str = "",
    date: str = "",
    company: str = "",
    color_scheme: str = "modern_blue",
    presentation: Optional[str] = None
) -> str:
    """
    Create a professional title slide with enhanced styling.
    
    Creates a title slide with professional formatting, optional author information,
    and customizable color schemes. Perfect for presentation openings.
    
    Args:
        title: Main title of the presentation
        subtitle: Optional subtitle or tagline
        author: Optional author name
        date: Optional date (e.g., "December 2024")
        company: Optional company or organization name
        color_scheme: Color scheme to use (modern_blue, corporate_gray, elegant_green, warm_orange)
        presentation: Name of presentation to add slide to (uses current if not specified)
        
    Returns:
        Success message confirming slide creation
        
    Example:
        await pptx_create_title_slide(
            title="Strategic Plan 2025",
            subtitle="Driving Innovation and Growth",
            author="Jane Smith",
            date="January 2025",
            company="Tech Corp",
            color_scheme="corporate_gray"
        )
    """
    def _create_title():
        prs = manager.get(presentation)
        if not prs:
            return "Error: No presentation found. Create one first with pptx_create()"
        
        slide = create_title_slide(
            prs, title, subtitle, author, date, color_scheme
        )
        
        # Update in VFS if enabled
        manager.update(presentation)
        
        slide_idx = len(prs.slides) - 1
        return f"Created title slide at index {slide_idx}"
    
    return await asyncio.get_event_loop().run_in_executor(None, _create_title)


@mcp.tool
async def pptx_create_comparison_slide(
    title: str,
    left_title: str,
    left_items: List[str],
    right_title: str,
    right_items: List[str],
    color_scheme: str = "modern_blue",
    presentation: Optional[str] = None
) -> str:
    """
    Create a two-column comparison slide.
    
    Creates a slide with two columns for comparing items, features, pros/cons, etc.
    Each column has its own title and list of items.
    
    Args:
        title: Main slide title
        left_title: Title for the left column
        left_items: List of items for the left column
        right_title: Title for the right column
        right_items: List of items for the right column
        color_scheme: Color scheme to use
        presentation: Name of presentation to add slide to (uses current if not specified)
        
    Returns:
        Success message confirming slide creation
        
    Example:
        await pptx_create_comparison_slide(
            title="Solution Comparison",
            left_title="Option A",
            left_items=["Lower cost", "Faster implementation", "Local support"],
            right_title="Option B",
            right_items=["More features", "Better scalability", "Cloud-based"],
            color_scheme="modern_blue"
        )
    """
    def _create_comparison():
        prs = manager.get(presentation)
        if not prs:
            return "Error: No presentation found. Create one first with pptx_create()"
        
        slide = create_comparison_slide(
            prs, title, left_title, left_items, right_title, right_items, color_scheme
        )
        
        # Update in VFS if enabled
        manager.update(presentation)
        
        slide_idx = len(prs.slides) - 1
        return f"Created comparison slide at index {slide_idx}"
    
    return await asyncio.get_event_loop().run_in_executor(None, _create_comparison)


@mcp.tool
async def pptx_create_metrics_slide(
    title: str,
    metrics: List[dict],
    color_scheme: str = "modern_blue",
    presentation: Optional[str] = None
) -> str:
    """
    Create a key metrics dashboard slide.
    
    Creates a slide displaying key metrics in a visually appealing grid layout.
    Each metric shows a large value and a description label.
    
    Args:
        title: Slide title
        metrics: List of metric dictionaries with 'value' and 'label' keys
        color_scheme: Color scheme to use
        presentation: Name of presentation to add slide to (uses current if not specified)
        
    Returns:
        Success message confirming slide creation
        
    Example:
        await pptx_create_metrics_slide(
            title="Q4 Performance Metrics",
            metrics=[
                {"value": "$2.5M", "label": "Revenue"},
                {"value": "147%", "label": "Growth"},
                {"value": "92", "label": "NPS Score"},
                {"value": "45ms", "label": "Response Time"}
            ],
            color_scheme="corporate_gray"
        )
    """
    def _create_metrics():
        prs = manager.get(presentation)
        if not prs:
            return "Error: No presentation found. Create one first with pptx_create()"
        
        slide = create_key_metrics_slide(prs, title, metrics, color_scheme)
        
        # Update in VFS if enabled
        manager.update(presentation)
        
        slide_idx = len(prs.slides) - 1
        return f"Created metrics slide at index {slide_idx}"
    
    return await asyncio.get_event_loop().run_in_executor(None, _create_metrics)


# Additional async tools for shapes, text extraction, etc.

# Note: pptx_extract_all_text is now provided by text_tools.py
# The function is registered via register_text_tools()

# Note: pptx_add_data_table is now provided by table_tools.py with layout validation
# The function is registered via register_table_tools()


@mcp.tool
async def pptx_list_templates() -> str:
    """
    List all available slide templates and color schemes.
    
    Returns information about available slide templates and color schemes
    that can be used when creating slides.
    
    Returns:
        JSON object with templates and color schemes
        
    Example:
        templates = await pptx_list_templates()
        # Returns JSON with available templates and color schemes
    """
    def _list():
        templates = list_templates()
        schemes = list_color_schemes()
        
        return json.dumps({
            "templates": templates,
            "color_schemes": schemes
        }, indent=2)
    
    return await asyncio.get_event_loop().run_in_executor(None, _list)


# Run the server
if __name__ == "__main__":
    import sys
    
    # When run directly, default to stdio mode for testing
    print(f"Starting PowerPoint MCP Server...", file=sys.stderr)
    print(f"VFS Mode: {USE_VFS}", file=sys.stderr)
    if USE_VFS:
        print(f"VFS Path: {VFS_BASE_PATH}", file=sys.stderr)
    
    # Run in stdio mode when executed directly
    mcp.run(stdio=True)