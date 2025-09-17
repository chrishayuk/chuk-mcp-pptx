"""
MCP tools for component-based PowerPoint creation.
Provides async tools that use the component system.
"""

import asyncio
from typing import Dict, Any, Optional, List
from pptx.util import Inches

from ..components import (
    Button, Card, MetricCard, FeatureCard,
    CodeBlock, Terminal
)
from ..components.chart import (
    BarChart, LineChart, PieChart, ScatterChart, AreaChart
)
from ..themes.theme_manager import ThemeManager


def register_component_tools(mcp, manager):
    """
    Register component-based tools with the MCP server.
    
    Args:
        mcp: ChukMCPServer instance
        manager: PresentationManager instance
    
    Returns:
        Dictionary of registered tools
    """
    tools = {}
    theme_manager = ThemeManager()
    
    @mcp.tool
    async def pptx_add_button(
        slide_index: int,
        text: str,
        left: float,
        top: float,
        variant: str = "primary",
        size: str = "md",
        width: Optional[float] = None,
        height: Optional[float] = None,
        theme: Optional[str] = None
    ) -> str:
        """
        Add a themed button component to a slide.
        
        Components provide consistent, theme-aware UI elements similar to shadcn/ui.
        
        Args:
            slide_index: Index of the slide (0-based)
            text: Button text
            left: Left position in inches
            top: Top position in inches
            variant: Button style (primary, secondary, outline, ghost, destructive)
            size: Button size (sm, md, lg)
            width: Optional width override in inches
            height: Optional height override in inches
            theme: Theme name to use (uses slide theme if not specified)
        
        Returns:
            Success message
        
        Example:
            await pptx_add_button(
                slide_index=0,
                text="Get Started",
                left=4.0,
                top=3.0,
                variant="primary",
                size="lg"
            )
        """
        def _add_button():
            prs = manager.get_current()
            if not prs:
                raise ValueError("No active presentation")
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide {slide_index} does not exist")
            
            slide = prs.slides[slide_index]
            
            # Get theme
            theme_obj = theme_manager.get_theme(theme) if theme else None
            if not theme_obj:
                theme_obj = theme_manager.get_theme("dark")
            
            # Create and render button
            button = Button(text, variant, size, theme_obj.__dict__)
            button.render(slide, left, top, width, height)
            
            return f"Added {variant} button '{text}' to slide {slide_index}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_button)
    
    @mcp.tool
    async def pptx_add_card(
        slide_index: int,
        left: float,
        top: float,
        width: float = 3.0,
        height: float = 2.0,
        title: Optional[str] = None,
        description: Optional[str] = None,
        variant: str = "default",
        theme: Optional[str] = None
    ) -> str:
        """
        Add a themed card component to a slide.
        
        Cards are container components for grouping related content.
        
        Args:
            slide_index: Index of the slide (0-based)
            left: Left position in inches
            top: Top position in inches
            width: Card width in inches
            height: Card height in inches
            title: Optional card title
            description: Optional card description
            variant: Card style (default, bordered, elevated)
            theme: Theme name to use
        
        Returns:
            Success message
        
        Example:
            await pptx_add_card(
                slide_index=0,
                left=2.0,
                top=2.0,
                title="Features",
                description="Key capabilities of our system",
                variant="elevated"
            )
        """
        def _add_card():
            prs = manager.get_current()
            if not prs:
                raise ValueError("No active presentation")
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide {slide_index} does not exist")
            
            slide = prs.slides[slide_index]
            
            # Get theme
            theme_obj = theme_manager.get_theme(theme) if theme else None
            if not theme_obj:
                theme_obj = theme_manager.get_theme("dark")
            
            # Create and render card
            card = Card(title, description, variant, theme_obj.__dict__)
            card.render(slide, left, top, width, height)
            
            return f"Added {variant} card to slide {slide_index}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_card)
    
    @mcp.tool
    async def pptx_add_metric_card(
        slide_index: int,
        label: str,
        value: str,
        left: float,
        top: float,
        change: Optional[str] = None,
        trend: Optional[str] = None,
        width: float = 2.0,
        height: float = 1.5,
        theme: Optional[str] = None
    ) -> str:
        """
        Add a metric/KPI card component to a slide.
        
        Metric cards display key performance indicators with optional trends.
        
        Args:
            slide_index: Index of the slide (0-based)
            label: Metric label (e.g., "Revenue")
            value: Metric value (e.g., "$1.2M")
            left: Left position in inches
            top: Top position in inches
            change: Optional change value (e.g., "+12%")
            trend: Trend direction (up, down, neutral)
            width: Card width in inches
            height: Card height in inches
            theme: Theme name to use
        
        Returns:
            Success message
        
        Example:
            await pptx_add_metric_card(
                slide_index=0,
                label="Monthly Revenue",
                value="$45.2K",
                left=2.0,
                top=2.0,
                change="+12%",
                trend="up"
            )
        """
        def _add_metric():
            prs = manager.get_current()
            if not prs:
                raise ValueError("No active presentation")
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide {slide_index} does not exist")
            
            slide = prs.slides[slide_index]
            
            # Get theme
            theme_obj = theme_manager.get_theme(theme) if theme else None
            if not theme_obj:
                theme_obj = theme_manager.get_theme("dark")
            
            # Create and render metric card
            metric = MetricCard(label, value, change, trend, theme_obj.__dict__)
            metric.render(slide, left, top, width, height)
            
            return f"Added metric card '{label}' to slide {slide_index}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_metric)
    
    @mcp.tool
    async def pptx_add_code_component(
        slide_index: int,
        code: str,
        language: str,
        left: float,
        top: float,
        width: float = 6.0,
        height: float = 3.0,
        show_line_numbers: bool = False,
        theme: Optional[str] = None
    ) -> str:
        """
        Add a themed code block component to a slide.
        
        Code blocks display syntax-highlighted code with language labels.
        
        Args:
            slide_index: Index of the slide (0-based)
            code: Code content
            language: Programming language (python, javascript, etc.)
            left: Left position in inches
            top: Top position in inches
            width: Code block width in inches
            height: Code block height in inches
            show_line_numbers: Whether to show line numbers
            theme: Theme name to use
        
        Returns:
            Success message
        
        Example:
            await pptx_add_code_component(
                slide_index=0,
                code="def hello():\\n    print('Hello, World!')",
                language="python",
                left=2.0,
                top=2.0,
                show_line_numbers=True
            )
        """
        def _add_code():
            prs = manager.get_current()
            if not prs:
                raise ValueError("No active presentation")
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide {slide_index} does not exist")
            
            slide = prs.slides[slide_index]
            
            # Get theme
            theme_obj = theme_manager.get_theme(theme) if theme else None
            if not theme_obj:
                theme_obj = theme_manager.get_theme("dark")
            
            # Create and render code block
            code_block = CodeBlock(code, language, show_line_numbers, theme_obj.__dict__)
            code_block.render(slide, left, top, width, height)
            
            return f"Added {language} code block to slide {slide_index}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_code)
    
    @mcp.tool
    async def pptx_add_chart_component(
        slide_index: int,
        chart_type: str,
        title: Optional[str],
        categories: List[str],
        series: Dict[str, List[float]],
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.5,
        theme: Optional[str] = None
    ) -> str:
        """
        Add a themed chart component to a slide.
        
        Charts use theme colors automatically for consistent visualization.
        
        Args:
            slide_index: Index of the slide (0-based)
            chart_type: Type of chart (bar, line, pie, area, scatter)
            title: Chart title
            categories: Category labels
            series: Data series dictionary
            left: Left position in inches
            top: Top position in inches
            width: Chart width in inches
            height: Chart height in inches
            theme: Theme name to use
        
        Returns:
            Success message
        
        Example:
            await pptx_add_chart_component(
                slide_index=0,
                chart_type="bar",
                title="Q4 Sales",
                categories=["Oct", "Nov", "Dec"],
                series={"Revenue": [100, 120, 140], "Profit": [20, 25, 30]}
            )
        """
        async def _add_chart():
            prs = manager.get_current()
            if not prs:
                raise ValueError("No active presentation")
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide {slide_index} does not exist")
            
            slide = prs.slides[slide_index]
            
            # Get theme
            theme_obj = theme_manager.get_theme(theme) if theme else None
            if not theme_obj:
                theme_obj = theme_manager.get_theme("dark")
            
            # Create appropriate chart component
            if chart_type.lower() == "bar":
                chart = BarChart(title=title, categories=categories, series=series, theme=theme_obj.__dict__)
            elif chart_type.lower() == "line":
                chart = LineChart(title=title, categories=categories, series=series, theme=theme_obj.__dict__)
            elif chart_type.lower() == "pie":
                # For pie charts, use first series
                values = list(series.values())[0] if series else []
                chart = PieChart(title=title, categories=categories, values=values, theme=theme_obj.__dict__)
            elif chart_type.lower() == "area":
                chart = AreaChart(title=title, categories=categories, series=series, theme=theme_obj.__dict__)
            else:
                raise ValueError(f"Unsupported chart type: {chart_type}")
            
            # Render chart
            await chart.render(slide, left, top, width, height)
            
            return f"Added {chart_type} chart to slide {slide_index}"
        
        return await _add_chart()
    
    @mcp.tool
    async def pptx_apply_component_theme(
        slide_index: int,
        theme: str
    ) -> str:
        """
        Apply a component theme to a slide.
        
        Themes provide consistent color schemes and styling across all components.
        Available themes:
        - Dark themes: dark, dark-blue, dark-violet, dark-green, dark-orange, dark-purple
        - Light themes: light, light-blue, light-violet, light-green, light-warm
        - Special themes: cyberpunk, sunset, ocean, aurora, minimal, corporate
        
        Args:
            slide_index: Index of the slide (0-based)
            theme: Theme name to apply
        
        Returns:
            Success message
        
        Example:
            await pptx_apply_component_theme(
                slide_index=0,
                theme="dark-violet"
            )
        """
        def _apply_theme():
            prs = manager.get_current()
            if not prs:
                raise ValueError("No active presentation")
            
            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide {slide_index} does not exist")
            
            slide = prs.slides[slide_index]
            
            # Apply theme
            theme_obj = theme_manager.get_theme(theme)
            if not theme_obj:
                raise ValueError(f"Theme '{theme}' not found")
            
            theme_obj.apply_to_slide(slide)
            
            return f"Applied '{theme}' theme to slide {slide_index}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _apply_theme)
    
    @mcp.tool
    async def pptx_list_component_themes() -> str:
        """
        List all available component themes.
        
        Returns:
            List of available theme names
        
        Example:
            themes = await pptx_list_component_themes()
        """
        themes = theme_manager.list_themes()
        return f"Available themes: {', '.join(themes)}"
    
    # Store tools for return
    tools['pptx_add_button'] = pptx_add_button
    tools['pptx_add_card'] = pptx_add_card
    tools['pptx_add_metric_card'] = pptx_add_metric_card
    tools['pptx_add_code_component'] = pptx_add_code_component
    tools['pptx_add_chart_component'] = pptx_add_chart_component
    tools['pptx_apply_component_theme'] = pptx_apply_component_theme
    tools['pptx_list_component_themes'] = pptx_list_component_themes
    
    return tools