"""
Shape and SmartArt Tools for PowerPoint MCP Server

Provides async MCP tools for creating shapes, connectors, and SmartArt-like diagrams.
"""
import asyncio
from typing import List, Optional, Tuple, Dict
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def register_shape_tools(mcp, manager):
    """Register all shape and SmartArt tools with the MCP server."""

    from ..components.core import Shape, Connector, ProcessFlow, CycleDiagram, HierarchyDiagram
    from ..layout.helpers import validate_position
    from ..components.code import CodeBlock
    from ..themes.theme_manager import ThemeManager
    
    @mcp.tool
    async def pptx_add_shape(
        slide_index: int,
        shape_type: str,
        left: float = 2.0,
        top: float = 2.0,
        width: float = 3.0,
        height: float = 2.0,
        text: Optional[str] = None,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
        line_width: float = 1.0,
        presentation: Optional[str] = None
    ) -> str:
        """
        Add a shape to a slide with customizable properties.
        
        Creates various shapes like rectangles, circles, arrows, stars, etc.
        
        Args:
            slide_index: Index of the slide to add shape to (0-based)
            shape_type: Type of shape. Options:
                - "rectangle", "rounded_rectangle"
                - "oval", "diamond"
                - "triangle", "arrow"
                - "star", "hexagon"
                - "chevron", "plus"
                - "callout"
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Optional text to display inside shape
            fill_color: Fill color as hex string (e.g., "#FF0000")
            line_color: Border color as hex string
            line_width: Border width in points
            presentation: Name of presentation (uses current if not specified)
            
        Returns:
            Success message confirming shape addition
            
        Example:
            await pptx_add_shape(
                slide_index=1,
                shape_type="star",
                left=3.0,
                top=2.0,
                width=2.0,
                height=2.0,
                text="Success!",
                fill_color="#FFD700"
            )
        """
        def _add_shape():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found. Create one first with pptx_create()"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            # Remove content placeholders that might interfere with shapes
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type'):
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, 'placeholder_format'):
                            if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                                continue
                        # Mark content placeholders for removal
                        placeholders_to_remove.append(shape)
            
            # Remove the placeholders
            for placeholder in placeholders_to_remove:
                sp = placeholder._element
                sp.getparent().remove(sp)
            
            # Validate position
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )

            try:
                # Get theme for shape
                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_default_theme()

                # Create shape using component
                shape_comp = Shape(
                    shape_type=shape_type,
                    text=text,
                    fill_color=fill_color,  # Component handles hex colors
                    line_color=line_color,
                    line_width=line_width,
                    theme=theme_obj
                )
                shape = shape_comp.render(
                    slide,
                    validated_left,
                    validated_top,
                    validated_width,
                    validated_height
                )
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                # Report if position was adjusted
                position_note = ""
                if (validated_left != left or validated_top != top or 
                    validated_width != width or validated_height != height):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"
                
                return f"Added {shape_type} shape to slide {slide_index}{position_note}"
                
            except Exception as e:
                return f"Error adding shape: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_shape)
    
    @mcp.tool
    async def pptx_add_arrow(
        slide_index: int,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        connector_type: str = "straight",
        line_color: Optional[str] = "#000000",
        line_width: float = 2.0,
        arrow_start: bool = False,
        arrow_end: bool = True,
        presentation: Optional[str] = None
    ) -> str:
        """
        Add an arrow or connector line to a slide.
        
        Creates connectors between points with optional arrowheads.
        
        Args:
            slide_index: Index of the slide to add arrow to
            start_x: Starting X position in inches
            start_y: Starting Y position in inches
            end_x: Ending X position in inches
            end_y: Ending Y position in inches
            connector_type: Type of connector ("straight", "elbow", "curved")
            line_color: Line color as hex string
            line_width: Line width in points
            arrow_start: Whether to add arrowhead at start
            arrow_end: Whether to add arrowhead at end
            presentation: Name of presentation
            
        Returns:
            Success message confirming arrow addition
            
        Example:
            await pptx_add_arrow(
                slide_index=1,
                start_x=2.0,
                start_y=2.0,
                end_x=5.0,
                end_y=3.0,
                connector_type="straight",
                arrow_end=True
            )
        """
        def _add_arrow():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            try:
                # Get theme for connector
                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_default_theme()

                # Create connector using component
                connector_comp = Connector(
                    start_x=start_x,
                    start_y=start_y,
                    end_x=end_x,
                    end_y=end_y,
                    connector_type=connector_type,
                    line_color=line_color,  # Component handles hex colors
                    line_width=line_width,
                    arrow_start=arrow_start,
                    arrow_end=arrow_end,
                    theme=theme_obj
                )
                connector = connector_comp.render(slide)
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                return f"Added {connector_type} arrow to slide {slide_index}"
                
            except Exception as e:
                return f"Error adding arrow: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_arrow)
    
    @mcp.tool
    async def pptx_add_smart_art(
        slide_index: int,
        art_type: str,
        items: List[str],
        title: Optional[str] = None,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 3.0,
        color_scheme: str = "modern_blue",
        presentation: Optional[str] = None
    ) -> str:
        """
        Add a SmartArt-style diagram to a slide.
        
        Creates professional diagrams like process flows, cycles, hierarchies, etc.
        
        Args:
            slide_index: Index of the slide to add SmartArt to
            art_type: Type of SmartArt diagram:
                - "process" - Sequential process flow
                - "cycle" - Circular/cyclical process
                - "hierarchy" - Organizational hierarchy
                - "list" - Bulleted list with shapes
                - "relationship" - Relationship diagram
                - "pyramid" - Pyramid hierarchy
            items: List of text items for the diagram
            title: Optional title for the diagram
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            color_scheme: Color scheme ("modern_blue", "corporate_gray", "warm_orange")
            presentation: Name of presentation
            
        Returns:
            Success message confirming SmartArt addition
            
        Example:
            await pptx_add_smart_art(
                slide_index=2,
                art_type="process",
                items=["Research", "Design", "Develop", "Test", "Deploy"],
                title="Development Process",
                color_scheme="modern_blue"
            )
        """
        def _add_smart():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            # Remove content placeholders that might interfere with SmartArt
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type'):
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, 'placeholder_format'):
                            if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                                continue
                        # Mark content placeholders for removal
                        placeholders_to_remove.append(shape)
            
            # Remove the placeholders
            for placeholder in placeholders_to_remove:
                sp = placeholder._element
                sp.getparent().remove(sp)
            
            # Validate position
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )
            
            # Add title if provided
            if title:
                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_default_theme()

                title_comp = Shape(
                    shape_type="rectangle",
                    text=title,
                    fill_color="transparent",
                    line_color="transparent",
                    theme=theme_obj
                )
                title_shape = title_comp.render(
                    slide,
                    validated_left,
                    validated_top - 0.5,
                    validated_width,
                    0.4
                )
                # Adjust top position for the SmartArt
                validated_top += 0.2
                validated_height -= 0.7
            
            try:
                # Get theme for SmartArt
                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_default_theme()

                # Map art_type to component
                art_components = {
                    "process": ProcessFlow,
                    "cycle": CycleDiagram,
                    "hierarchy": HierarchyDiagram,
                }

                component_class = art_components.get(art_type)
                if not component_class:
                    return f"Error: Unsupported art type '{art_type}'. Supported: {', '.join(art_components.keys())}"

                # Create SmartArt using component
                smart_art_comp = component_class(
                    items=items,
                    theme=theme_obj
                )
                shapes = smart_art_comp.render(
                    slide,
                    validated_left,
                    validated_top,
                    validated_width,
                    validated_height
                )
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                # Report if position was adjusted
                position_note = ""
                if (validated_left != left or validated_top != top or 
                    validated_width != width or validated_height != height):
                    position_note = f" (position adjusted)"
                
                return f"Added {art_type} SmartArt with {len(items)} items to slide {slide_index}{position_note}"
                
            except Exception as e:
                return f"Error adding SmartArt: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_smart)
    
    @mcp.tool
    async def pptx_add_code_block(
        slide_index: int,
        code: str,
        language: str = "python",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 3.0,
        theme: str = "dark_modern",
        presentation: Optional[str] = None
    ) -> str:
        """
        Add a code block to a slide with syntax highlighting appearance.
        
        Creates a formatted code block with monospace font and theme colors.
        
        Args:
            slide_index: Index of the slide to add code to (0-based)
            code: The code content to display
            language: Programming language (for label)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            theme: Theme to use for code block styling
            presentation: Name of presentation (uses current if not specified)
            
        Returns:
            Success message confirming code block addition
            
        Example:
            await pptx_add_code_block(
                slide_index=1,
                code="def hello_world():\\n    print('Hello, World!')",
                language="python",
                theme="dark_purple"
            )
        """
        def _add_code():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            try:
                # Remove content placeholders
                from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
                placeholders_to_remove = []
                for shape in slide.shapes:
                    if hasattr(shape, 'shape_type'):
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            if hasattr(shape, 'placeholder_format'):
                                if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                                    continue
                            placeholders_to_remove.append(shape)
                
                for shape in placeholders_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)
                
                # Add code block using CodeBlock component
                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_theme(theme) if theme else None

                code_component = CodeBlock(
                    code_text=code,
                    language=language,
                    theme=theme_obj
                )
                code_shape = code_component.render(
                    slide,
                    left=left,
                    top=top,
                    width=width,
                    height=height
                )

                # Update in VFS if enabled
                manager.update(presentation)

                return f"Added {language} code block to slide {slide_index}"
                
            except Exception as e:
                return f"Error adding code block: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_code)
    
    @mcp.tool
    async def pptx_apply_theme(
        slide_index: Optional[int] = None,
        theme: str = "default-light",
        presentation: Optional[str] = None
    ) -> str:
        """
        Apply a beautiful theme to slides.

        Available themes: Use pptx_list_themes() to see all available themes.

        Args:
            slide_index: Index of slide to theme (None for all slides)
            theme: Name of theme to apply (e.g., "default-light", "ocean-dark", etc.)
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming theme application
        """
        def _apply_theme():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"

            theme_manager = ThemeManager()
            available_themes = theme_manager.list_themes()

            if theme not in available_themes:
                return f"Error: Unknown theme '{theme}'. Available: {', '.join(available_themes[:10])}"

            try:
                theme_obj = theme_manager.get_theme(theme)

                if slide_index is not None:
                    if slide_index >= len(prs.slides):
                        return f"Error: Slide index {slide_index} out of range"
                    slides = [prs.slides[slide_index]]
                else:
                    slides = prs.slides

                # Apply background color from theme
                for slide in slides:
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    # Get background color from theme
                    bg_color_hex = theme_obj.colors.get("background", {}).get("DEFAULT", "#FFFFFF")
                    if bg_color_hex.startswith("#"):
                        bg_color_hex = bg_color_hex[1:]
                    from pptx.dml.color import RGBColor
                    fill.fore_color.rgb = RGBColor(
                        int(bg_color_hex[0:2], 16),
                        int(bg_color_hex[2:4], 16),
                        int(bg_color_hex[4:6], 16)
                    )

                # Update in VFS if enabled
                manager.update(presentation)

                slide_msg = f"slide {slide_index}" if slide_index is not None else "all slides"
                return f"Applied {theme} theme to {slide_msg}"

            except Exception as e:
                return f"Error applying theme: {str(e)}"

        return await asyncio.get_event_loop().run_in_executor(None, _apply_theme)
    
    @mcp.tool
    async def pptx_list_themes() -> str:
        """
        List all available themes with descriptions.

        Returns:
            List of available themes and their characteristics
        """
        theme_manager = ThemeManager()
        themes = theme_manager.list_themes()

        theme_list = []
        for theme_name in themes:
            theme_obj = theme_manager.get_theme(theme_name)
            mode = theme_obj.mode
            primary = theme_obj.colors.get("primary", {}).get("DEFAULT", "N/A")
            theme_list.append(f"• {theme_name} ({mode}): Primary: {primary}")

        return "Available themes:\n" + "\n".join(theme_list)
    
    # Return tools for external access
    return {
        'pptx_add_shape': pptx_add_shape,
        'pptx_add_arrow': pptx_add_arrow,
        'pptx_add_smart_art': pptx_add_smart_art,
        'pptx_add_code_block': pptx_add_code_block,
        'pptx_apply_theme': pptx_apply_theme,
        'pptx_list_themes': pptx_list_themes
    }