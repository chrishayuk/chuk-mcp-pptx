"""
Button component for PowerPoint presentations.
Provides various button styles and variants.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from .base import Component


class Button(Component):
    """
    Button component with multiple variants and styles.
    
    Variants:
    - default: Standard button
    - primary: Primary action button
    - secondary: Secondary action button
    - outline: Outlined button
    - ghost: Minimal button
    - destructive: Destructive action button
    """
    
    def __init__(self, 
                 text: str,
                 variant: str = "default",
                 size: str = "md",
                 theme: Optional[Dict[str, Any]] = None):
        """
        Initialize button component.
        
        Args:
            text: Button text
            variant: Button variant
            size: Button size (sm, md, lg)
            theme: Theme configuration
        """
        super().__init__(theme)
        self.text = text
        self.variant = variant
        self.size = size
    
    def get_size_config(self) -> Dict[str, Any]:
        """Get size configuration for button."""
        sizes = {
            "sm": {
                "width": 1.5,
                "height": 0.4,
                "font_size": 12,
                "padding": self.get_padding("sm"),
            },
            "md": {
                "width": 2.0,
                "height": 0.5,
                "font_size": 14,
                "padding": self.get_padding("md"),
            },
            "lg": {
                "width": 2.5,
                "height": 0.6,
                "font_size": 16,
                "padding": self.get_padding("lg"),
            },
        }
        return sizes.get(self.size, sizes["md"])
    
    def get_variant_colors(self) -> Dict[str, Any]:
        """Get colors for button variant."""
        variants = {
            "default": {
                "bg": self.get_color("background.secondary"),
                "fg": self.get_color("foreground.DEFAULT"),
                "border": self.get_color("border.DEFAULT"),
                "hover_bg": self.get_color("background.tertiary"),
            },
            "primary": {
                "bg": self.get_color("primary.DEFAULT"),
                "fg": self.get_color("primary.foreground"),
                "border": self.get_color("primary.DEFAULT"),
                "hover_bg": self.get_color("primary.hover"),
            },
            "secondary": {
                "bg": self.get_color("secondary.DEFAULT"),
                "fg": self.get_color("secondary.foreground"),
                "border": self.get_color("border.DEFAULT"),
                "hover_bg": self.get_color("secondary.hover"),
            },
            "outline": {
                "bg": None,  # Transparent
                "fg": self.get_color("foreground.DEFAULT"),
                "border": self.get_color("border.DEFAULT"),
                "hover_bg": self.get_color("background.secondary"),
            },
            "ghost": {
                "bg": None,  # Transparent
                "fg": self.get_color("foreground.DEFAULT"),
                "border": None,  # No border
                "hover_bg": self.get_color("background.secondary"),
            },
            "destructive": {
                "bg": self.get_color("destructive.DEFAULT"),
                "fg": self.get_color("destructive.foreground"),
                "border": self.get_color("destructive.DEFAULT"),
                "hover_bg": self.tokens["destructive"]["DEFAULT"],
            },
        }
        return variants.get(self.variant, variants["default"])
    
    def render(self, slide, left: float, top: float, 
               width: Optional[float] = None, 
               height: Optional[float] = None) -> Any:
        """
        Render button to slide.
        
        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Optional width override
            height: Optional height override
        
        Returns:
            Shape object representing the button
        """
        size_config = self.get_size_config()
        colors = self.get_variant_colors()
        
        # Use provided dimensions or defaults
        btn_width = width or size_config["width"]
        btn_height = height or size_config["height"]
        
        # Create button shape
        button = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(btn_width),
            Inches(btn_height)
        )
        
        # Apply background
        if colors["bg"]:
            button.fill.solid()
            button.fill.fore_color.rgb = colors["bg"]
        else:
            button.fill.background()  # Transparent
        
        # Apply border
        if colors["border"]:
            button.line.color.rgb = colors["border"]
            button.line.width = Pt(1)
        else:
            button.line.fill.background()  # No border
        
        # Add text
        text_frame = button.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(size_config["padding"])
        text_frame.margin_right = Inches(size_config["padding"])
        text_frame.margin_top = Inches(size_config["padding"] / 2)
        text_frame.margin_bottom = Inches(size_config["padding"] / 2)
        
        paragraph = text_frame.paragraphs[0]
        paragraph.text = self.text
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(size_config["font_size"])
        paragraph.font.color.rgb = colors["fg"]
        paragraph.font.name = self.theme.get("font_family", "Inter")
        
        # Make text bold for primary and destructive variants
        if self.variant in ["primary", "destructive"]:
            paragraph.font.bold = True
        
        return button


class IconButton(Button):
    """
    Icon button component - a button with an icon instead of text.
    Note: PowerPoint doesn't have native icon support, so we use Unicode symbols.
    """
    
    # Common icon mappings using Unicode
    ICONS = {
        "play": "▶",
        "pause": "⏸",
        "stop": "⏹",
        "next": "⏭",
        "previous": "⏮",
        "plus": "+",
        "minus": "-",
        "close": "✕",
        "check": "✓",
        "star": "★",
        "heart": "♥",
        "settings": "⚙",
        "menu": "☰",
        "search": "🔍",
        "download": "⬇",
        "upload": "⬆",
        "refresh": "↻",
        "share": "⤴",
        "edit": "✎",
        "delete": "🗑",
    }
    
    def __init__(self,
                 icon: str,
                 variant: str = "ghost",
                 size: str = "md",
                 theme: Optional[Dict[str, Any]] = None):
        """
        Initialize icon button.
        
        Args:
            icon: Icon name or Unicode character
            variant: Button variant
            size: Button size
            theme: Theme configuration
        """
        # Get icon character or use provided Unicode
        icon_char = self.ICONS.get(icon, icon)
        super().__init__(icon_char, variant, size, theme)
    
    def get_size_config(self) -> Dict[str, Any]:
        """Get size configuration for icon button (square)."""
        sizes = {
            "sm": {
                "width": 0.4,
                "height": 0.4,
                "font_size": 14,
                "padding": self.get_padding("xs"),
            },
            "md": {
                "width": 0.5,
                "height": 0.5,
                "font_size": 16,
                "padding": self.get_padding("sm"),
            },
            "lg": {
                "width": 0.6,
                "height": 0.6,
                "font_size": 20,
                "padding": self.get_padding("sm"),
            },
        }
        return sizes.get(self.size, sizes["md"])