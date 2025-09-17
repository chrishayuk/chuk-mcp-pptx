"""
Chart components for PowerPoint presentations.
Theme-aware data visualization components.
"""

from typing import Dict, Any, List, Optional, Tuple
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches
import asyncio

from .base import AsyncComponent


class Chart(AsyncComponent):
    """
    Base chart component for data visualization.
    """
    
    def __init__(self, 
                 title: Optional[str] = None,
                 categories: List[str] = None,
                 series: Dict[str, List[float]] = None,
                 theme: Optional[Dict[str, Any]] = None):
        """
        Initialize chart component.
        
        Args:
            title: Chart title
            categories: Category labels
            series: Data series dictionary
            theme: Theme configuration
        """
        super().__init__(theme)
        self.title = title
        self.categories = categories or []
        self.series = series or {}
        self.chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    
    def _prepare_chart_data(self) -> CategoryChartData:
        """Prepare chart data for PowerPoint."""
        chart_data = CategoryChartData()
        chart_data.categories = self.categories
        
        for series_name, values in self.series.items():
            chart_data.add_series(series_name, values)
        
        return chart_data
    
    def _apply_theme_to_chart(self, chart):
        """Apply theme colors to chart."""
        # Get theme colors
        chart_colors = self.tokens.get("chart", [])
        
        # Apply colors to series
        for idx, series in enumerate(chart.series):
            if idx < len(chart_colors):
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = self.get_color(f"chart.{idx}")
    
    def _render_sync(self, slide, left: float = 1.0, top: float = 2.0,
                     width: float = 8.0, height: float = 4.5) -> Any:
        """
        Synchronous render implementation.
        
        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Chart width in inches
            height: Chart height in inches
        
        Returns:
            Chart object
        """
        chart_data = self._prepare_chart_data()
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            self.chart_type,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart
        
        # Set title
        if self.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = self.title
        
        # Configure legend
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        # Apply theme
        self._apply_theme_to_chart(chart)
        
        return chart


class BarChart(Chart):
    """Bar chart component."""
    
    def __init__(self, horizontal: bool = False, **kwargs):
        """
        Initialize bar chart.
        
        Args:
            horizontal: Whether bars should be horizontal
            **kwargs: Chart parameters
        """
        super().__init__(**kwargs)
        self.chart_type = (
            XL_CHART_TYPE.BAR_CLUSTERED if horizontal 
            else XL_CHART_TYPE.COLUMN_CLUSTERED
        )


class LineChart(Chart):
    """Line chart component."""
    
    def __init__(self, smooth: bool = True, markers: bool = True, **kwargs):
        """
        Initialize line chart.
        
        Args:
            smooth: Whether lines should be smoothed
            markers: Whether to show data point markers
            **kwargs: Chart parameters
        """
        super().__init__(**kwargs)
        self.smooth = smooth
        self.markers = markers
        self.chart_type = XL_CHART_TYPE.LINE
    
    def _render_sync(self, slide, **kwargs):
        """Render line chart with additional configuration."""
        chart = super()._render_sync(slide, **kwargs)
        
        # Configure line styles
        for series in chart.series:
            series.smooth = self.smooth
            if self.markers:
                series.marker.style = 1  # Circle markers
                series.marker.size = 7
        
        return chart


class PieChart(Chart):
    """Pie chart component."""
    
    def __init__(self, 
                 categories: List[str] = None,
                 values: List[float] = None,
                 explode: Optional[int] = None,
                 **kwargs):
        """
        Initialize pie chart.
        
        Args:
            categories: Category labels
            values: Data values
            explode: Index of slice to explode
            **kwargs: Chart parameters
        """
        # Convert to series format for base class
        series = {"Values": values} if values else {}
        super().__init__(categories=categories, series=series, **kwargs)
        self.chart_type = XL_CHART_TYPE.PIE
        self.explode = explode
    
    def _render_sync(self, slide, **kwargs):
        """Render pie chart with additional configuration."""
        chart = super()._render_sync(slide, **kwargs)
        
        # Explode slice if specified
        if self.explode is not None and 0 <= self.explode < len(self.categories):
            chart.plots[0].series[0].points[self.explode].explosion = 20
        
        # Show data labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_category_name = True
        
        return chart


class AreaChart(Chart):
    """Area chart component."""
    
    def __init__(self, stacked: bool = False, **kwargs):
        """
        Initialize area chart.
        
        Args:
            stacked: Whether areas should be stacked
            **kwargs: Chart parameters
        """
        super().__init__(**kwargs)
        self.chart_type = (
            XL_CHART_TYPE.AREA_STACKED if stacked
            else XL_CHART_TYPE.AREA
        )


class ScatterChart(AsyncComponent):
    """Scatter plot component."""
    
    def __init__(self,
                 title: Optional[str] = None,
                 data_points: List[Tuple[float, float]] = None,
                 series_name: str = "Series 1",
                 theme: Optional[Dict[str, Any]] = None):
        """
        Initialize scatter chart.
        
        Args:
            title: Chart title
            data_points: List of (x, y) tuples
            series_name: Name of data series
            theme: Theme configuration
        """
        super().__init__(theme)
        self.title = title
        self.data_points = data_points or []
        self.series_name = series_name
    
    def _prepare_chart_data(self) -> XyChartData:
        """Prepare XY chart data."""
        chart_data = XyChartData()
        
        series = chart_data.add_series(self.series_name)
        for x, y in self.data_points:
            series.add_data_point(x, y)
        
        return chart_data
    
    def _render_sync(self, slide, left: float = 1.0, top: float = 2.0,
                     width: float = 8.0, height: float = 4.5) -> Any:
        """Render scatter chart."""
        chart_data = self._prepare_chart_data()
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart
        
        # Set title
        if self.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = self.title
        
        # Apply theme colors
        fill = chart.series[0].format.fill
        fill.solid()
        fill.fore_color.rgb = self.get_color("primary.DEFAULT")
        
        return chart


class DoughnutChart(PieChart):
    """Doughnut chart component."""
    
    def __init__(self, hole_size: int = 50, **kwargs):
        """
        Initialize doughnut chart.
        
        Args:
            hole_size: Size of center hole (0-90)
            **kwargs: Chart parameters
        """
        super().__init__(**kwargs)
        self.chart_type = XL_CHART_TYPE.DOUGHNUT
        self.hole_size = max(0, min(90, hole_size))
    
    def _render_sync(self, slide, **kwargs):
        """Render doughnut chart."""
        chart = super()._render_sync(slide, **kwargs)
        
        # Note: python-pptx doesn't directly support hole size adjustment
        # This would require XML manipulation
        
        return chart