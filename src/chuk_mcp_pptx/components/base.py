"""
Base component class for all PowerPoint components.
Provides common functionality and theme integration.
"""

from typing import Dict, Any, Optional, Tuple
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import asyncio

from ..tokens.colors import get_semantic_tokens, PALETTE
from ..tokens.typography import get_text_style, FONT_FAMILIES
from ..tokens.spacing import SPACING, PADDING, MARGINS, RADIUS


class Component:
    """
    Base class for all PowerPoint components.
    Handles theme integration and common properties.
    """
    
    def __init__(self, theme: Optional[Dict[str, Any]] = None):
        """
        Initialize component with theme.
        
        Args:
            theme: Theme configuration or None for default
        """
        self.theme = theme or self.get_default_theme()
        self.tokens = get_semantic_tokens(
            self.theme.get("primary_hue", "blue"),
            self.theme.get("mode", "dark")
        )
    
    @staticmethod
    def get_default_theme() -> Dict[str, Any]:
        """Get default theme configuration."""
        return {
            "name": "default",
            "mode": "dark",
            "primary_hue": "blue",
            "font_family": "Inter",
            "radius": "md",
            "spacing": "default",
        }
    
    def hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def get_color(self, color_path: str) -> RGBColor:
        """
        Get color from theme tokens.
        
        Args:
            color_path: Dot-separated path to color (e.g., "primary.DEFAULT")
        
        Returns:
            RGBColor object
        """
        parts = color_path.split('.')
        value = self.tokens
        
        for part in parts:
            if isinstance(value, dict):
                value = value.get(part, "#000000")
            else:
                break
        
        if isinstance(value, str):
            return RGBColor(*self.hex_to_rgb(value))
        return RGBColor(0, 0, 0)
    
    def get_spacing(self, size: str) -> float:
        """
        Get spacing value in inches.
        
        Args:
            size: Spacing size key (e.g., "4", "md")
        
        Returns:
            Spacing value in inches
        """
        # Check if it's a preset like "md"
        if size in MARGINS:
            return MARGINS[size]
        # Otherwise get from spacing scale
        return SPACING.get(size, SPACING["4"])
    
    def get_padding(self, size: str) -> float:
        """Get padding value in inches."""
        return PADDING.get(size, PADDING["md"])
    
    def get_text_style(self, variant: str) -> Dict[str, Any]:
        """Get text style configuration."""
        return get_text_style(variant)
    
    def apply_text_style(self, text_frame, variant: str = "body"):
        """
        Apply text style to a text frame.
        
        Args:
            text_frame: PowerPoint text frame object
            variant: Text style variant
        """
        style = self.get_text_style(variant)
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = style.get("font_family", "Inter")
            paragraph.font.size = Pt(style.get("font_size", 14))
            
            # Apply font weight (PowerPoint has limited support)
            weight = style.get("font_weight", 400)
            paragraph.font.bold = weight >= 600
            
            # Apply color
            paragraph.font.color.rgb = self.get_color("foreground.DEFAULT")
    
    def apply_shape_style(self, shape, style_type: str = "card"):
        """
        Apply theme styling to a shape.
        
        Args:
            shape: PowerPoint shape object
            style_type: Style type (card, primary, secondary, accent)
        """
        # Determine colors based on style type
        if style_type == "primary":
            bg_path = "primary.DEFAULT"
            fg_path = "primary.foreground"
            border_path = "primary.DEFAULT"
        elif style_type == "secondary":
            bg_path = "secondary.DEFAULT"
            fg_path = "secondary.foreground"
            border_path = "border.DEFAULT"
        elif style_type == "accent":
            bg_path = "accent.DEFAULT"
            fg_path = "accent.foreground"
            border_path = "accent.DEFAULT"
        elif style_type == "muted":
            bg_path = "muted.DEFAULT"
            fg_path = "muted.foreground"
            border_path = "border.secondary"
        else:  # card
            bg_path = "card.DEFAULT"
            fg_path = "card.foreground"
            border_path = "border.DEFAULT"
        
        # Apply fill
        if hasattr(shape, 'fill'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.get_color(bg_path)
        
        # Apply border
        if hasattr(shape, 'line'):
            line = shape.line
            line.color.rgb = self.get_color(border_path)
            line.width = Pt(1)
        
        # Apply text color
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self.get_color(fg_path)
    
    async def render(self, slide, **kwargs):
        """
        Render component to slide (to be implemented by subclasses).
        
        Args:
            slide: PowerPoint slide object
            **kwargs: Component-specific parameters
        """
        raise NotImplementedError("Subclasses must implement render method")


class AsyncComponent(Component):
    """Async version of Component base class."""
    
    async def render(self, slide, **kwargs):
        """Async render method."""
        # Run synchronous operations in thread pool
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(None, lambda: self._render_sync(slide, **kwargs))
    
    def _render_sync(self, slide, **kwargs):
        """Synchronous render implementation."""
        raise NotImplementedError("Subclasses must implement _render_sync method")