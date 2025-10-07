"""
WhatsApp chat components for PowerPoint presentations.

Provides realistic WhatsApp-style chat interfaces.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component


class WhatsAppBubble(Component):
    """
    WhatsApp-style chat bubble component.

    Replicates WhatsApp's distinctive look with:
    - Green bubbles for sent messages (right-aligned)
    - White bubbles for received messages (left-aligned)
    - Tail on bubbles
    - Timestamp and check marks inside bubbles
    - WhatsApp green (#25D366)

    Examples:
        # Sent message
        msg = WhatsAppBubble(
            text="Hey! How are you?",
            variant="sent",
            timestamp="10:30",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=7)

        # Received message with sender (group chat)
        msg = WhatsAppBubble(
            text="I'm good, thanks!",
            sender="John",
            variant="received",
            timestamp="10:31",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=7)
    """

    def __init__(self,
                 text: str,
                 sender: Optional[str] = None,
                 variant: str = "received",
                 timestamp: Optional[str] = None,
                 show_checkmarks: bool = True,
                 theme: Optional[Dict[str, Any]] = None):
        """
        Initialize WhatsApp bubble.

        Args:
            text: Message text
            sender: Sender name (for group chats)
            variant: Message variant (sent, received)
            timestamp: Optional timestamp (e.g., "10:30")
            show_checkmarks: Show double checkmarks for sent messages
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.sender = sender
        self.variant = variant
        self.timestamp = timestamp
        self.show_checkmarks = show_checkmarks

    def _get_bubble_color(self) -> RGBColor:
        """Get WhatsApp bubble color."""
        if self.variant == "sent":
            # WhatsApp green (#DCF8C6) - light green for sent
            return RGBColor(220, 248, 198)
        else:
            # White for received
            return RGBColor(255, 255, 255)

    def _get_text_color(self) -> RGBColor:
        """Get text color."""
        return RGBColor(0, 0, 0)  # Black text for both

    def _calculate_bubble_height(self, width: float) -> float:
        """Estimate bubble height."""
        chars_per_line = int(width * 10)  # More generous
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        line_height = 0.2  # Increased from 0.18
        padding = 0.35  # More padding
        sender_height = 0.22 if self.sender else 0  # More space for sender
        timestamp_height = 0.18  # More space for inline timestamp
        return (lines * line_height) + padding + sender_height + timestamp_height

    def render(self, slide, left: float, top: float, width: float = 6.0) -> list:
        """Render WhatsApp bubble."""
        shapes = []

        # Calculate dimensions
        bubble_width = min(width * 0.7, width - 1.0)
        bubble_height = self._calculate_bubble_height(bubble_width)

        # Position based on variant
        if self.variant == "sent":
            bubble_left = left + width - bubble_width
        else:
            bubble_left = left

        # Create bubble
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bubble_left),
            Inches(top),
            Inches(bubble_width),
            Inches(bubble_height)
        )

        # Style bubble - WhatsApp style
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = self._get_bubble_color()

        # Light border for white bubbles
        if self.variant == "received":
            bubble.line.color.rgb = RGBColor(230, 230, 230)
            bubble.line.width = Pt(0.5)
        else:
            bubble.line.fill.background()

        # Subtle shadow
        bubble.shadow.visible = True
        bubble.shadow.blur_radius = Pt(2)
        bubble.shadow.distance = Pt(1)
        bubble.shadow.angle = 90
        bubble.shadow.transparency = 0.5

        # Text frame
        text_frame = bubble.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        padding = 0.1
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        current_p = text_frame.paragraphs[0]

        # Sender name (for group chats, received messages)
        if self.sender and self.variant == "received":
            current_p.text = self.sender
            current_p.alignment = PP_ALIGN.LEFT
            current_p.font.size = Pt(13)  # Larger sender name
            current_p.font.bold = True
            current_p.font.color.rgb = RGBColor(6, 124, 98)  # Stronger WhatsApp teal
            current_p.space_after = Pt(4)  # More space after sender
            current_p = text_frame.add_paragraph()

        # Message text
        current_p.text = self.text
        current_p.alignment = PP_ALIGN.LEFT
        current_p.font.size = Pt(15)  # Larger font
        current_p.font.name = "Helvetica Neue"  # WhatsApp font
        current_p.font.color.rgb = self._get_text_color()
        current_p.line_spacing = 1.3  # Better line spacing

        shapes.append(bubble)

        # Timestamp and checkmarks (inline, bottom right of bubble)
        if self.timestamp or (self.variant == "sent" and self.show_checkmarks):
            ts_text = ""
            if self.timestamp:
                ts_text = self.timestamp
            if self.variant == "sent" and self.show_checkmarks:
                ts_text += " ✓✓"  # Double checkmark

            if ts_text:
                ts_box = slide.shapes.add_textbox(
                    Inches(bubble_left + bubble_width - 0.7),
                    Inches(top + bubble_height - 0.22),
                    Inches(0.6),
                    Inches(0.18)
                )
                ts_frame = ts_box.text_frame
                ts_frame.text = ts_text
                ts_frame.margin_top = 0
                ts_frame.margin_right = Inches(0.05)
                ts_p = ts_frame.paragraphs[0]
                ts_p.alignment = PP_ALIGN.RIGHT
                ts_p.font.size = Pt(9)
                ts_p.font.color.rgb = RGBColor(150, 150, 150)  # Gray
                shapes.append(ts_box)

        return shapes


class WhatsAppConversation(Component):
    """
    WhatsApp-style conversation component.

    Displays a complete WhatsApp conversation with automatic spacing.

    Examples:
        messages = [
            {"text": "Hey!", "sender": "Alice", "variant": "received", "timestamp": "10:30"},
            {"text": "Hi there!", "variant": "sent", "timestamp": "10:31"}
        ]
        conversation = WhatsAppConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=8)
    """

    def __init__(self,
                 messages: List[Dict[str, Any]],
                 spacing: float = 0.12,
                 theme: Optional[Dict[str, Any]] = None):
        """
        Initialize WhatsApp conversation.

        Args:
            messages: List of message dicts with WhatsAppBubble params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render WhatsApp conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = WhatsAppBubble(
                text=msg_data.get("text", ""),
                sender=msg_data.get("sender"),
                variant=msg_data.get("variant", "received"),
                timestamp=msg_data.get("timestamp"),
                show_checkmarks=msg_data.get("show_checkmarks", True),
                theme=self.theme
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            bubble_width = min(width * 0.7, width - 1.0)
            msg_height = message._calculate_bubble_height(bubble_width)
            current_top += msg_height + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# WhatsAppBubble - Variants: sent, received
# WhatsAppConversation - Full conversation flow
