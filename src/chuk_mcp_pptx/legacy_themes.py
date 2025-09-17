"""
Beautiful themes for PowerPoint presentations.
Inspired by modern design systems like shadcn/ui.
"""

from typing import Dict, Tuple, Any
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Modern theme definitions inspired by shadcn/ui and contemporary design
THEMES = {
    "dark_modern": {
        "name": "Dark Modern",
        "background": "#0a0a0a",
        "foreground": "#fafafa",
        "card": "#1a1a1a",
        "card_foreground": "#fafafa",
        "primary": "#f97316",  # Orange accent
        "primary_foreground": "#fafafa",
        "secondary": "#27272a",
        "secondary_foreground": "#fafafa",
        "muted": "#27272a",
        "muted_foreground": "#a1a1aa",
        "accent": "#f97316",
        "accent_foreground": "#0a0a0a",
        "border": "#27272a",
        "chart_colors": ["#f97316", "#06b6d4", "#8b5cf6", "#10b981", "#f43f5e", "#6366f1"],
        "font_heading": "Segoe UI Semibold",
        "font_body": "Segoe UI",
        "font_code": "Cascadia Code"
    },
    
    "dark_blue": {
        "name": "Dark Blue",
        "background": "#0f172a",
        "foreground": "#f1f5f9",
        "card": "#1e293b",
        "card_foreground": "#f1f5f9",
        "primary": "#3b82f6",  # Blue accent
        "primary_foreground": "#ffffff",
        "secondary": "#334155",
        "secondary_foreground": "#f1f5f9",
        "muted": "#334155",
        "muted_foreground": "#94a3b8",
        "accent": "#3b82f6",
        "accent_foreground": "#ffffff",
        "border": "#334155",
        "chart_colors": ["#3b82f6", "#10b981", "#f59e0b", "#ef4444", "#8b5cf6", "#ec4899"],
        "font_heading": "Inter Bold",
        "font_body": "Inter",
        "font_code": "JetBrains Mono"
    },
    
    "dark_purple": {
        "name": "Dark Purple",
        "background": "#0c0a1d",
        "foreground": "#e9e3ff",
        "card": "#1a1630",
        "card_foreground": "#e9e3ff",
        "primary": "#8b5cf6",  # Purple accent
        "primary_foreground": "#ffffff",
        "secondary": "#2d2650",
        "secondary_foreground": "#e9e3ff",
        "muted": "#2d2650",
        "muted_foreground": "#a394f5",
        "accent": "#8b5cf6",
        "accent_foreground": "#ffffff",
        "border": "#2d2650",
        "chart_colors": ["#8b5cf6", "#ec4899", "#3b82f6", "#10b981", "#f59e0b", "#ef4444"],
        "font_heading": "Poppins SemiBold",
        "font_body": "Poppins",
        "font_code": "Fira Code"
    },
    
    "dark_green": {
        "name": "Dark Green",
        "background": "#0a0f0a",
        "foreground": "#e6f4e6",
        "card": "#162316",
        "card_foreground": "#e6f4e6",
        "primary": "#10b981",  # Green accent
        "primary_foreground": "#ffffff",
        "secondary": "#1f3a1f",
        "secondary_foreground": "#e6f4e6",
        "muted": "#1f3a1f",
        "muted_foreground": "#86efac",
        "accent": "#10b981",
        "accent_foreground": "#ffffff",
        "border": "#1f3a1f",
        "chart_colors": ["#10b981", "#06b6d4", "#3b82f6", "#8b5cf6", "#f59e0b", "#ef4444"],
        "font_heading": "Montserrat SemiBold",
        "font_body": "Montserrat",
        "font_code": "Source Code Pro"
    },
    
    "light_minimal": {
        "name": "Light Minimal",
        "background": "#ffffff",
        "foreground": "#0a0a0a",
        "card": "#f4f4f5",
        "card_foreground": "#0a0a0a",
        "primary": "#18181b",
        "primary_foreground": "#fafafa",
        "secondary": "#f4f4f5",
        "secondary_foreground": "#18181b",
        "muted": "#f4f4f5",
        "muted_foreground": "#71717a",
        "accent": "#18181b",
        "accent_foreground": "#fafafa",
        "border": "#e4e4e7",
        "chart_colors": ["#18181b", "#3b82f6", "#10b981", "#f59e0b", "#ef4444", "#8b5cf6"],
        "font_heading": "Helvetica Neue Medium",
        "font_body": "Helvetica Neue",
        "font_code": "Monaco"
    },
    
    "light_warm": {
        "name": "Light Warm",
        "background": "#fffbf7",
        "foreground": "#1c1917",
        "card": "#fef3e7",
        "card_foreground": "#1c1917",
        "primary": "#ea580c",
        "primary_foreground": "#ffffff",
        "secondary": "#fed7aa",
        "secondary_foreground": "#1c1917",
        "muted": "#fed7aa",
        "muted_foreground": "#78716c",
        "accent": "#ea580c",
        "accent_foreground": "#ffffff",
        "border": "#fed7aa",
        "chart_colors": ["#ea580c", "#dc2626", "#16a34a", "#2563eb", "#7c3aed", "#db2777"],
        "font_heading": "Playfair Display",
        "font_body": "Lato",
        "font_code": "Consolas"
    },
    
    "cyberpunk": {
        "name": "Cyberpunk",
        "background": "#0a0014",
        "foreground": "#00ffff",
        "card": "#1a0028",
        "card_foreground": "#00ffff",
        "primary": "#ff00ff",
        "primary_foreground": "#0a0014",
        "secondary": "#2a0042",
        "secondary_foreground": "#00ffff",
        "muted": "#2a0042",
        "muted_foreground": "#00cccc",
        "accent": "#ffff00",
        "accent_foreground": "#0a0014",
        "border": "#ff00ff",
        "chart_colors": ["#ff00ff", "#00ffff", "#ffff00", "#ff1493", "#00ff00", "#ff4500"],
        "font_heading": "Orbitron Bold",
        "font_body": "Exo 2",
        "font_code": "Hack"
    },
    
    "gradient_sunset": {
        "name": "Gradient Sunset",
        "background": "#1a1625",
        "foreground": "#fff5f5",
        "card": "#2d2438",
        "card_foreground": "#fff5f5",
        "primary": "#ff6b6b",
        "primary_foreground": "#ffffff",
        "secondary": "#4ecdc4",
        "secondary_foreground": "#1a1625",
        "muted": "#3d3456",
        "muted_foreground": "#c9b8db",
        "accent": "#f7b731",
        "accent_foreground": "#1a1625",
        "border": "#4a3d5c",
        "chart_colors": ["#ff6b6b", "#4ecdc4", "#f7b731", "#a8e6cf", "#ff8b94", "#c7ceea"],
        "font_heading": "Raleway Bold",
        "font_body": "Open Sans",
        "font_code": "IBM Plex Mono"
    }
}

def apply_theme(slide, theme_name: str = "dark_modern"):
    """
    Apply a theme to a slide.
    
    Args:
        slide: The slide to apply theme to
        theme_name: Name of the theme to apply
    """
    theme = THEMES.get(theme_name, THEMES["dark_modern"])
    
    # Set slide background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["background"]))
    
    return theme

def apply_theme_to_shape(shape, theme: Dict[str, Any], shape_type: str = "card"):
    """
    Apply theme colors to a shape.
    
    Args:
        shape: The shape to style
        theme: Theme dictionary
        shape_type: Type of shape (card, primary, secondary, accent)
    """
    if shape_type == "card":
        bg_color = theme["card"]
        fg_color = theme["card_foreground"]
    elif shape_type == "primary":
        bg_color = theme["primary"]
        fg_color = theme["primary_foreground"]
    elif shape_type == "secondary":
        bg_color = theme["secondary"]
        fg_color = theme["secondary_foreground"]
    elif shape_type == "accent":
        bg_color = theme["accent"]
        fg_color = theme["accent_foreground"]
    else:
        bg_color = theme["muted"]
        fg_color = theme["muted_foreground"]
    
    # Apply fill
    if hasattr(shape, 'fill'):
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_color))
    
    # Apply border
    if hasattr(shape, 'line'):
        line = shape.line
        line.color.rgb = RGBColor(*hex_to_rgb(theme["border"]))
        line.width = Pt(1)
    
    # Apply text color
    if hasattr(shape, 'text_frame') and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*hex_to_rgb(fg_color))

def create_code_block(slide, code: str, language: str, 
                      left: float, top: float, width: float, height: float,
                      theme_name: str = "dark_modern"):
    """
    Create a code block with syntax highlighting appearance.
    
    Args:
        slide: Slide to add code block to
        code: Code content
        language: Programming language
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        theme_name: Theme to use
    """
    from pptx.enum.shapes import MSO_SHAPE
    
    theme = THEMES.get(theme_name, THEMES["dark_modern"])
    
    # Create container shape
    code_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    
    # Style as code block
    fill = code_shape.fill
    fill.solid()
    # Use card color for code background
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["card"]))
    
    # Add border
    line = code_shape.line
    line.color.rgb = RGBColor(*hex_to_rgb(theme["border"]))
    line.width = Pt(1)
    
    # Add code text
    text_frame = code_shape.text_frame
    text_frame.clear()
    text_frame.margin_left = Inches(0.2)
    text_frame.margin_right = Inches(0.2)
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_bottom = Inches(0.15)
    
    # Add language label
    p = text_frame.paragraphs[0]
    p.text = f"// {language}"
    p.font.name = theme.get("font_code", "Cascadia Code")
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(*hex_to_rgb(theme["muted_foreground"]))
    
    # Add code content
    for line in code.split('\n'):
        p = text_frame.add_paragraph()
        p.text = line
        p.font.name = theme.get("font_code", "Cascadia Code")
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(*hex_to_rgb(theme["foreground"]))
        p.level = 0
    
    return code_shape

def create_gradient_background(slide, colors: list, direction: str = "diagonal"):
    """
    Create a gradient background for a slide.
    
    Args:
        slide: Slide to apply gradient to
        colors: List of hex colors for gradient
        direction: Gradient direction (horizontal, vertical, diagonal)
    """
    # Note: python-pptx doesn't directly support gradients,
    # but we can simulate with shapes
    from pptx.enum.shapes import MSO_SHAPE
    
    # For now, apply solid color
    # Full gradient implementation would require XML manipulation
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(colors[0]))
    
    return slide

def get_chart_colors(theme_name: str = "dark_modern") -> list:
    """
    Get chart colors for a theme.
    
    Args:
        theme_name: Name of the theme
        
    Returns:
        List of RGB tuples for chart colors
    """
    theme = THEMES.get(theme_name, THEMES["dark_modern"])
    return [hex_to_rgb(color) for color in theme["chart_colors"]]

def get_slide_theme(slide) -> Dict[str, Any]:
    """
    Detect which theme is applied to a slide based on background color.
    
    Args:
        slide: The slide to check
        
    Returns:
        Theme dictionary or default theme
    """
    # Try to detect theme from background color
    try:
        bg_fill = slide.background.fill
        if hasattr(bg_fill, 'fore_color') and bg_fill.fore_color:
            if hasattr(bg_fill.fore_color, 'rgb'):
                rgb = bg_fill.fore_color.rgb
                bg_hex = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                
                # Find matching theme
                for theme_name, theme_data in THEMES.items():
                    if theme_data["background"].lower() == bg_hex.lower():
                        return theme_data
    except:
        pass
    
    # Return default if no match
    return THEMES["dark_modern"]

def get_theme_colors_for_shape(slide, shape_style: str = "primary") -> Tuple[str, str]:
    """
    Get appropriate colors for a shape based on slide theme.
    
    Args:
        slide: The slide containing the shape
        shape_style: Style of shape (primary, secondary, accent, muted)
        
    Returns:
        Tuple of (fill_color_hex, text_color_hex)
    """
    theme = get_slide_theme(slide)
    
    if shape_style == "primary":
        return theme["primary"], theme["primary_foreground"]
    elif shape_style == "secondary":
        return theme["secondary"], theme["secondary_foreground"]
    elif shape_style == "accent":
        return theme["accent"], theme["accent_foreground"]
    elif shape_style == "muted":
        return theme["muted"], theme["muted_foreground"]
    else:
        # Default to card colors
        return theme["card"], theme["card_foreground"]