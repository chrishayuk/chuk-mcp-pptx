"""
Presentation Manager for PowerPoint MCP Server

Manages PowerPoint presentations with support for chuk-artifacts integration.
Each presentation is stored as a BLOB namespace for persistence and multi-server access.

Uses chuk-mcp-server's built-in artifact store context for storage.
Uses Pydantic models throughout for type safety and validation.
"""

from __future__ import annotations

import asyncio
import base64
import io
import logging
from datetime import datetime
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from ..models import (
    PresentationMetadata,
    SlideMetadata,
    PresentationInfo,
    ListPresentationsResponse,
)

logger = logging.getLogger(__name__)


class PresentationManager:
    """
    Manages PowerPoint presentations with chuk-artifacts integration.

    Uses chuk-mcp-server's built-in artifact store context for flexible storage
    (memory, filesystem, sqlite, s3). Each presentation is stored as a BLOB
    namespace with automatic session management.
    Presentations and metadata are Pydantic models for type safety.
    """

    # MIME type for PowerPoint presentations
    PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def __init__(self, base_path: str = "presentations") -> None:
        """
        Initialize the presentation manager.

        Args:
            base_path: Base path prefix for presentation names
        """
        self.base_path = base_path
        self._presentations: dict[str, PresentationType] = {}
        self._metadata: dict[str, PresentationMetadata] = {}
        self._namespace_ids: dict[str, str] = {}  # name -> namespace_id mapping
        self._current_presentation: str | None = None
        logger.info(f"PresentationManager initialized, base path: {base_path}")

    def _get_store(self):
        """Get the artifact store from context."""
        from chuk_mcp_server import get_artifact_store, has_artifact_store

        if has_artifact_store():
            return get_artifact_store()
        return None

    def _sanitize_name(self, name: str) -> str:
        """Sanitize presentation name to prevent directory traversal."""
        safe_name = "".join(c for c in name if c.isalnum() or c in ("-", "_"))
        if not safe_name:
            safe_name = "presentation"
        return safe_name

    def get_namespace_id(self, name: str) -> str | None:
        """Get the namespace ID for a presentation by name."""
        return self._namespace_ids.get(name)

    def get_artifact_uri(self, name: str) -> str | None:
        """
        Get the artifact URI for a presentation.

        Args:
            name: Presentation name

        Returns:
            Artifact URI string or None if not found
        """
        namespace_id = self._namespace_ids.get(name)
        if namespace_id:
            return f"artifact://chuk-mcp-pptx/{self.base_path}/{name}"
        return None

    async def _save_to_store(self, name: str, prs: PresentationType) -> bool:
        """
        Save presentation to artifact store.

        Args:
            name: Presentation name
            prs: Presentation object

        Returns:
            True if successful, False otherwise
        """
        store = self._get_store()
        if not store:
            logger.debug("No artifact store available, skipping persistence")
            return False

        from chuk_mcp_server import NamespaceType, StorageScope

        try:
            # Convert presentation to bytes (wrap blocking I/O)
            buffer = io.BytesIO()
            await asyncio.to_thread(prs.save, buffer)
            buffer.seek(0)
            data = buffer.read()

            # Check if namespace already exists
            namespace_id = self._namespace_ids.get(name)

            if namespace_id:
                # Update existing namespace
                await store.write_namespace(namespace_id, data=data)
                logger.info(f"Updated presentation in artifact store: {name} ({namespace_id})")
            else:
                # Create new BLOB namespace
                safe_name = self._sanitize_name(name)
                namespace_info = await store.create_namespace(
                    type=NamespaceType.BLOB,
                    scope=StorageScope.SESSION,
                    name=f"{self.base_path}/{safe_name}",
                    metadata={
                        "mime_type": self.PPTX_MIME_TYPE,
                        "presentation_name": name,
                        "file_extension": ".pptx",
                    },
                )
                self._namespace_ids[name] = namespace_info.namespace_id

                # Write the presentation data
                await store.write_namespace(namespace_info.namespace_id, data=data)
                logger.info(
                    f"Saved presentation to artifact store: {name} ({namespace_info.namespace_id})"
                )

            return True
        except Exception as e:
            logger.error(f"Failed to save to artifact store: {e}")
            return False

    async def _load_from_store(self, name: str) -> PresentationType | None:
        """
        Load presentation from artifact store.

        Args:
            name: Presentation name

        Returns:
            Presentation object or None if not found
        """
        store = self._get_store()
        if not store:
            logger.debug("No artifact store available")
            return None

        try:
            namespace_id = self._namespace_ids.get(name)
            if not namespace_id:
                logger.debug(f"Presentation not found in namespace mapping: {name}")
                return None

            # Read from artifact store
            data = await store.read_namespace(namespace_id)
            if data is None:
                logger.debug(f"No data found for namespace: {namespace_id}")
                return None

            buffer = io.BytesIO(data)
            prs = Presentation(buffer)
            logger.info(f"Loaded presentation from artifact store: {name} ({namespace_id})")
            return prs
        except Exception as e:
            logger.error(f"Failed to load from artifact store: {e}")
            return None

    async def _delete_from_store(self, name: str) -> bool:
        """
        Delete presentation from artifact store.

        Args:
            name: Presentation name

        Returns:
            True if successful, False otherwise
        """
        store = self._get_store()
        if not store:
            logger.debug("No artifact store available")
            return False

        try:
            namespace_id = self._namespace_ids.get(name)
            if not namespace_id:
                logger.warning(f"Presentation not found in namespace mapping for deletion: {name}")
                return False

            await store.destroy_namespace(namespace_id)
            del self._namespace_ids[name]
            logger.info(f"Deleted presentation from artifact store: {name} ({namespace_id})")
            return True
        except Exception as e:
            logger.error(f"Failed to delete from artifact store: {e}")  # nosec B608
            return False

    async def _load_template_from_store(self, template_name: str) -> bytes | None:
        """
        Load a template from artifact store by name.

        Args:
            template_name: Name of the template presentation in artifact store

        Returns:
            Template data as bytes or None if not found
        """
        store = self._get_store()
        if not store:
            logger.error("No artifact store available for template loading")
            return None

        try:
            # Check if template exists in our namespace mapping
            namespace_id = self._namespace_ids.get(template_name)
            if namespace_id:
                data = await store.read_namespace(namespace_id)
                if data:
                    logger.info(f"Loaded template from artifact store: {template_name}")
                    return data

            # Try to find template by searching for it
            logger.warning(f"Template not found in namespace mapping: {template_name}")
            return None
        except Exception as e:
            logger.error(f"Failed to load template from artifact store: {e}")
            return None

    async def create(
        self,
        name: str,
        theme: str | None = None,
        template_name: str | None = None,
    ) -> PresentationMetadata:
        """
        Create a new presentation, optionally from a template.

        Args:
            name: Presentation name
            theme: Optional theme to apply
            template_name: Optional name of a template presentation in artifact store to use as base

        Returns:
            PresentationMetadata for the new presentation
        """
        if template_name:
            # First check builtin templates via TemplateManager
            from ..templates import TemplateManager
            template_manager = TemplateManager()
            template_data = await template_manager.get_template_data(template_name)

            if template_data:
                # Create presentation from builtin template
                buffer = io.BytesIO(template_data)
                prs = await asyncio.to_thread(Presentation, buffer)
                logger.info(f"Created presentation from builtin template: {template_name}")

                # Remove all existing slides from the template - we only want the layouts/master
                # Users will add slides using pptx_add_slide_from_template
                slide_ids_to_remove = list(range(len(prs.slides) - 1, -1, -1))
                for slide_idx in slide_ids_to_remove:
                    rId = prs.slides._sldIdLst[slide_idx].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[slide_idx]
                logger.info(f"Removed {len(slide_ids_to_remove)} template slides, keeping only layouts")
            else:
                # Fallback to artifact store template
                template_data = await self._load_template_from_store(template_name)
                if template_data:
                    # Create presentation from template bytes
                    buffer = io.BytesIO(template_data)
                    prs = await asyncio.to_thread(Presentation, buffer)
                    logger.info(f"Created presentation from artifact store template: {template_name}")

                    # Remove all existing slides from the template - we only want the layouts/master
                    slide_ids_to_remove = list(range(len(prs.slides) - 1, -1, -1))
                    for slide_idx in slide_ids_to_remove:
                        rId = prs.slides._sldIdLst[slide_idx].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[slide_idx]
                    logger.info(f"Removed {len(slide_ids_to_remove)} template slides, keeping only layouts")
                else:
                    # Fallback to blank presentation if template not found
                    logger.warning(f"Template {template_name} not found in builtin or artifact store, creating blank presentation")
                    prs = Presentation()
        else:
            # Create blank presentation
            prs = Presentation()
            logger.info("Created blank presentation")

        self._presentations[name] = prs
        self._current_presentation = name

        # Auto-save to artifact store (creates namespace)
        saved = await self._save_to_store(name, prs)

        # Create metadata
        metadata = PresentationMetadata(
            name=name,
            slide_count=len(prs.slides),
            theme=theme,
            vfs_path=self.get_artifact_uri(name),
            namespace_id=self.get_namespace_id(name),
            is_saved=saved,
            template_path=template_name,
        )
        self._metadata[name] = metadata

        return metadata

    async def get(
        self, name: str | None = None
    ) -> tuple[PresentationType, PresentationMetadata] | None:
        """
        Get a presentation and its metadata by name.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Tuple of (Presentation, PresentationMetadata) or None if not found
        """
        pres_name = name or self._current_presentation
        if not pres_name:
            return None

        # Check memory first
        if pres_name in self._presentations:
            prs = self._presentations[pres_name]
            metadata = self._metadata.get(pres_name)

            # Create metadata if missing (for backwards compatibility)
            if not metadata:
                metadata = PresentationMetadata(
                    name=pres_name,
                    slide_count=len(prs.slides),
                    vfs_path=self.get_artifact_uri(pres_name),
                    namespace_id=self.get_namespace_id(pres_name),
                    is_saved=True,
                )
                self._metadata[pres_name] = metadata

            return (prs, metadata)

        # Try loading from artifact store
        loaded_prs = await self._load_from_store(pres_name)
        if loaded_prs is not None:
            self._presentations[pres_name] = loaded_prs

            # Create metadata
            metadata = PresentationMetadata(
                name=pres_name,
                slide_count=len(loaded_prs.slides),
                vfs_path=self.get_artifact_uri(pres_name),
                namespace_id=self.get_namespace_id(pres_name),
                is_saved=True,
            )
            self._metadata[pres_name] = metadata

            return (loaded_prs, metadata)

        return None

    def get_presentation(self, name: str | None = None) -> PresentationType | None:
        """
        Get just the presentation object (synchronous).

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Presentation object or None if not found
        """
        pres_name = name or self._current_presentation
        if not pres_name:
            return None
        return self._presentations.get(pres_name)

    def get_metadata(self, name: str | None = None) -> PresentationMetadata | None:
        """
        Get just the metadata (synchronous).

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            PresentationMetadata or None if not found
        """
        pres_name = name or self._current_presentation
        if not pres_name:
            return None
        return self._metadata.get(pres_name)

    async def save(self, name: str | None = None) -> bool:
        """
        Save presentation to artifact store.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            True if successful, False otherwise
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return False

        return await self._save_to_store(pres_name, self._presentations[pres_name])

    async def update(self, name: str | None = None) -> bool:
        """
        Update presentation in artifact store after modifications.

        This should be called after any modification to ensure persistence.
        Also updates metadata (slide count, modified time, etc.).

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            True if successful, False otherwise
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return False

        # Update metadata before saving
        prs = self._presentations[pres_name]
        metadata = self._metadata.get(pres_name)
        if metadata:
            metadata.slide_count = len(prs.slides)
            metadata.modified_at = datetime.now()

        return await self._save_to_store(pres_name, prs)

    async def delete(self, name: str) -> bool:
        """
        Delete a presentation from memory and artifact store.

        Args:
            name: Presentation name

        Returns:
            True if successful, False otherwise
        """
        if name not in self._presentations:
            return False

        del self._presentations[name]
        if name in self._metadata:
            del self._metadata[name]

        # Update current if we deleted it
        if self._current_presentation == name:
            self._current_presentation = (
                next(iter(self._presentations), None) if self._presentations else None
            )

        # Delete from artifact store
        await self._delete_from_store(name)

        return True

    async def list_presentations(self) -> ListPresentationsResponse:
        """
        List all presentations with metadata.

        Returns:
            ListPresentationsResponse with presentation info
        """
        presentations: list[PresentationInfo] = []

        # List from memory (artifact store presentations are tracked via _namespace_ids)
        for name, prs in self._presentations.items():
            metadata = self._metadata.get(name)
            presentations.append(
                PresentationInfo(
                    name=name,
                    slide_count=len(prs.slides),
                    is_current=(name == self._current_presentation),
                    file_path=self.get_artifact_uri(name)
                    if metadata and metadata.is_saved
                    else None,
                    namespace_id=self.get_namespace_id(name),
                )
            )

        return ListPresentationsResponse(
            presentations=presentations,
            total=len(presentations),
            current=self._current_presentation,
        )

    async def set_current(self, name: str) -> bool:
        """
        Set the current presentation.

        Args:
            name: Presentation name

        Returns:
            True if successful, False if presentation not found
        """
        if name not in self._presentations:
            # Try loading from artifact store
            prs = await self._load_from_store(name)
            if prs:
                self._presentations[name] = prs
            else:
                return False

        self._current_presentation = name
        return True

    def get_current_name(self) -> str | None:
        """Get the name of the current presentation."""
        return self._current_presentation

    def update_slide_metadata(self, slide_index: int) -> None:
        """
        Update metadata for a slide after modifications.

        Args:
            slide_index: Index of the slide to update
        """
        if not self._current_presentation:
            return

        prs = self._presentations.get(self._current_presentation)
        metadata = self._metadata.get(self._current_presentation)

        if not prs or not metadata:
            return

        # Ensure we have enough slide metadata entries
        while len(metadata.slides) <= slide_index:
            metadata.slides.append(SlideMetadata(index=len(metadata.slides), layout="Blank"))

        # Get the slide
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            slide_meta = metadata.slides[slide_index]

            # Update metadata from slide
            slide_meta.shape_count = len(slide.shapes)
            slide_meta.has_title = slide.shapes.title is not None
            if slide_meta.has_title and slide.shapes.title:
                slide_meta.title_text = slide.shapes.title.text

            # Check for charts, tables, images
            for shape in slide.shapes:
                if hasattr(shape, "chart") and shape.has_chart:
                    slide_meta.has_chart = True
                if hasattr(shape, "table") and shape.has_table:
                    slide_meta.has_table = True
                if shape.shape_type == 13:  # Picture type
                    slide_meta.has_images = True

            metadata.update_modified()

    async def export_base64(self, name: str | None = None) -> str | None:
        """
        Export presentation as base64.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Base64-encoded presentation data or None
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return None

        prs = self._presentations[pres_name]
        try:
            buffer = io.BytesIO()
            await asyncio.to_thread(prs.save, buffer)
            buffer.seek(0)
            return base64.b64encode(buffer.read()).decode("utf-8")
        except Exception as e:
            logger.error(f"Failed to export as base64: {e}")
            return None

    async def import_base64(self, data: str, name: str, as_template: bool = False) -> bool:
        """
        Import presentation from base64.

        Args:
            data: Base64-encoded presentation data
            name: Name for the imported presentation
            as_template: If True, saves as a template (doesn't set as current)

        Returns:
            True if successful, False otherwise
        """
        try:
            buffer = io.BytesIO(base64.b64decode(data))
            prs = await asyncio.to_thread(Presentation, buffer)
            self._presentations[name] = prs

            if not as_template:
                self._current_presentation = name

            # Auto-save to artifact store
            await self._save_to_store(name, prs)

            # Create metadata
            metadata = PresentationMetadata(
                name=name,
                slide_count=len(prs.slides),
                vfs_path=self.get_artifact_uri(name),
                namespace_id=self.get_namespace_id(name),
                is_saved=True,
            )
            self._metadata[name] = metadata

            logger.info(f"Imported presentation: {name}" + (" (as template)" if as_template else ""))
            return True
        except Exception as e:
            logger.error(f"Failed to import from base64: {e}")
            return False

    async def import_template(self, file_path: str, template_name: str) -> bool:
        """
        Import a PowerPoint file as a template into the artifact store.

        Args:
            file_path: Path to the PowerPoint file
            template_name: Name to save the template as

        Returns:
            True if successful, False otherwise
        """
        store = self._get_store()
        if not store:
            logger.error("No artifact store available for template import")
            return False

        from chuk_mcp_server import NamespaceType, StorageScope

        try:
            # Read the PowerPoint file
            with open(file_path, "rb") as f:
                data = await asyncio.to_thread(f.read)

            # Verify it's a valid presentation
            buffer = io.BytesIO(data)
            prs = await asyncio.to_thread(Presentation, buffer)
            logger.info(f"Validated template file: {file_path} ({len(prs.slides)} slides)")

            # Create namespace for template
            safe_name = self._sanitize_name(template_name)
            namespace_info = await store.create_namespace(
                type=NamespaceType.BLOB,
                scope=StorageScope.SESSION,
                name=f"{self.base_path}/templates/{safe_name}",
                metadata={
                    "mime_type": self.PPTX_MIME_TYPE,
                    "template_name": template_name,
                    "file_extension": ".pptx",
                    "is_template": True,
                    "slide_count": len(prs.slides),
                },
            )
            self._namespace_ids[template_name] = namespace_info.namespace_id

            # Write template data
            await store.write_namespace(namespace_info.namespace_id, data=data)

            # Store in memory for immediate use
            self._presentations[template_name] = prs

            # Create metadata
            metadata = PresentationMetadata(
                name=template_name,
                slide_count=len(prs.slides),
                vfs_path=self.get_artifact_uri(template_name),
                namespace_id=namespace_info.namespace_id,
                is_saved=True,
            )
            self._metadata[template_name] = metadata

            logger.info(f"Imported template: {template_name} ({namespace_info.namespace_id})")
            return True
        except Exception as e:
            logger.error(f"Failed to import template: {e}")
            return False

    def clear_all(self) -> None:
        """Clear all presentations from memory."""
        self._presentations.clear()
        self._metadata.clear()
        self._namespace_ids.clear()
        self._current_presentation = None

        # Note: This doesn't delete from artifact store, only clears memory
        # Artifact store presentations persist based on session/scope
