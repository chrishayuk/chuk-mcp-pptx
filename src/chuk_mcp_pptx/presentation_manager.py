"""
Presentation Manager for PowerPoint MCP Server

Manages PowerPoint presentations with support for virtual filesystem integration.
Each presentation can be auto-saved to VFS for persistence and multi-server access.

Uses Pydantic models throughout for type safety and validation.
"""
from __future__ import annotations

import base64
import io
import json
import logging
from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING
from pptx import Presentation

from .models import (
    PresentationMetadata,
    SlideMetadata,
    PresentationInfo,
    ListPresentationsResponse,
)

if TYPE_CHECKING:
    from chuk_virtual_fs import AsyncVirtualFileSystem

logger = logging.getLogger(__name__)


class PresentationManager:
    """
    Manages PowerPoint presentations with VFS integration.

    Uses chuk-virtual-fs for flexible storage (file, memory, sqlite, s3).
    Presentations and metadata are Pydantic models for type safety.
    """

    def __init__(self, vfs: "AsyncVirtualFileSystem", base_path: str = "presentations"):
        """
        Initialize the presentation manager.

        Args:
            vfs: Virtual filesystem instance for file operations
            base_path: Base directory in VFS for storing presentations
        """
        self.vfs = vfs
        self.base_path = base_path
        self._presentations: dict[str, Presentation] = {}
        self._metadata: dict[str, PresentationMetadata] = {}
        self._current_presentation: str | None = None
        self._vfs_initialized: bool = False
        logger.info(f"PresentationManager initialized with VFS, base path: {base_path}")

    async def _ensure_vfs_initialized(self) -> None:
        """Ensure VFS is initialized before use."""
        if not self._vfs_initialized:
            await self.vfs.initialize()
            self._vfs_initialized = True
            logger.info("VFS initialized")

    def _get_vfs_path(self, name: str) -> str:
        """Get the VFS path for a presentation."""
        # Sanitize the name to prevent directory traversal
        safe_name = "".join(c for c in name if c.isalnum() or c in ('-', '_'))
        if not safe_name:
            safe_name = "presentation"
        return f"{self.base_path}/{safe_name}.pptx"

    async def _save_to_vfs(self, name: str, prs: Presentation) -> bool:
        """
        Save presentation to VFS.

        Args:
            name: Presentation name
            prs: Presentation object

        Returns:
            True if successful, False otherwise
        """
        try:
            # Ensure VFS is initialized
            await self._ensure_vfs_initialized()

            # Ensure base directory exists
            if not await self.vfs.exists(self.base_path):
                await self.vfs.mkdir(self.base_path)
                logger.info(f"Created VFS directory: {self.base_path}")

            # Convert presentation to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            data = buffer.read()

            # Save to VFS
            file_path = self._get_vfs_path(name)
            await self.vfs.write_file(file_path, data)
            logger.info(f"Saved presentation to VFS: {file_path}")
            return True
        except Exception as e:
            logger.error(f"Failed to save to VFS: {e}")
            return False

    async def _load_from_vfs(self, name: str) -> Presentation | None:
        """
        Load presentation from VFS.

        Args:
            name: Presentation name

        Returns:
            Presentation object or None if not found
        """
        try:
            # Ensure VFS is initialized
            await self._ensure_vfs_initialized()

            file_path = self._get_vfs_path(name)

            # Check if file exists
            if not await self.vfs.exists(file_path):
                logger.debug(f"Presentation not found in VFS: {file_path}")
                return None

            # Read from VFS
            data = await self.vfs.read_file(file_path)
            buffer = io.BytesIO(data)
            prs = Presentation(buffer)
            logger.info(f"Loaded presentation from VFS: {file_path}")
            return prs
        except Exception as e:
            logger.error(f"Failed to load from VFS: {e}")
            return None

    async def _delete_from_vfs(self, name: str) -> bool:
        """
        Delete presentation from VFS.

        Args:
            name: Presentation name

        Returns:
            True if successful, False otherwise
        """
        try:
            file_path = self._get_vfs_path(name)

            if await self.vfs.exists(file_path):
                await self.vfs.delete(file_path)
                logger.info(f"Deleted presentation from VFS: {file_path}")
                return True
            else:
                logger.warning(f"Presentation not found in VFS for deletion: {file_path}")
                return False
        except Exception as e:
            logger.error(f"Failed to delete from VFS: {e}")
            return False

    async def create(self, name: str, theme: str | None = None) -> PresentationMetadata:
        """
        Create a new presentation.

        Args:
            name: Presentation name
            theme: Optional theme to apply

        Returns:
            PresentationMetadata for the new presentation
        """
        prs = Presentation()
        self._presentations[name] = prs
        self._current_presentation = name

        # Create metadata
        metadata = PresentationMetadata(
            name=name,
            slide_count=len(prs.slides),
            theme=theme,
            vfs_path=self._get_vfs_path(name),
            is_saved=False,
        )
        self._metadata[name] = metadata

        # Auto-save to VFS
        saved = await self._save_to_vfs(name, prs)
        if saved:
            metadata.is_saved = True

        return metadata

    async def get(self, name: str | None = None) -> tuple[Presentation, PresentationMetadata] | None:
        """
        Get a presentation and its metadata by name.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Tuple of (Presentation, PresentationMetadata) or None if not found
        """
        pres_name = name or self._current_presentation
        if not pres_name:
            return None

        # Check memory first
        if pres_name in self._presentations:
            prs = self._presentations[pres_name]
            metadata = self._metadata.get(pres_name)

            # Create metadata if missing (for backwards compatibility)
            if not metadata:
                metadata = PresentationMetadata(
                    name=pres_name,
                    slide_count=len(prs.slides),
                    vfs_path=self._get_vfs_path(pres_name),
                    is_saved=True,
                )
                self._metadata[pres_name] = metadata

            return (prs, metadata)

        # Try loading from VFS
        prs = await self._load_from_vfs(pres_name)
        if prs:
            self._presentations[pres_name] = prs

            # Create metadata
            metadata = PresentationMetadata(
                name=pres_name,
                slide_count=len(prs.slides),
                vfs_path=self._get_vfs_path(pres_name),
                is_saved=True,
            )
            self._metadata[pres_name] = metadata

            return (prs, metadata)

        return None

    def get_presentation(self, name: str | None = None) -> Presentation | None:
        """
        Get just the presentation object (synchronous).

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Presentation object or None if not found
        """
        pres_name = name or self._current_presentation
        if not pres_name:
            return None
        return self._presentations.get(pres_name)

    def get_metadata(self, name: str | None = None) -> PresentationMetadata | None:
        """
        Get just the metadata (synchronous).

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            PresentationMetadata or None if not found
        """
        pres_name = name or self._current_presentation
        if not pres_name:
            return None
        return self._metadata.get(pres_name)

    async def save(self, name: str | None = None) -> bool:
        """
        Save presentation to VFS.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            True if successful, False otherwise
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return False

        return await self._save_to_vfs(pres_name, self._presentations[pres_name])

    async def update(self, name: str | None = None) -> bool:
        """
        Update presentation in VFS after modifications.

        This should be called after any modification to ensure VFS persistence.
        Also updates metadata (slide count, modified time, etc.).

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            True if successful, False otherwise
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return False

        # Update metadata before saving
        prs = self._presentations[pres_name]
        metadata = self._metadata.get(pres_name)
        if metadata:
            metadata.slide_count = len(prs.slides)
            metadata.modified_at = datetime.now()

        return await self._save_to_vfs(pres_name, prs)

    async def delete(self, name: str) -> bool:
        """
        Delete a presentation from memory and VFS.

        Args:
            name: Presentation name

        Returns:
            True if successful, False otherwise
        """
        if name not in self._presentations:
            return False

        del self._presentations[name]

        # Update current if we deleted it
        if self._current_presentation == name:
            self._current_presentation = next(iter(self._presentations), None) if self._presentations else None

        # Delete from VFS
        await self._delete_from_vfs(name)

        return True

    async def list_presentations(self) -> ListPresentationsResponse:
        """
        List all presentations with metadata.

        Returns:
            ListPresentationsResponse with presentation info
        """
        presentations: list[PresentationInfo] = []

        # List from memory
        for name, prs in self._presentations.items():
            metadata = self._metadata.get(name)
            presentations.append(
                PresentationInfo(
                    name=name,
                    slide_count=len(prs.slides),
                    is_current=(name == self._current_presentation),
                    file_path=self._get_vfs_path(name) if metadata and metadata.is_saved else None,
                )
            )

        # List from VFS (presentations not yet loaded into memory)
        try:
            if await self.vfs.exists(self.base_path):
                files = await self.vfs.ls(self.base_path)
                existing_names = {p.name for p in presentations}

                for file in files:
                    if file.endswith('.pptx'):
                        name = file[:-5]  # Remove .pptx extension
                        if name not in existing_names:
                            # Load to get slide count
                            prs = await self._load_from_vfs(name)
                            if prs:
                                self._presentations[name] = prs
                                metadata = PresentationMetadata(
                                    name=name,
                                    slide_count=len(prs.slides),
                                    vfs_path=self._get_vfs_path(name),
                                    is_saved=True,
                                )
                                self._metadata[name] = metadata

                                presentations.append(
                                    PresentationInfo(
                                        name=name,
                                        slide_count=len(prs.slides),
                                        is_current=False,
                                        file_path=self._get_vfs_path(name),
                                    )
                                )
        except Exception as e:
            logger.error(f"Failed to list VFS presentations: {e}")

        return ListPresentationsResponse(
            presentations=presentations,
            total=len(presentations),
            current=self._current_presentation,
        )

    async def set_current(self, name: str) -> bool:
        """
        Set the current presentation.

        Args:
            name: Presentation name

        Returns:
            True if successful, False if presentation not found
        """
        if name not in self._presentations:
            # Try loading from VFS
            prs = await self._load_from_vfs(name)
            if prs:
                self._presentations[name] = prs
            else:
                return False

        self._current_presentation = name
        return True

    def get_current_name(self) -> str | None:
        """Get the name of the current presentation."""
        return self._current_presentation

    def update_slide_metadata(self, slide_index: int) -> None:
        """
        Update metadata for a slide after modifications.

        Args:
            slide_index: Index of the slide to update
        """
        if not self._current_presentation:
            return

        prs = self._presentations.get(self._current_presentation)
        metadata = self._metadata.get(self._current_presentation)

        if not prs or not metadata:
            return

        # Ensure we have enough slide metadata entries
        while len(metadata.slides) <= slide_index:
            metadata.slides.append(
                SlideMetadata(index=len(metadata.slides), layout="Blank")
            )

        # Get the slide
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            slide_meta = metadata.slides[slide_index]

            # Update metadata from slide
            slide_meta.shape_count = len(slide.shapes)
            slide_meta.has_title = slide.shapes.title is not None
            if slide_meta.has_title and slide.shapes.title:
                slide_meta.title_text = slide.shapes.title.text

            # Check for charts, tables, images
            for shape in slide.shapes:
                if hasattr(shape, "chart") and shape.has_chart:
                    slide_meta.has_chart = True
                if hasattr(shape, "table") and shape.has_table:
                    slide_meta.has_table = True
                if shape.shape_type == 13:  # Picture type
                    slide_meta.has_images = True

            metadata.update_modified()

    def export_base64(self, name: str | None = None) -> str | None:
        """
        Export presentation as base64.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Base64-encoded presentation data or None
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return None

        prs = self._presentations[pres_name]
        try:
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return base64.b64encode(buffer.read()).decode('utf-8')
        except Exception as e:
            logger.error(f"Failed to export as base64: {e}")
            return None

    async def import_base64(self, data: str, name: str) -> bool:
        """
        Import presentation from base64.

        Args:
            data: Base64-encoded presentation data
            name: Name for the imported presentation

        Returns:
            True if successful, False otherwise
        """
        try:
            buffer = io.BytesIO(base64.b64decode(data))
            prs = Presentation(buffer)
            self._presentations[name] = prs
            self._current_presentation = name

            # Auto-save to VFS
            await self._save_to_vfs(name, prs)

            return True
        except Exception as e:
            logger.error(f"Failed to import from base64: {e}")
            return False

    def clear_all(self):
        """Clear all presentations from memory."""
        self._presentations.clear()
        self._current_presentation = None

        # Note: This doesn't delete from VFS, only clears memory
        # VFS presentations can still be loaded on demand

    async def export_to_vfs(self, name: str, file_path: str) -> str:
        """
        Export presentation to a specific VFS path.

        Args:
            name: Presentation name
            file_path: Destination path in VFS

        Returns:
            Success message or error
        """
        if name not in self._presentations:
            return f"Error: Presentation '{name}' not found"

        try:
            prs = self._presentations[name]
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            data = buffer.read()

            await self.vfs.write_file(file_path, data)
            return f"Exported presentation '{name}' to {file_path}"
        except Exception as e:
            logger.error(f"Failed to export to VFS: {e}")
            return f"Error: {e}"

    async def import_from_vfs(self, file_path: str, name: str) -> bool:
        """
        Import presentation from a specific VFS path.

        Args:
            file_path: Source path in VFS
            name: Name for the imported presentation

        Returns:
            True if successful, False otherwise
        """
        try:
            if not await self.vfs.exists(file_path):
                logger.error(f"File not found in VFS: {file_path}")
                return False

            data = await self.vfs.read_file(file_path)
            buffer = io.BytesIO(data)
            prs = Presentation(buffer)
            self._presentations[name] = prs
            self._current_presentation = name

            # Save to standard location
            await self._save_to_vfs(name, prs)

            return True
        except Exception as e:
            logger.error(f"Failed to import from VFS: {e}")
            return False
