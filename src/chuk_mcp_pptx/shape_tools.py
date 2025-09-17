"""
Shape and SmartArt Tools for PowerPoint MCP Server

Provides async MCP tools for creating shapes, connectors, and SmartArt-like diagrams.
"""
import asyncio
from typing import List, Optional, Tuple, Dict
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def register_shape_tools(mcp, manager):
    """Register all shape and SmartArt tools with the MCP server."""
    
    from .shape_utils import add_shape, add_connector, add_smart_art
    from .layout_helpers import validate_position
    from .legacy_themes import create_code_block, apply_theme, THEMES
    
    @mcp.tool
    async def pptx_add_shape(
        slide_index: int,
        shape_type: str,
        left: float = 2.0,
        top: float = 2.0,
        width: float = 3.0,
        height: float = 2.0,
        text: Optional[str] = None,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
        line_width: float = 1.0,
        presentation: Optional[str] = None
    ) -> str:
        """
        Add a shape to a slide with customizable properties.
        
        Creates various shapes like rectangles, circles, arrows, stars, etc.
        
        Args:
            slide_index: Index of the slide to add shape to (0-based)
            shape_type: Type of shape. Options:
                - "rectangle", "rounded_rectangle"
                - "oval", "diamond"
                - "triangle", "arrow"
                - "star", "hexagon"
                - "chevron", "plus"
                - "callout"
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Optional text to display inside shape
            fill_color: Fill color as hex string (e.g., "#FF0000")
            line_color: Border color as hex string
            line_width: Border width in points
            presentation: Name of presentation (uses current if not specified)
            
        Returns:
            Success message confirming shape addition
            
        Example:
            await pptx_add_shape(
                slide_index=1,
                shape_type="star",
                left=3.0,
                top=2.0,
                width=2.0,
                height=2.0,
                text="Success!",
                fill_color="#FFD700"
            )
        """
        def _add_shape():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found. Create one first with pptx_create()"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            # Remove content placeholders that might interfere with shapes
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type'):
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, 'placeholder_format'):
                            if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                                continue
                        # Mark content placeholders for removal
                        placeholders_to_remove.append(shape)
            
            # Remove the placeholders
            for placeholder in placeholders_to_remove:
                sp = placeholder._element
                sp.getparent().remove(sp)
            
            # Validate position
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )
            
            # Convert hex colors to RGB tuples
            fill_rgb = None
            if fill_color:
                color_str = fill_color
                if color_str.startswith("#"):
                    color_str = color_str[1:]
                fill_rgb = (
                    int(color_str[0:2], 16),
                    int(color_str[2:4], 16),
                    int(color_str[4:6], 16)
                )
            
            line_rgb = None
            if line_color:
                color_str = line_color
                if color_str.startswith("#"):
                    color_str = color_str[1:]
                line_rgb = (
                    int(color_str[0:2], 16),
                    int(color_str[2:4], 16),
                    int(color_str[4:6], 16)
                )
            
            try:
                shape = add_shape(
                    slide, shape_type,
                    validated_left, validated_top,
                    validated_width, validated_height,
                    text=text,
                    fill_color=fill_rgb,
                    line_color=line_rgb,
                    line_width=line_width
                )
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                # Report if position was adjusted
                position_note = ""
                if (validated_left != left or validated_top != top or 
                    validated_width != width or validated_height != height):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"
                
                return f"Added {shape_type} shape to slide {slide_index}{position_note}"
                
            except Exception as e:
                return f"Error adding shape: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_shape)
    
    @mcp.tool
    async def pptx_add_arrow(
        slide_index: int,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        connector_type: str = "straight",
        line_color: Optional[str] = "#000000",
        line_width: float = 2.0,
        arrow_start: bool = False,
        arrow_end: bool = True,
        presentation: Optional[str] = None
    ) -> str:
        """
        Add an arrow or connector line to a slide.
        
        Creates connectors between points with optional arrowheads.
        
        Args:
            slide_index: Index of the slide to add arrow to
            start_x: Starting X position in inches
            start_y: Starting Y position in inches
            end_x: Ending X position in inches
            end_y: Ending Y position in inches
            connector_type: Type of connector ("straight", "elbow", "curved")
            line_color: Line color as hex string
            line_width: Line width in points
            arrow_start: Whether to add arrowhead at start
            arrow_end: Whether to add arrowhead at end
            presentation: Name of presentation
            
        Returns:
            Success message confirming arrow addition
            
        Example:
            await pptx_add_arrow(
                slide_index=1,
                start_x=2.0,
                start_y=2.0,
                end_x=5.0,
                end_y=3.0,
                connector_type="straight",
                arrow_end=True
            )
        """
        def _add_arrow():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            # Convert hex color to RGB
            line_rgb = (0, 0, 0)
            if line_color and line_color.startswith("#"):
                color = line_color[1:]
                line_rgb = (
                    int(color[0:2], 16),
                    int(color[2:4], 16),
                    int(color[4:6], 16)
                )
            
            try:
                connector = add_connector(
                    slide,
                    start_x, start_y,
                    end_x, end_y,
                    connector_type=connector_type,
                    line_color=line_rgb,
                    line_width=line_width,
                    arrow_start=arrow_start,
                    arrow_end=arrow_end
                )
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                return f"Added {connector_type} arrow to slide {slide_index}"
                
            except Exception as e:
                return f"Error adding arrow: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_arrow)
    
    @mcp.tool
    async def pptx_add_smart_art(
        slide_index: int,
        art_type: str,
        items: List[str],
        title: Optional[str] = None,
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 3.0,
        color_scheme: str = "modern_blue",
        presentation: Optional[str] = None
    ) -> str:
        """
        Add a SmartArt-style diagram to a slide.
        
        Creates professional diagrams like process flows, cycles, hierarchies, etc.
        
        Args:
            slide_index: Index of the slide to add SmartArt to
            art_type: Type of SmartArt diagram:
                - "process" - Sequential process flow
                - "cycle" - Circular/cyclical process
                - "hierarchy" - Organizational hierarchy
                - "list" - Bulleted list with shapes
                - "relationship" - Relationship diagram
                - "pyramid" - Pyramid hierarchy
            items: List of text items for the diagram
            title: Optional title for the diagram
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            color_scheme: Color scheme ("modern_blue", "corporate_gray", "warm_orange")
            presentation: Name of presentation
            
        Returns:
            Success message confirming SmartArt addition
            
        Example:
            await pptx_add_smart_art(
                slide_index=2,
                art_type="process",
                items=["Research", "Design", "Develop", "Test", "Deploy"],
                title="Development Process",
                color_scheme="modern_blue"
            )
        """
        def _add_smart():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            # Remove content placeholders that might interfere with SmartArt
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type'):
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, 'placeholder_format'):
                            if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                                continue
                        # Mark content placeholders for removal
                        placeholders_to_remove.append(shape)
            
            # Remove the placeholders
            for placeholder in placeholders_to_remove:
                sp = placeholder._element
                sp.getparent().remove(sp)
            
            # Validate position
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )
            
            # Add title if provided
            if title:
                from .shape_utils import add_shape
                title_shape = add_shape(
                    slide, "rectangle",
                    validated_left, validated_top - 0.5,
                    validated_width, 0.4,
                    text=title,
                    fill_color=None,
                    line_color=None
                )
                # Adjust top position for the SmartArt
                validated_top += 0.2
                validated_height -= 0.7
            
            try:
                shapes = add_smart_art(
                    slide, art_type, items,
                    validated_left, validated_top,
                    validated_width, validated_height,
                    color_scheme=color_scheme
                )
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                # Report if position was adjusted
                position_note = ""
                if (validated_left != left or validated_top != top or 
                    validated_width != width or validated_height != height):
                    position_note = f" (position adjusted)"
                
                return f"Added {art_type} SmartArt with {len(items)} items to slide {slide_index}{position_note}"
                
            except Exception as e:
                return f"Error adding SmartArt: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_smart)
    
    @mcp.tool
    async def pptx_add_code_block(
        slide_index: int,
        code: str,
        language: str = "python",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 3.0,
        theme: str = "dark_modern",
        presentation: Optional[str] = None
    ) -> str:
        """
        Add a code block to a slide with syntax highlighting appearance.
        
        Creates a formatted code block with monospace font and theme colors.
        
        Args:
            slide_index: Index of the slide to add code to (0-based)
            code: The code content to display
            language: Programming language (for label)
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            theme: Theme to use for code block styling
            presentation: Name of presentation (uses current if not specified)
            
        Returns:
            Success message confirming code block addition
            
        Example:
            await pptx_add_code_block(
                slide_index=1,
                code="def hello_world():\\n    print('Hello, World!')",
                language="python",
                theme="dark_purple"
            )
        """
        def _add_code():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range"
            
            slide = prs.slides[slide_index]
            
            try:
                # Remove content placeholders
                from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
                placeholders_to_remove = []
                for shape in slide.shapes:
                    if hasattr(shape, 'shape_type'):
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            if hasattr(shape, 'placeholder_format'):
                                if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                                    continue
                            placeholders_to_remove.append(shape)
                
                for shape in placeholders_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)
                
                # Add code block
                code_shape = create_code_block(
                    slide, code, language,
                    left, top, width, height,
                    theme_name=theme
                )
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                return f"Added {language} code block to slide {slide_index} with {theme} theme"
                
            except Exception as e:
                return f"Error adding code block: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_code)
    
    @mcp.tool
    async def pptx_apply_theme(
        slide_index: Optional[int] = None,
        theme: str = "dark_modern",
        presentation: Optional[str] = None
    ) -> str:
        """
        Apply a beautiful theme to slides.
        
        Available themes:
        - dark_modern: Dark with orange accents
        - dark_blue: Dark with blue accents  
        - dark_purple: Dark with purple accents
        - dark_green: Dark with green accents
        - light_minimal: Clean light theme
        - light_warm: Warm light theme
        - cyberpunk: Neon cyberpunk style
        - gradient_sunset: Gradient sunset colors
        
        Args:
            slide_index: Index of slide to theme (None for all slides)
            theme: Name of theme to apply
            presentation: Name of presentation (uses current if not specified)
            
        Returns:
            Success message confirming theme application
        """
        def _apply_theme():
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found"
            
            if theme not in THEMES:
                return f"Error: Unknown theme '{theme}'. Available: {', '.join(THEMES.keys())}"
            
            try:
                if slide_index is not None:
                    if slide_index >= len(prs.slides):
                        return f"Error: Slide index {slide_index} out of range"
                    slides = [prs.slides[slide_index]]
                else:
                    slides = prs.slides
                
                for slide in slides:
                    apply_theme(slide, theme)
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                slide_msg = f"slide {slide_index}" if slide_index is not None else "all slides"
                return f"Applied {theme} theme to {slide_msg}"
                
            except Exception as e:
                return f"Error applying theme: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _apply_theme)
    
    @mcp.tool
    async def pptx_list_themes() -> str:
        """
        List all available themes with descriptions.
        
        Returns:
            List of available themes and their characteristics
        """
        theme_list = []
        for key, theme in THEMES.items():
            theme_list.append(f"• {key}: {theme['name']} - Primary: {theme['primary']}, Fonts: {theme['font_body']}")
        
        return "Available themes:\n" + "\n".join(theme_list)
    
    # Return tools for external access
    return {
        'pptx_add_shape': pptx_add_shape,
        'pptx_add_arrow': pptx_add_arrow,
        'pptx_add_smart_art': pptx_add_smart_art,
        'pptx_add_code_block': pptx_add_code_block,
        'pptx_apply_theme': pptx_apply_theme,
        'pptx_list_themes': pptx_list_themes
    }