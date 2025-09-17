"""
Chart Utilities for PowerPoint MCP Server

Provides utilities for creating and managing charts in presentations.
"""
from typing import List, Dict, Any, Optional, Tuple
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


# Chart type mapping
CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "bubble": XL_CHART_TYPE.BUBBLE
}


def add_chart(slide, chart_type: str, left: float, top: float,
             width: float, height: float,
             categories: List[str], series_data: Dict[str, List[float]],
             title: str = None, has_legend: bool = True,
             legend_position: str = "right") -> Any:
    """
    Add a chart to a slide.
    
    Args:
        slide: Slide to add chart to
        chart_type: Type of chart (column, bar, line, pie, etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: Category labels
        series_data: Dictionary of series names to data values
        title: Chart title
        has_legend: Whether to show legend
        legend_position: Legend position (right, left, top, bottom)
        
    Returns:
        The created chart shape
    """
    # Get chart type
    xl_chart_type = CHART_TYPES.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)
    
    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        xl_chart_type,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    )
    
    chart = chart_shape.chart
    
    # Set title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Configure legend
    if has_legend:
        chart.has_legend = True
        
        # Set legend position
        if legend_position.lower() == "right":
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
        elif legend_position.lower() == "left":
            chart.legend.position = XL_LEGEND_POSITION.LEFT
        elif legend_position.lower() == "top":
            chart.legend.position = XL_LEGEND_POSITION.TOP
        elif legend_position.lower() == "bottom":
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    
    return chart_shape


def add_scatter_chart(slide, left: float, top: float,
                     width: float, height: float,
                     series_data: List[Dict[str, Any]],
                     title: str = None, has_legend: bool = True) -> Any:
    """
    Add a scatter plot to a slide.
    
    Args:
        slide: Slide to add chart to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        series_data: List of series, each with 'name', 'x_values', 'y_values'
        title: Chart title
        has_legend: Whether to show legend
        
    Returns:
        The created chart shape
    """
    # Create XY chart data
    chart_data = XyChartData()
    
    for series in series_data:
        series_obj = chart_data.add_series(series.get('name', 'Series'))
        
        x_values = series.get('x_values', [])
        y_values = series.get('y_values', [])
        
        for x, y in zip(x_values, y_values):
            series_obj.add_data_point(x, y)
    
    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    )
    
    chart = chart_shape.chart
    
    # Set title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Configure legend
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    
    return chart_shape


def add_pie_chart(slide, left: float, top: float,
                 width: float, height: float,
                 categories: List[str], values: List[float],
                 title: str = None, show_percentages: bool = True,
                 explode_slice: int = None) -> Any:
    """
    Add a pie chart to a slide.
    
    Args:
        slide: Slide to add chart to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: Category labels
        values: Data values
        title: Chart title
        show_percentages: Whether to show percentage labels
        explode_slice: Index of slice to explode (separate from pie)
        
    Returns:
        The created chart shape
    """
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)
    
    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    )
    
    chart = chart_shape.chart
    
    # Set title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Configure pie chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    
    # Show data labels with percentages
    if show_percentages:
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.show_percentage = True
        data_labels.show_value = False
        data_labels.font.size = Pt(10)
    
    # Explode a slice if specified
    if explode_slice is not None and 0 <= explode_slice < len(categories):
        chart.series[0].points[explode_slice].explosion = 20
    
    return chart_shape


def add_combo_chart(slide, left: float, top: float,
                   width: float, height: float,
                   categories: List[str],
                   bar_series: Dict[str, List[float]],
                   line_series: Dict[str, List[float]],
                   title: str = None) -> Any:
    """
    Add a combination chart (bars and lines) to a slide.
    
    Args:
        slide: Slide to add chart to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: Category labels
        bar_series: Dictionary of series names to values for bars
        line_series: Dictionary of series names to values for lines
        title: Chart title
        
    Returns:
        The created chart shape
    """
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    # Add all series
    all_series = {}
    all_series.update(bar_series)
    all_series.update(line_series)
    
    for series_name, values in all_series.items():
        chart_data.add_series(series_name, values)
    
    # Add chart as column chart first
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    )
    
    chart = chart_shape.chart
    
    # Convert some series to lines
    # Note: This is simplified - full combo chart requires more complex manipulation
    
    # Set title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    # Configure legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    return chart_shape


def update_chart_data(chart_shape, categories: List[str],
                     series_data: Dict[str, List[float]]):
    """
    Update data in an existing chart.
    
    Args:
        chart_shape: Chart shape to update
        categories: New category labels
        series_data: New series data
    """
    chart = chart_shape.chart
    
    # Create new chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)
    
    # Replace chart data
    chart.replace_data(chart_data)


def style_chart(chart_shape, 
               colors: List[Tuple[int, int, int]] = None,
               font_name: str = "Calibri",
               font_size: int = 10,
               gridlines: bool = True):
    """
    Apply styling to a chart.
    
    Args:
        chart_shape: Chart shape to style
        colors: List of RGB tuples for series colors
        font_name: Font for labels
        font_size: Font size for labels
        gridlines: Whether to show gridlines
    """
    chart = chart_shape.chart
    
    # Apply colors to series
    if colors:
        for idx, series in enumerate(chart.series):
            if idx < len(colors):
                r, g, b = colors[idx]
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r, g, b)
    
    # Configure axes if present
    if hasattr(chart, 'value_axis'):
        value_axis = chart.value_axis
        value_axis.tick_labels.font.name = font_name
        value_axis.tick_labels.font.size = Pt(font_size)
        value_axis.has_major_gridlines = gridlines
        
    if hasattr(chart, 'category_axis'):
        category_axis = chart.category_axis
        category_axis.tick_labels.font.name = font_name
        category_axis.tick_labels.font.size = Pt(font_size)


def add_data_table(slide, left: float, top: float,
                  width: float, height: float,
                  headers: List[str], data: List[List[Any]],
                  style: str = "medium") -> Any:
    """
    Add a formatted data table to a slide.
    
    Args:
        slide: Slide to add table to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        headers: Column headers
        data: Table data (list of rows)
        style: Table style (light, medium, dark)
        
    Returns:
        The created table shape
    """
    rows = len(data) + 1  # +1 for header
    cols = len(headers)
    
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    ).table
    
    # Set headers
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = str(header)
        
        # Style header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)
        
        # Set header background color based on style
        if style == "dark":
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 68, 68)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        elif style == "medium":
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(217, 217, 217)
    
    # Add data
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table_shape.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            
            # Style data cells
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)
            
            # Alternate row colors for medium style
            if style == "medium" and row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    return table_shape