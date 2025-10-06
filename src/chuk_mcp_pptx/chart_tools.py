"""
Chart Tools for PowerPoint MCP Server

Provides a unified async MCP tool for creating all chart types in presentations.
Optimized for AI/LLM tool invocation with clear parameter structure.
"""
import asyncio
from typing import Dict, List, Optional, Union, Any
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def register_chart_tools(mcp, manager):
    """Register the unified chart tool with the MCP server."""
    
    from .chart_utils import add_chart, add_pie_chart, add_scatter_chart
    from .layout.helpers import validate_position, get_safe_content_area
    
    @mcp.tool
    async def pptx_add_chart(
        slide_index: int,
        chart_type: str,
        data: dict,
        title: str = "",
        left: float = 1.0,
        top: float = 2.0,
        width: float = 8.0,
        height: float = 4.5,
        options: Optional[dict] = None,
        presentation: Optional[str] = None
    ) -> str:
        """
        Add any type of chart to a slide with a unified interface.
        
        This single tool handles all chart types, making it easier for AI/LLM invocation.
        The data structure automatically adapts based on the chart_type.
        
        Args:
            slide_index: Index of the slide to add chart to (0-based)
            
            chart_type: Type of chart to create. Options:
                - "column" - Vertical bars comparing values
                - "column_stacked" - Stacked vertical bars
                - "bar" - Horizontal bars
                - "bar_stacked" - Stacked horizontal bars
                - "line" - Line graph showing trends
                - "line_markers" - Line graph with data point markers
                - "pie" - Pie chart showing proportions
                - "doughnut" - Doughnut chart (pie with hollow center)
                - "area" - Area chart showing magnitude over time
                - "area_stacked" - Stacked area chart
                - "scatter" - XY scatter plot for correlations
                - "bubble" - Bubble chart for 3D data (x, y, size)
                - "radar" - Radar/spider chart for multi-criteria
                - "radar_filled" - Filled radar chart
                - "waterfall" - Waterfall chart for incremental changes
                
            data: Chart data structure. Format depends on chart_type:
            
                For column/bar/line/area charts:
                {
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": {
                        "Revenue": [100, 120, 140, 160],
                        "Profit": [20, 25, 30, 35]
                    }
                }
                
                For pie/doughnut charts:
                {
                    "categories": ["Product A", "Product B", "Product C"],
                    "values": [45, 30, 25]
                }
                
                For scatter charts:
                {
                    "series": [
                        {
                            "name": "Dataset 1",
                            "x_values": [1, 2, 3, 4, 5],
                            "y_values": [2, 4, 6, 8, 10]
                        }
                    ]
                }
                
                For bubble charts:
                {
                    "series": [
                        {
                            "name": "Markets",
                            "points": [[10, 20, 5], [15, 25, 8], [20, 30, 12]]
                        }
                    ]
                }
                
                For radar charts:
                {
                    "categories": ["Speed", "Reliability", "Comfort", "Design"],
                    "series": {
                        "Model A": [8, 7, 9, 8],
                        "Model B": [7, 9, 7, 6]
                    }
                }
                
                For waterfall charts:
                {
                    "categories": ["Start", "Sales", "Costs", "Tax", "End"],
                    "values": [100, 50, -30, -10, 110]
                }
            
            title: Optional chart title
            left: Left position in inches (will be validated)
            top: Top position in inches (will be validated)
            width: Width in inches (will be validated)
            height: Height in inches (will be validated)
            
            options: Optional chart-specific options:
                {
                    "show_percentages": true,  # For pie charts
                    "show_legend": true,        # Show/hide legend
                    "legend_position": "right", # Legend position: right, left, top, bottom
                    "colors": ["#FF5733", "#33FF57", "#3357FF"]  # Custom colors
                }
            
            presentation: Name of presentation (uses current if not specified)
            
        Returns:
            Success message confirming chart addition, or error with guidance
            
        Examples:
            # Column chart
            await pptx_add_chart(
                slide_index=1,
                chart_type="column",
                data={
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": {
                        "Revenue": [100, 120, 140, 160],
                        "Profit": [20, 25, 30, 35]
                    }
                },
                title="Quarterly Performance"
            )
            
            # Pie chart
            await pptx_add_chart(
                slide_index=2,
                chart_type="pie",
                data={
                    "categories": ["Product A", "Product B", "Product C"],
                    "values": [45, 30, 25]
                },
                title="Market Share",
                options={"show_percentages": True}
            )
            
            # Scatter plot
            await pptx_add_chart(
                slide_index=3,
                chart_type="scatter",
                data={
                    "series": [
                        {
                            "name": "Sales Data",
                            "x_values": [10, 20, 30, 40, 50],
                            "y_values": [15, 25, 45, 35, 55]
                        }
                    ]
                },
                title="Price vs Sales Correlation"
            )
        """
        def _add_unified_chart():
            nonlocal options
            
            prs = manager.get(presentation)
            if not prs:
                return "Error: No presentation found. Create one first with pptx_create()"
            
            if slide_index >= len(prs.slides):
                return f"Error: Slide index {slide_index} out of range. Presentation has {len(prs.slides)} slides."
            
            slide = prs.slides[slide_index]
            
            # Get safe content area
            safe_area = get_safe_content_area(has_title=bool(slide.shapes.title))
            
            # Validate and adjust position to fit within slide
            validated_left, validated_top, validated_width, validated_height = validate_position(
                left, top, width, height
            )
            
            # Further adjust if position is too close to title area
            if slide.shapes.title and validated_top < safe_area['top']:
                validated_top = safe_area['top']
            
            # Default options
            if options is None:
                options = {}
            
            # Remove any overlapping placeholders (except title)
            placeholders_to_remove = []
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type'):
                    # Check if it's a placeholder (but not a title)
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Skip title placeholders
                        if hasattr(shape, 'placeholder_format'):
                            from pptx.enum.shapes import PP_PLACEHOLDER
                            if shape.placeholder_format.type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                                continue
                        
                        # Check if placeholder overlaps with chart area
                        shape_left = shape.left.inches if hasattr(shape.left, 'inches') else 0
                        shape_top = shape.top.inches if hasattr(shape.top, 'inches') else 0
                        shape_right = shape_left + (shape.width.inches if hasattr(shape.width, 'inches') else 0)
                        shape_bottom = shape_top + (shape.height.inches if hasattr(shape.height, 'inches') else 0)
                        
                        chart_right = validated_left + validated_width
                        chart_bottom = validated_top + validated_height
                        
                        # Check for overlap
                        if not (shape_right < validated_left or shape_left > chart_right or 
                               shape_bottom < validated_top or shape_top > chart_bottom):
                            placeholders_to_remove.append(shape)
            
            # Remove overlapping placeholders
            for placeholder in placeholders_to_remove:
                slide.shapes._spTree.remove(placeholder.element)
            
            try:
                # Validate chart type
                valid_types = [
                    "column", "column_stacked", "bar", "bar_stacked",
                    "line", "line_markers", "pie", "doughnut",
                    "area", "area_stacked", "scatter", "bubble",
                    "radar", "radar_filled", "waterfall"
                ]
                
                if chart_type not in valid_types:
                    return f"Error: Invalid chart_type '{chart_type}'. Valid types: {', '.join(valid_types)}"
                
                # Handle different chart types based on their data requirements
                if chart_type in ["column", "column_stacked", "bar", "bar_stacked", 
                                 "line", "line_markers", "area", "area_stacked"]:
                    # Standard category charts
                    if "categories" not in data or "series" not in data:
                        return f"Error: {chart_type} charts require 'categories' and 'series' in data. Got: {list(data.keys())}"
                    
                    categories = data["categories"]
                    series_data = data["series"]
                    
                    # Map radar_filled to the enum name
                    chart_type_mapped = chart_type
                    
                    chart_shape = add_chart(
                        slide, chart_type_mapped, validated_left, validated_top, validated_width, validated_height,
                        categories, series_data, title,
                        has_legend=options.get("show_legend", True),
                        legend_position=options.get("legend_position", "right")
                    )
                    
                elif chart_type in ["pie", "doughnut"]:
                    # Pie-type charts
                    if "categories" not in data or "values" not in data:
                        return f"Error: {chart_type} charts require 'categories' and 'values' in data. Got: {list(data.keys())}"
                    
                    categories = data["categories"]
                    values = data["values"]
                    
                    if chart_type == "pie":
                        chart_shape = add_pie_chart(
                            slide, validated_left, validated_top, validated_width, validated_height,
                            categories, values, title,
                            show_percentages=options.get("show_percentages", True)
                        )
                    else:  # doughnut
                        # Use the general add_chart for doughnut
                        series_data = {"Values": values}
                        chart_shape = add_chart(
                            slide, "doughnut", validated_left, validated_top, validated_width, validated_height,
                            categories, series_data, title
                        )
                    
                elif chart_type == "scatter":
                    # Scatter charts
                    if "series" not in data:
                        return f"Error: Scatter charts require 'series' in data with 'x_values' and 'y_values'. Got: {list(data.keys())}"
                    
                    series_data = data["series"]
                    chart_shape = add_scatter_chart(
                        slide, validated_left, validated_top, validated_width, validated_height,
                        series_data, title,
                        has_legend=options.get("show_legend", True)
                    )
                    
                elif chart_type == "bubble":
                    # Bubble charts
                    if "series" not in data:
                        return f"Error: Bubble charts require 'series' in data with 'points' as [x, y, size]. Got: {list(data.keys())}"
                    
                    from pptx.chart.data import BubbleChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    from pptx.util import Inches, Pt
                    
                    chart_data = BubbleChartData()
                    
                    for series in data["series"]:
                        series_obj = chart_data.add_series(series.get("name", "Series"))
                        for point in series.get("points", []):
                            if len(point) == 3:
                                series_obj.add_data_point(point[0], point[1], point[2])
                            else:
                                return f"Error: Bubble chart points must be [x, y, size]. Got: {point}"
                    
                    chart_shape = slide.shapes.add_chart(
                        XL_CHART_TYPE.BUBBLE,
                        Inches(validated_left), Inches(validated_top),
                        Inches(validated_width), Inches(validated_height),
                        chart_data
                    )
                    
                    if title:
                        chart = chart_shape.chart
                        chart.has_title = True
                        chart.chart_title.text_frame.text = title
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
                        chart.chart_title.text_frame.paragraphs[0].font.bold = True
                    
                elif chart_type in ["radar", "radar_filled"]:
                    # Radar charts
                    if "categories" not in data or "series" not in data:
                        return f"Error: Radar charts require 'categories' and 'series' in data. Got: {list(data.keys())}"
                    
                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
                    from pptx.util import Inches, Pt
                    
                    chart_data = CategoryChartData()
                    chart_data.categories = data["categories"]
                    
                    for series_name, values in data["series"].items():
                        chart_data.add_series(series_name, values)
                    
                    chart_enum = XL_CHART_TYPE.RADAR_FILLED if chart_type == "radar_filled" else XL_CHART_TYPE.RADAR
                    
                    chart_shape = slide.shapes.add_chart(
                        chart_enum,
                        Inches(validated_left), Inches(validated_top),
                        Inches(validated_width), Inches(validated_height),
                        chart_data
                    )
                    
                    chart = chart_shape.chart
                    if title:
                        chart.has_title = True
                        chart.chart_title.text_frame.text = title
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
                        chart.chart_title.text_frame.paragraphs[0].font.bold = True
                    
                    chart.has_legend = options.get("show_legend", True)
                    if chart.has_legend:
                        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                        chart.legend.include_in_layout = False
                    
                elif chart_type == "waterfall":
                    # Waterfall charts
                    if "categories" not in data or "values" not in data:
                        return f"Error: Waterfall charts require 'categories' and 'values' in data. Got: {list(data.keys())}"
                    
                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    from pptx.util import Inches, Pt
                    from pptx.dml.color import RGBColor
                    
                    # Calculate cumulative values for waterfall effect
                    values = data["values"]
                    cumulative = []
                    running_total = 0
                    for val in values:
                        running_total += val
                        cumulative.append(running_total)
                    
                    chart_data = CategoryChartData()
                    chart_data.categories = data["categories"]
                    chart_data.add_series('Values', cumulative)
                    
                    chart_shape = slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(validated_left), Inches(validated_top),
                        Inches(validated_width), Inches(validated_height),
                        chart_data
                    )
                    
                    chart = chart_shape.chart
                    
                    if title:
                        chart.has_title = True
                        chart.chart_title.text_frame.text = title
                        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(18)
                        chart.chart_title.text_frame.paragraphs[0].font.bold = True
                    
                    # Color bars based on positive/negative values
                    series = chart.series[0]
                    for idx, val in enumerate(values):
                        point = series.points[idx]
                        fill = point.format.fill
                        fill.solid()
                        if val > 0:
                            fill.fore_color.rgb = RGBColor(0, 176, 80)  # Green
                        elif val < 0:
                            fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
                        else:
                            fill.fore_color.rgb = RGBColor(128, 128, 128)  # Gray
                    
                    chart.has_legend = False  # No legend for waterfall
                
                else:
                    return f"Error: Unsupported chart_type '{chart_type}'"
                
                # Update in VFS if enabled
                manager.update(presentation)
                
                # Report if position was adjusted
                position_note = ""
                if (validated_left != left or validated_top != top or 
                    validated_width != width or validated_height != height):
                    position_note = f" (position adjusted to fit: {validated_left:.1f}, {validated_top:.1f}, {validated_width:.1f}x{validated_height:.1f})"
                
                return f"Added {chart_type} chart to slide {slide_index}{position_note}"
                
            except KeyError as e:
                return f"Error: Missing required data field: {str(e)}. Check the data structure for {chart_type} charts."
            except ValueError as e:
                return f"Error: Invalid data values: {str(e)}"
            except Exception as e:
                return f"Error adding {chart_type} chart: {str(e)}"
        
        return await asyncio.get_event_loop().run_in_executor(None, _add_unified_chart)
    
    # Return the tool for external access
    return {'pptx_add_chart': pptx_add_chart}