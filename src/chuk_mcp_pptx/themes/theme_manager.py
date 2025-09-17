"""
Theme manager for PowerPoint presentations.
Central system for managing and applying themes.
"""

from typing import Dict, Any, Optional, List
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ..tokens.colors import get_semantic_tokens, GRADIENTS, PALETTE
from ..tokens.typography import FONT_FAMILIES, get_text_style
from ..tokens.spacing import SPACING, RADIUS


class ThemeManager:
    """
    Manages themes for PowerPoint presentations.
    Provides theme registration, selection, and application.
    """
    
    def __init__(self):
        """Initialize theme manager with built-in themes."""
        self.themes = {}
        self._register_builtin_themes()
        self.current_theme = None
    
    def _register_builtin_themes(self):
        """Register all built-in themes."""
        # Dark themes
        self.register_theme(Theme("dark", primary_hue="blue", mode="dark"))
        self.register_theme(Theme("dark-blue", primary_hue="blue", mode="dark"))
        self.register_theme(Theme("dark-violet", primary_hue="violet", mode="dark"))
        self.register_theme(Theme("dark-green", primary_hue="emerald", mode="dark"))
        self.register_theme(Theme("dark-orange", primary_hue="orange", mode="dark"))
        self.register_theme(Theme("dark-red", primary_hue="red", mode="dark"))
        self.register_theme(Theme("dark-pink", primary_hue="pink", mode="dark"))
        self.register_theme(Theme("dark-purple", primary_hue="purple", mode="dark"))
        
        # Light themes
        self.register_theme(Theme("light", primary_hue="blue", mode="light"))
        self.register_theme(Theme("light-blue", primary_hue="blue", mode="light"))
        self.register_theme(Theme("light-violet", primary_hue="violet", mode="light"))
        self.register_theme(Theme("light-green", primary_hue="emerald", mode="light"))
        self.register_theme(Theme("light-orange", primary_hue="orange", mode="light"))
        self.register_theme(Theme("light-warm", primary_hue="amber", mode="light"))
        
        # Special themes
        self.register_theme(CyberpunkTheme())
        self.register_theme(GradientTheme("sunset", GRADIENTS["sunset"]))
        self.register_theme(GradientTheme("ocean", GRADIENTS["ocean"]))
        self.register_theme(GradientTheme("aurora", GRADIENTS["aurora"]))
        self.register_theme(MinimalTheme())
        self.register_theme(CorporateTheme())
    
    def register_theme(self, theme: 'Theme'):
        """Register a theme."""
        self.themes[theme.name] = theme
    
    def get_theme(self, name: str) -> Optional['Theme']:
        """Get theme by name."""
        return self.themes.get(name)
    
    def set_current_theme(self, name: str):
        """Set the current active theme."""
        theme = self.get_theme(name)
        if theme:
            self.current_theme = theme
            return theme
        raise ValueError(f"Theme '{name}' not found")
    
    def list_themes(self) -> List[str]:
        """List all available theme names."""
        return list(self.themes.keys())
    
    def apply_to_slide(self, slide, theme_name: Optional[str] = None):
        """
        Apply theme to a slide.
        
        Args:
            slide: PowerPoint slide object
            theme_name: Theme name or None for current theme
        """
        theme = self.get_theme(theme_name) if theme_name else self.current_theme
        if not theme:
            theme = self.get_theme("dark")  # Default
        
        theme.apply_to_slide(slide)


class Theme:
    """
    Base theme class.
    """
    
    def __init__(self, name: str, primary_hue: str = "blue", 
                 mode: str = "dark", font_family: str = "Inter"):
        """
        Initialize theme.
        
        Args:
            name: Theme name
            primary_hue: Primary color hue
            mode: Color mode (dark/light)
            font_family: Primary font family
        """
        self.name = name
        self.primary_hue = primary_hue
        self.mode = mode
        self.font_family = font_family
        self.tokens = get_semantic_tokens(primary_hue, mode)
    
    def hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex to RGB."""
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def get_color(self, path: str) -> RGBColor:
        """Get color from tokens."""
        parts = path.split('.')
        value = self.tokens
        
        for part in parts:
            if isinstance(value, dict):
                value = value.get(part, "#000000")
            else:
                break
        
        if isinstance(value, str):
            return RGBColor(*self.hex_to_rgb(value))
        return RGBColor(0, 0, 0)
    
    def apply_to_slide(self, slide):
        """Apply theme to slide."""
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.get_color("background.DEFAULT")
    
    def apply_to_shape(self, shape, style: str = "card"):
        """Apply theme to shape."""
        style_map = {
            "card": ("card.DEFAULT", "card.foreground", "border.DEFAULT"),
            "primary": ("primary.DEFAULT", "primary.foreground", "primary.DEFAULT"),
            "secondary": ("secondary.DEFAULT", "secondary.foreground", "border.DEFAULT"),
            "accent": ("accent.DEFAULT", "accent.foreground", "accent.DEFAULT"),
            "muted": ("muted.DEFAULT", "muted.foreground", "border.secondary"),
        }
        
        bg_path, fg_path, border_path = style_map.get(style, style_map["card"])
        
        # Apply fill
        if hasattr(shape, 'fill'):
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.get_color(bg_path)
        
        # Apply border
        if hasattr(shape, 'line'):
            shape.line.color.rgb = self.get_color(border_path)
            shape.line.width = Pt(1)
        
        # Apply text color
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self.get_color(fg_path)
                    run.font.name = self.font_family
    
    def get_chart_colors(self) -> List[RGBColor]:
        """Get chart colors for data visualization."""
        chart_colors = self.tokens.get("chart", [])
        return [RGBColor(*self.hex_to_rgb(color)) for color in chart_colors]


class CyberpunkTheme(Theme):
    """Cyberpunk theme with neon colors."""
    
    def __init__(self):
        super().__init__("cyberpunk", primary_hue="violet", mode="dark")
        
        # Override with cyberpunk colors
        self.tokens.update({
            "background": {"DEFAULT": "#0a0014"},
            "foreground": {"DEFAULT": "#00ffff"},
            "primary": {
                "DEFAULT": "#ff00ff",
                "foreground": "#0a0014",
            },
            "accent": {
                "DEFAULT": "#ffff00",
                "foreground": "#0a0014",
            },
            "border": {"DEFAULT": "#ff00ff"},
            "chart": ["#ff00ff", "#00ffff", "#ffff00", "#ff1493", "#00ff00"],
        })
        self.font_family = "Orbitron"


class GradientTheme(Theme):
    """Theme with gradient backgrounds."""
    
    def __init__(self, name: str, gradient_colors: List[str]):
        super().__init__(name, mode="dark")
        self.gradient_colors = gradient_colors
    
    def apply_to_slide(self, slide):
        """Apply gradient background (using first color as fallback)."""
        # PowerPoint doesn't easily support gradients via python-pptx
        # Use first color as solid background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(self.gradient_colors[0]))


class MinimalTheme(Theme):
    """Minimal black and white theme."""
    
    def __init__(self):
        super().__init__("minimal", primary_hue="zinc", mode="light")
        self.tokens = {
            "background": {"DEFAULT": "#ffffff"},
            "foreground": {"DEFAULT": "#000000"},
            "card": {
                "DEFAULT": "#f8f8f8",
                "foreground": "#000000",
            },
            "primary": {
                "DEFAULT": "#000000",
                "foreground": "#ffffff",
            },
            "secondary": {
                "DEFAULT": "#f0f0f0",
                "foreground": "#000000",
            },
            "accent": {
                "DEFAULT": "#000000",
                "foreground": "#ffffff",
            },
            "muted": {
                "DEFAULT": "#f5f5f5",
                "foreground": "#666666",
            },
            "border": {
                "DEFAULT": "#e0e0e0",
                "secondary": "#d0d0d0",
            },
            "chart": ["#000000", "#333333", "#666666", "#999999", "#cccccc"],
        }
        self.font_family = "Helvetica Neue"


class CorporateTheme(Theme):
    """Professional corporate theme."""
    
    def __init__(self):
        super().__init__("corporate", primary_hue="blue", mode="light")
        self.tokens.update({
            "background": {"DEFAULT": "#f7f9fb"},
            "foreground": {"DEFAULT": "#1a202c"},
            "primary": {
                "DEFAULT": "#2b6cb0",
                "foreground": "#ffffff",
            },
            "secondary": {
                "DEFAULT": "#e2e8f0",
                "foreground": "#2d3748",
            },
            "accent": {
                "DEFAULT": "#48bb78",
                "foreground": "#ffffff",
            },
            "chart": ["#2b6cb0", "#48bb78", "#ed8936", "#9f7aea", "#38b2ac", "#f56565"],
        })
        self.font_family = "Segoe UI"