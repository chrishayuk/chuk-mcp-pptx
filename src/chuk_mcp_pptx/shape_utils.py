"""
Shape and Formatting Utilities for PowerPoint MCP Server

Provides utilities for shapes, images, and visual effects.
"""
from typing import Dict, List, Any, Optional, Tuple
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import base64
import io
from pathlib import Path


def add_shape(slide, shape_type: str, left: float, top: float,
             width: float, height: float,
             text: str = None, fill_color: Tuple[int, int, int] = None,
             line_color: Tuple[int, int, int] = None,
             line_width: float = 1.0,
             use_theme: bool = True) -> Any:
    """
    Add a shape to a slide.
    
    Args:
        slide: Slide to add shape to
        shape_type: Type of shape (rectangle, oval, arrow, etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Optional text to add to shape
        fill_color: RGB tuple for fill color (if None, uses theme)
        line_color: RGB tuple for line color (if None, uses theme)
        line_width: Line width in points
        use_theme: Whether to use theme colors when colors not specified
        
    Returns:
        The created shape
    """
    # Import here to avoid circular dependency
    from .legacy_themes import get_slide_theme, hex_to_rgb
    # Map shape types to MSO_SHAPE constants
    shape_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "diamond": MSO_SHAPE.DIAMOND,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "star": MSO_SHAPE.STAR_5_POINT,
        "hexagon": MSO_SHAPE.HEXAGON,
        "chevron": MSO_SHAPE.CHEVRON,
        "plus": MSO_SHAPE.MATH_PLUS,
        "callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT
    }
    
    mso_shape = shape_map.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)
    
    # Add the shape
    shape = slide.shapes.add_shape(
        mso_shape,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    
    # Determine colors based on theme if not specified
    if use_theme and (not fill_color or not line_color):
        theme = get_slide_theme(slide)
        
        # Determine style based on text or shape type
        if text and text.lower() in ['primary', 'accent']:
            theme_fill = theme["primary"]
            theme_text = theme["primary_foreground"]
            theme_line = theme["primary"]
        elif text and text.lower() in ['secondary', 'muted']:
            theme_fill = theme["secondary"]
            theme_text = theme["secondary_foreground"]
            theme_line = theme["border"]
        else:
            # Use accent colors for shapes by default
            theme_fill = theme["accent"]
            theme_text = theme["accent_foreground"]
            theme_line = theme["border"]
        
        if not fill_color:
            fill_color = hex_to_rgb(theme_fill)
        if not line_color:
            line_color = hex_to_rgb(theme_line)
    
    # Set fill color
    if fill_color:
        fill = shape.fill
        fill.solid()
        r, g, b = fill_color
        fill.fore_color.rgb = RGBColor(r, g, b)
    
    # Set line properties
    line = shape.line
    if line_color:
        line.color.rgb = RGBColor(*line_color)
    line.width = Pt(line_width)
    
    # Add text if provided
    if text and shape.has_text_frame:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        # Center the text
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Set text color based on theme
        if use_theme:
            theme = get_slide_theme(slide)
            text_color = hex_to_rgb(theme["foreground"])
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*text_color)
    
    return shape


def add_connector(slide, start_x: float, start_y: float,
                 end_x: float, end_y: float,
                 connector_type: str = "straight",
                 line_color: Tuple[int, int, int] = None,
                 line_width: float = 2.0,
                 arrow_start: bool = False,
                 arrow_end: bool = True,
                 use_theme: bool = True) -> Any:
    """
    Add a connector line between two points.
    
    Args:
        slide: Slide to add connector to
        start_x: Starting X position in inches
        start_y: Starting Y position in inches
        end_x: Ending X position in inches
        end_y: Ending Y position in inches
        connector_type: Type of connector (straight, elbow, curved)
        line_color: RGB tuple for line color
        line_width: Line width in points
        arrow_start: Whether to add arrow at start
        arrow_end: Whether to add arrow at end
        
    Returns:
        The created connector shape
    """
    connector_map = {
        "straight": MSO_CONNECTOR.STRAIGHT,
        "elbow": MSO_CONNECTOR.ELBOW,
        "curved": MSO_CONNECTOR.CURVE
    }
    
    mso_connector = connector_map.get(connector_type.lower(), MSO_CONNECTOR.STRAIGHT)
    
    # Use theme color if not specified
    if use_theme and not line_color:
        from .legacy_themes import get_slide_theme, hex_to_rgb
        theme = get_slide_theme(slide)
        line_color = hex_to_rgb(theme["muted_foreground"])
    elif not line_color:
        line_color = (0, 0, 0)  # Default black
    
    # Add connector
    connector = slide.shapes.add_connector(
        mso_connector,
        Inches(start_x), Inches(start_y),
        Inches(end_x), Inches(end_y)
    )
    
    # Style the connector
    line = connector.line
    line.color.rgb = RGBColor(*line_color)
    line.width = Pt(line_width)
    
    # Add arrows using XML manipulation since python-pptx doesn't have direct support
    from pptx.oxml import parse_xml
    
    # Get or add line element
    line_elem = connector._element.spPr.ln if hasattr(connector._element.spPr, 'ln') else None
    if line_elem is None:
        # If no line element exists, we need to add one
        from pptx.oxml.ns import qn
        line_elem = connector._element.spPr._add_ln()
    
    # Add arrow heads as needed
    if arrow_end:
        # Add arrow at the end
        headEnd = parse_xml(
            '<a:headEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
        line_elem.append(headEnd)
    
    if arrow_start:
        # Add arrow at the beginning
        tailEnd = parse_xml(
            '<a:tailEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )
        line_elem.append(tailEnd)
    
    return connector


def add_image(slide, image_source: str, left: float, top: float,
             width: float = None, height: float = None,
             maintain_ratio: bool = True) -> Any:
    """
    Add an image to a slide.
    
    Args:
        slide: Slide to add image to
        image_source: Path to image file or base64 data
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional, maintains ratio if not specified)
        height: Height in inches (optional, maintains ratio if not specified)
        maintain_ratio: Whether to maintain aspect ratio
        
    Returns:
        The created picture shape
    """
    # Handle base64 image data
    if image_source.startswith("data:image/"):
        header, encoded = image_source.split(",", 1)
        image_data = base64.b64decode(encoded)
        image_stream = io.BytesIO(image_data)
        
        if width and height:
            pic = slide.shapes.add_picture(
                image_stream,
                Inches(left), Inches(top),
                width=Inches(width), height=Inches(height)
            )
        elif width:
            pic = slide.shapes.add_picture(
                image_stream,
                Inches(left), Inches(top),
                width=Inches(width)
            )
        elif height:
            pic = slide.shapes.add_picture(
                image_stream,
                Inches(left), Inches(top),
                height=Inches(height)
            )
        else:
            pic = slide.shapes.add_picture(
                image_stream,
                Inches(left), Inches(top)
            )
    # Handle file path
    elif Path(image_source).exists():
        if width and height:
            pic = slide.shapes.add_picture(
                image_source,
                Inches(left), Inches(top),
                width=Inches(width), height=Inches(height)
            )
        elif width:
            pic = slide.shapes.add_picture(
                image_source,
                Inches(left), Inches(top),
                width=Inches(width)
            )
        elif height:
            pic = slide.shapes.add_picture(
                image_source,
                Inches(left), Inches(top),
                height=Inches(height)
            )
        else:
            pic = slide.shapes.add_picture(
                image_source,
                Inches(left), Inches(top)
            )
    else:
        raise FileNotFoundError(f"Image not found: {image_source}")
    
    return pic


def apply_picture_effects(picture_shape,
                         shadow: bool = False,
                         reflection: bool = False,
                         glow: bool = False,
                         soft_edges: bool = False,
                         brightness: float = 0.0,
                         contrast: float = 0.0):
    """
    Apply visual effects to a picture.
    
    Args:
        picture_shape: Picture shape to apply effects to
        shadow: Whether to add shadow
        reflection: Whether to add reflection
        glow: Whether to add glow effect
        soft_edges: Whether to add soft edges
        brightness: Brightness adjustment (-1.0 to 1.0)
        contrast: Contrast adjustment (-1.0 to 1.0)
    """
    # Note: Some effects require XML manipulation in python-pptx
    # This is a simplified version
    
    # Adjust brightness and contrast
    if brightness != 0.0 or contrast != 0.0:
        picture = picture_shape._element
        # This would require direct XML manipulation
        # which is complex in python-pptx
    
    # Add shadow
    if shadow:
        # Simplified shadow effect
        shadow_format = picture_shape.shadow
        shadow_format.inherit = False
        shadow_format.visible = True
        shadow_format.distance = Pt(4)
        shadow_format.blur_radius = Pt(4)
        shadow_format.transparency = 0.5
        shadow_format.angle = 45


def add_text_box_with_style(slide, left: float, top: float,
                           width: float, height: float,
                           text: str, style_preset: str = "default") -> Any:
    """
    Add a styled text box using preset styles.
    
    Args:
        slide: Slide to add text box to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        style_preset: Style preset name
        
    Returns:
        The created text box shape
    """
    # Style presets
    styles = {
        "title": {
            "font_size": 32,
            "bold": True,
            "color": (0, 120, 215),
            "alignment": "center"
        },
        "subtitle": {
            "font_size": 24,
            "bold": False,
            "color": (68, 68, 68),
            "alignment": "center"
        },
        "heading": {
            "font_size": 20,
            "bold": True,
            "color": (51, 51, 51),
            "alignment": "left"
        },
        "body": {
            "font_size": 16,
            "bold": False,
            "color": (51, 51, 51),
            "alignment": "left"
        },
        "caption": {
            "font_size": 12,
            "bold": False,
            "color": (136, 136, 136),
            "alignment": "left"
        },
        "quote": {
            "font_size": 18,
            "bold": False,
            "italic": True,
            "color": (68, 68, 68),
            "alignment": "center"
        },
        "highlight": {
            "font_size": 18,
            "bold": True,
            "color": (255, 140, 0),
            "alignment": "center"
        },
        "default": {
            "font_size": 16,
            "bold": False,
            "color": (0, 0, 0),
            "alignment": "left"
        }
    }
    
    style = styles.get(style_preset, styles["default"])
    
    # Add text box
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    
    # Apply style
    from pptx.enum.text import PP_ALIGN
    
    for paragraph in text_frame.paragraphs:
        # Set alignment
        alignment = style.get("alignment", "left")
        if alignment == "center":
            paragraph.alignment = PP_ALIGN.CENTER
        elif alignment == "right":
            paragraph.alignment = PP_ALIGN.RIGHT
        elif alignment == "justify":
            paragraph.alignment = PP_ALIGN.JUSTIFY
        else:
            paragraph.alignment = PP_ALIGN.LEFT
        
        # Format font
        font = paragraph.font
        font.size = Pt(style.get("font_size", 16))
        font.bold = style.get("bold", False)
        font.italic = style.get("italic", False)
        
        color = style.get("color", (0, 0, 0))
        font.color.rgb = RGBColor(*color)
    
    return text_box


def add_smart_art(slide, layout_type: str, items: List[str],
                 left: float, top: float, width: float, height: float,
                 color_scheme: str = "modern_blue",
                 use_theme: bool = True) -> List[Any]:
    """
    Create a SmartArt-like diagram using shapes.
    
    Args:
        slide: Slide to add diagram to
        layout_type: Type of layout (process, cycle, hierarchy, list)
        items: List of items for the diagram
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        color_scheme: Color scheme to use (ignored if use_theme is True)
        use_theme: Whether to use slide theme colors
        
    Returns:
        List of created shapes
    """
    from .slide_templates import COLOR_SCHEMES
    from .legacy_themes import get_slide_theme, hex_to_rgb
    
    shapes = []
    
    # Use theme colors if available
    if use_theme:
        theme = get_slide_theme(slide)
        colors = {
            "primary": hex_to_rgb(theme["primary"]),
            "secondary": hex_to_rgb(theme["secondary"]),
            "accent": hex_to_rgb(theme["accent"]),
            "text": hex_to_rgb(theme["foreground"]),
            "background": hex_to_rgb(theme["card"])
        }
    else:
        colors = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES["modern_blue"])
    
    if layout_type == "process":
        # Create process flow diagram with proper spacing
        # Calculate spacing between items (minimum 0.2 inches)
        num_items = len(items)
        min_spacing = 0.2
        
        # Calculate item width considering minimum spacing
        total_spacing = min_spacing * (num_items - 1)
        available_width = width - total_spacing
        item_width = available_width / num_items
        
        # Ensure item width is reasonable (max 1.5 inches for readability)
        if item_width > 1.5:
            item_width = 1.5
            total_width_needed = item_width * num_items + total_spacing
            # Center the process flow if there's extra space
            start_offset = (width - total_width_needed) / 2
        else:
            start_offset = 0
        
        for idx, item in enumerate(items):
            x = left + start_offset + idx * (item_width + min_spacing)
            
            # Add shape
            shape = add_shape(
                slide, "chevron" if idx < len(items) - 1 else "rectangle",
                x, top, item_width, height * 0.8,  # Use 80% of height for better proportions
                text=item,
                fill_color=colors["primary"] if idx % 2 == 0 else colors["secondary"]
            )
            shapes.append(shape)
            
            # Add connector to next shape
            if idx < len(items) - 1:
                next_x = left + start_offset + (idx + 1) * (item_width + min_spacing)
                connector = add_connector(
                    slide,
                    x + item_width,  # Right edge of current shape
                    top + height * 0.4,  # Middle of shape height
                    next_x,  # Left edge of next shape
                    top + height * 0.4,  # Middle of shape height
                    connector_type="straight",
                    line_color=colors["text"],
                    arrow_end=True
                )
                shapes.append(connector)
            
    elif layout_type == "cycle":
        # Create cycle diagram with proper spacing
        import math
        center_x = left + width / 2
        center_y = top + height / 2
        
        # Adjust radius and shape size based on number of items
        num_items = len(items)
        if num_items <= 4:
            radius = min(width, height) / 2.5
            shape_width = 1.2
            shape_height = 0.8
        elif num_items <= 6:
            radius = min(width, height) / 2.2
            shape_width = 1.0
            shape_height = 0.7
        else:
            radius = min(width, height) / 2.0
            shape_width = 0.9
            shape_height = 0.6
        
        for idx, item in enumerate(items):
            angle = 2 * math.pi * idx / num_items - math.pi / 2
            x = center_x + radius * math.cos(angle) - shape_width/2
            y = center_y + radius * math.sin(angle) - shape_height/2
            
            shape = add_shape(
                slide, "rounded_rectangle",
                x, y, shape_width, shape_height,
                text=item,
                fill_color=colors["accent"] if idx % 2 == 0 else colors["secondary"]
            )
            shapes.append(shape)
            
            # Add connector to next item with cleaner edge calculation
            next_idx = (idx + 1) % num_items
            next_angle = 2 * math.pi * next_idx / num_items - math.pi / 2
            next_x = center_x + radius * math.cos(next_angle) - shape_width/2
            next_y = center_y + radius * math.sin(next_angle) - shape_height/2
            
            # Current shape center
            curr_cx = x + shape_width/2
            curr_cy = y + shape_height/2
            
            # Next shape center
            next_cx = next_x + shape_width/2
            next_cy = next_y + shape_height/2
            
            # Calculate direction vector from current to next
            dx = next_cx - curr_cx
            dy = next_cy - curr_cy
            dist = math.sqrt(dx*dx + dy*dy)
            
            if dist > 0:
                # Normalize direction
                dx /= dist
                dy /= dist
                
                # Calculate edge points
                # Start from current shape edge (towards next shape)
                start_x = curr_cx + dx * shape_width * 0.5
                start_y = curr_cy + dy * shape_height * 0.5
                
                # End at next shape edge (from current shape direction)
                end_x = next_cx - dx * shape_width * 0.5
                end_y = next_cy - dy * shape_height * 0.5
                
                connector = add_connector(
                    slide,
                    start_x, start_y,
                    end_x, end_y,
                    connector_type="curved",
                    line_color=colors["text"],
                    arrow_end=True
                )
                shapes.append(connector)
    
    elif layout_type == "list":
        # Create vertical list with shapes containing text
        item_height = height / len(items)
        
        for idx, item in enumerate(items):
            y = top + idx * item_height
            
            # Add rounded rectangle with text
            shape = add_shape(
                slide, "rounded_rectangle",
                left, y, width, item_height * 0.8,
                text=item,
                fill_color=colors["accent"] if idx % 2 == 0 else colors["secondary"],
                line_color=colors["text"]
            )
            shapes.append(shape)
    
    elif layout_type == "hierarchy":
        # Create hierarchy diagram with connectors
        if len(items) > 0:
            # Top level (CEO/Root)
            top_x = left + width/2 - 1.5
            top_y = top
            top_shape = add_shape(
                slide, "rounded_rectangle",
                top_x, top_y, 3.0, 0.8,
                text=items[0],
                fill_color=colors["primary"],
                line_color=colors["text"]
            )
            shapes.append(top_shape)
            
            # Second level items with connectors
            if len(items) > 1:
                remaining = items[1:]
                num_items = len(remaining)
                item_width = min(1.8, (width - 0.3 * (num_items - 1)) / num_items)
                
                # Calculate starting position to center the items
                total_width = num_items * item_width + (num_items - 1) * 0.3
                start_x = left + (width - total_width) / 2
                
                for idx, item in enumerate(remaining):
                    x = start_x + idx * (item_width + 0.3)
                    y = top + 2.0
                    
                    # Add shape
                    shape = add_shape(
                        slide, "rectangle",
                        x, y, item_width, 0.8,
                        text=item,
                        fill_color=colors["secondary"],
                        line_color=colors["text"]
                    )
                    shapes.append(shape)
                    
                    # Add connector from top to this shape
                    connector = add_connector(
                        slide,
                        top_x + 1.5,  # Center bottom of top shape
                        top_y + 0.8,  # Bottom of top shape
                        x + item_width/2,  # Center top of child shape
                        y,  # Top of child shape
                        connector_type="straight",
                        line_color=colors["text"],
                        arrow_end=True
                    )
                    shapes.append(connector)
    
    elif layout_type == "relationship":
        # Create relationship diagram with connected shapes and better spacing
        import math
        if len(items) >= 1:
            # Center shape
            center_width = 1.5
            center_height = 0.8
            center_x = left + width/2 - center_width/2
            center_y = top + height/2 - center_height/2
            
            center_shape = add_shape(
                slide, "oval",
                center_x, center_y, center_width, center_height,
                text=items[0] if items else "Center",
                fill_color=colors["primary"],
                line_color=colors["text"]
            )
            shapes.append(center_shape)
            
            # Surrounding shapes with proper spacing
            if len(items) > 1:
                remaining = items[1:]
                num_surrounding = len(remaining)
                
                # Calculate appropriate radius to avoid overlaps
                # Make radius larger to spread out shapes more
                radius_x = min((width - 1.5) / 2, 2.5)
                radius_y = min((height - 1.0) / 2, 1.5)
                
                for idx, item in enumerate(remaining):
                    angle = 2 * math.pi * idx / num_surrounding
                    
                    # Position shapes around the center
                    shape_width = 1.2
                    shape_height = 0.7
                    x = left + width/2 + radius_x * math.cos(angle) - shape_width/2
                    y = top + height/2 + radius_y * math.sin(angle) - shape_height/2
                    
                    shape = add_shape(
                        slide, "rounded_rectangle",
                        x, y, shape_width, shape_height,
                        text=item,
                        fill_color=colors["accent"],
                        line_color=colors["text"]
                    )
                    shapes.append(shape)
                    
                    # Add connector from center to this shape with proper edge calculation
                    # Calculate connection points on shape edges
                    # For center shape (oval), calculate point on edge
                    center_cx = center_x + center_width/2
                    center_cy = center_y + center_height/2
                    
                    # For surrounding shape center
                    shape_cx = x + shape_width/2
                    shape_cy = y + shape_height/2
                    
                    # Calculate unit vector from center to shape
                    dx = shape_cx - center_cx
                    dy = shape_cy - center_cy
                    distance = math.sqrt(dx*dx + dy*dy)
                    
                    if distance > 0:
                        # Normalize
                        dx /= distance
                        dy /= distance
                        
                        # Calculate edge points
                        # Start point on center shape edge (oval)
                        center_edge_x = center_cx + dx * center_width/2 * 0.9
                        center_edge_y = center_cy + dy * center_height/2 * 0.9
                        
                        # End point on surrounding shape edge
                        shape_edge_x = shape_cx - dx * shape_width/2 * 0.9
                        shape_edge_y = shape_cy - dy * shape_height/2 * 0.9
                        
                        connector = add_connector(
                            slide,
                            center_edge_x, center_edge_y,
                            shape_edge_x, shape_edge_y,
                            connector_type="straight",
                            line_color=colors["text"],
                            arrow_start=True,
                            arrow_end=True
                        )
                        shapes.append(connector)
    
    elif layout_type == "pyramid":
        # Create pyramid diagram
        if len(items) > 0:
            level_height = height / len(items)
            for idx, item in enumerate(items):
                # Width increases as we go down
                level_width = width * (idx + 1) / len(items)
                x = left + (width - level_width) / 2
                y = top + idx * level_height
                
                shape = add_shape(
                    slide, "triangle" if idx == 0 else "hexagon",
                    x, y, level_width, level_height * 0.8,
                    text=item,
                    fill_color=colors["primary"] if idx == 0 else colors["secondary"] if idx % 2 == 1 else colors["accent"],
                    line_color=colors["text"]
                )
                shapes.append(shape)
    
    return shapes