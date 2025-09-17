"""
Slide Templates for PowerPoint MCP Server

Provides pre-designed slide layouts and templates for quick presentation creation.
"""
from typing import Dict, List, Any, Optional
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Professional color schemes
COLOR_SCHEMES = {
    "modern_blue": {
        "primary": (0, 120, 215),      # Microsoft Blue
        "secondary": (0, 176, 240),    # Light Blue
        "accent": (255, 140, 0),        # Orange
        "text": (51, 51, 51),           # Dark Gray
        "light": (242, 242, 242)        # Light Gray
    },
    "corporate_gray": {
        "primary": (68, 68, 68),        # Dark Gray
        "secondary": (136, 136, 136),   # Medium Gray
        "accent": (0, 112, 192),        # Blue
        "text": (51, 51, 51),           # Text Gray
        "light": (240, 240, 240)        # Light Gray
    },
    "elegant_green": {
        "primary": (34, 139, 34),       # Forest Green
        "secondary": (144, 238, 144),   # Light Green
        "accent": (255, 215, 0),        # Gold
        "text": (51, 51, 51),           # Dark Gray
        "light": (245, 255, 250)        # Mint Cream
    },
    "warm_orange": {
        "primary": (255, 140, 0),       # Dark Orange
        "secondary": (255, 218, 185),   # Peach
        "accent": (220, 20, 60),        # Crimson
        "text": (51, 51, 51),           # Dark Gray
        "light": (255, 250, 240)        # Floral White
    }
}

# Slide layout templates
SLIDE_TEMPLATES = {
    "title_slide": {
        "name": "Title Slide",
        "description": "Professional title slide with subtitle and date",
        "layout_index": 0,
        "elements": ["title", "subtitle", "date"],
        "styling": {
            "title_font_size": 44,
            "subtitle_font_size": 24,
            "date_font_size": 18
        }
    },
    "section_header": {
        "name": "Section Header",
        "description": "Section divider with large title",
        "layout_index": 2,
        "elements": ["section_title", "section_subtitle"],
        "styling": {
            "title_font_size": 40,
            "subtitle_font_size": 20
        }
    },
    "content_slide": {
        "name": "Content Slide",
        "description": "Standard content slide with title and bullets",
        "layout_index": 1,
        "elements": ["title", "content"],
        "styling": {
            "title_font_size": 32,
            "content_font_size": 18,
            "bullet_level_indent": 0.5
        }
    },
    "two_column": {
        "name": "Two Column Layout",
        "description": "Two equal columns with headers",
        "layout_index": 3,
        "elements": ["title", "left_header", "left_content", "right_header", "right_content"],
        "styling": {
            "title_font_size": 32,
            "header_font_size": 24,
            "content_font_size": 16
        }
    },
    "comparison": {
        "name": "Comparison Slide",
        "description": "Before/After or comparison layout",
        "layout_index": 3,
        "elements": ["title", "left_title", "left_items", "right_title", "right_items"],
        "styling": {
            "title_font_size": 32,
            "subtitle_font_size": 24,
            "item_font_size": 16
        }
    },
    "image_with_caption": {
        "name": "Image with Caption",
        "description": "Large image with title and caption",
        "layout_index": 8,
        "elements": ["title", "image", "caption"],
        "styling": {
            "title_font_size": 32,
            "caption_font_size": 14
        }
    },
    "quote": {
        "name": "Quote Slide",
        "description": "Featured quote with attribution",
        "layout_index": 1,
        "elements": ["quote", "attribution"],
        "styling": {
            "quote_font_size": 28,
            "attribution_font_size": 18,
            "quote_italic": True
        }
    },
    "key_metrics": {
        "name": "Key Metrics Dashboard",
        "description": "Dashboard with 3-4 key metrics",
        "layout_index": 5,
        "elements": ["title", "metric1", "metric2", "metric3", "metric4"],
        "styling": {
            "title_font_size": 32,
            "metric_value_font_size": 48,
            "metric_label_font_size": 14
        }
    },
    "timeline": {
        "name": "Timeline",
        "description": "Horizontal timeline with milestones",
        "layout_index": 1,
        "elements": ["title", "milestones"],
        "styling": {
            "title_font_size": 32,
            "milestone_font_size": 14
        }
    },
    "closing": {
        "name": "Closing Slide",
        "description": "Thank you slide with contact info",
        "layout_index": 0,
        "elements": ["thank_you", "contact", "website"],
        "styling": {
            "title_font_size": 44,
            "contact_font_size": 20
        }
    }
}


def apply_color_scheme(slide, scheme_name: str = "modern_blue"):
    """
    Apply a color scheme to a slide background.
    
    Args:
        slide: The slide to apply the color scheme to
        scheme_name: Name of the color scheme to apply
    """
    if scheme_name not in COLOR_SCHEMES:
        return
        
    scheme = COLOR_SCHEMES[scheme_name]
    
    # Set slide background to a subtle gradient or solid color
    background = slide.background
    fill = background.fill
    fill.solid()
    
    # Use light color for background
    r, g, b = scheme["light"]
    fill.fore_color.rgb = RGBColor(r, g, b)
    

def create_title_slide(prs, title: str, subtitle: str = "", 
                      author: str = "", date: str = "",
                      color_scheme: str = "modern_blue") -> int:
    """
    Create a professional title slide.
    
    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle text
        author: Author name
        date: Presentation date
        color_scheme: Color scheme to apply
        
    Returns:
        Index of the created slide
    """
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply color scheme
    apply_color_scheme(slide, color_scheme)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Style the title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    font = title_para.font
    font.size = Pt(44)
    font.bold = True
    r, g, b = COLOR_SCHEMES[color_scheme]["primary"]
    font.color.rgb = RGBColor(r, g, b)
    
    # Add subtitle if provided
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        
        # Add author and date on separate lines if provided
        if author or date:
            text_frame = subtitle_shape.text_frame
            if author:
                p = text_frame.add_paragraph()
                p.text = author
                p.font.size = Pt(18)
            if date:
                p = text_frame.add_paragraph()
                p.text = date
                p.font.size = Pt(16)
                p.font.italic = True
    
    return len(prs.slides) - 1


def create_content_slide(prs, title: str, content: List[str],
                        color_scheme: str = "modern_blue") -> int:
    """
    Create a content slide with bullet points.
    
    Args:
        prs: Presentation object
        title: Slide title
        content: List of bullet points
        color_scheme: Color scheme to apply
        
    Returns:
        Index of the created slide
    """
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply color scheme
    apply_color_scheme(slide, color_scheme)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Style the title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    font = title_para.font
    font.size = Pt(32)
    font.bold = True
    r, g, b = COLOR_SCHEMES[color_scheme]["primary"]
    font.color.rgb = RGBColor(r, g, b)
    
    # Add content
    if len(slide.placeholders) > 1:
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        
        # Clear existing paragraph
        text_frame.clear()
        
        # Add bullet points
        for idx, bullet in enumerate(content):
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            
            # Add spacing between bullets
            p.space_after = Pt(12)
    
    return len(prs.slides) - 1


def create_comparison_slide(prs, title: str, 
                           left_title: str, left_items: List[str],
                           right_title: str, right_items: List[str],
                           color_scheme: str = "modern_blue") -> int:
    """
    Create a comparison slide with two columns.
    
    Args:
        prs: Presentation object
        title: Main slide title
        left_title: Left column title
        left_items: Left column items
        right_title: Right column title
        right_items: Right column items
        color_scheme: Color scheme to apply
        
    Returns:
        Index of the created slide
    """
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply color scheme
    apply_color_scheme(slide, color_scheme)
    scheme = COLOR_SCHEMES[color_scheme]
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.alignment = PP_ALIGN.CENTER
    font = title_para.font
    font.size = Pt(32)
    font.bold = True
    r, g, b = scheme["primary"]
    font.color.rgb = RGBColor(r, g, b)
    
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(4), Inches(4.5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    # Left title
    left_title_para = left_frame.paragraphs[0]
    left_title_para.text = left_title
    left_title_para.alignment = PP_ALIGN.CENTER
    font = left_title_para.font
    font.size = Pt(24)
    font.bold = True
    r, g, b = scheme["secondary"]
    font.color.rgb = RGBColor(r, g, b)
    
    # Left items
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(2), Inches(4), Inches(4.5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    # Right title
    right_title_para = right_frame.paragraphs[0]
    right_title_para.text = right_title
    right_title_para.alignment = PP_ALIGN.CENTER
    font = right_title_para.font
    font.size = Pt(24)
    font.bold = True
    r, g, b = scheme["accent"]
    font.color.rgb = RGBColor(r, g, b)
    
    # Right items
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_after = Pt(6)
    
    return len(prs.slides) - 1


def create_key_metrics_slide(prs, title: str, metrics: List[Dict[str, Any]],
                            color_scheme: str = "modern_blue") -> int:
    """
    Create a key metrics dashboard slide.
    
    Args:
        prs: Presentation object
        title: Slide title
        metrics: List of dicts with 'value' and 'label' keys
        color_scheme: Color scheme to apply
        
    Returns:
        Index of the created slide
    """
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply color scheme
    apply_color_scheme(slide, color_scheme)
    scheme = COLOR_SCHEMES[color_scheme]
    
    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.alignment = PP_ALIGN.CENTER
    font = title_para.font
    font.size = Pt(32)
    font.bold = True
    r, g, b = scheme["primary"]
    font.color.rgb = RGBColor(r, g, b)
    
    # Calculate metric box positions
    num_metrics = len(metrics)
    if num_metrics == 0:
        return len(prs.slides) - 1
        
    box_width = 8.0 / num_metrics
    start_x = 1.0
    
    # Add metric boxes
    for idx, metric in enumerate(metrics):
        x = start_x + (idx * box_width)
        
        # Metric box
        metric_box = slide.shapes.add_textbox(
            Inches(x), Inches(2.5), Inches(box_width - 0.2), Inches(2)
        )
        text_frame = metric_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Metric value
        value_para = text_frame.paragraphs[0]
        value_para.text = str(metric.get('value', ''))
        value_para.alignment = PP_ALIGN.CENTER
        font = value_para.font
        font.size = Pt(48)
        font.bold = True
        r, g, b = scheme["accent"]
        font.color.rgb = RGBColor(r, g, b)
        
        # Metric label
        label_para = text_frame.add_paragraph()
        label_para.text = metric.get('label', '')
        label_para.alignment = PP_ALIGN.CENTER
        font = label_para.font
        font.size = Pt(14)
        r, g, b = scheme["text"]
        font.color.rgb = RGBColor(r, g, b)
    
    return len(prs.slides) - 1


def get_template_info(template_id: str) -> Optional[Dict[str, Any]]:
    """
    Get information about a specific slide template.
    
    Args:
        template_id: ID of the template
        
    Returns:
        Template information or None if not found
    """
    return SLIDE_TEMPLATES.get(template_id)


def list_templates() -> List[Dict[str, str]]:
    """
    List all available slide templates.
    
    Returns:
        List of template information
    """
    return [
        {
            "id": template_id,
            "name": template["name"],
            "description": template["description"]
        }
        for template_id, template in SLIDE_TEMPLATES.items()
    ]


def list_color_schemes() -> List[str]:
    """
    List all available color schemes.
    
    Returns:
        List of color scheme names
    """
    return list(COLOR_SCHEMES.keys())