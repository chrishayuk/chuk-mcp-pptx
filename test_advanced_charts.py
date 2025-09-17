#!/usr/bin/env python3
"""Test the advanced chart types from slides 10-15."""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.charts import BubbleChart, RadarChart, SparklineChart
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


async def test_bubble_chart_advanced():
    """Test the advanced bubble chart with transparency and size_scale."""
    
    print("Testing BubbleChart with advanced options...")
    
    prs = Presentation()
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)
    
    try:
        bubble_chart = BubbleChart(
            series_data=[
                {
                    "name": "Product Portfolio",
                    "points": [
                        [25, 35, 18],   
                        [45, 55, 28],
                        [35, 45, 22],
                        [55, 65, 35],
                        [65, 75, 42]
                    ]
                }
            ],
            size_scale=2.0,
            transparency=25,
            theme=theme.__dict__
        )
        
        await bubble_chart.render(slide, left=1.0, top=1.2, width=8.0, height=4.0)
        
        prs.save("outputs/test_bubble_advanced.pptx")
        print("✅ Advanced bubble chart created successfully")
        
    except Exception as e:
        print(f"❌ Advanced bubble chart error: {e}")
        import traceback
        traceback.print_exc()


async def test_radar_chart_advanced():
    """Test the radar chart with filled variant."""
    
    print("\nTesting RadarChart with filled variant...")
    
    prs = Presentation()
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)
    
    try:
        radar_chart = RadarChart(
            categories=["Performance", "Reliability", "Security", "Usability"],
            series={
                "Our Product": [9.2, 9.5, 9.0, 8.8],
                "Competitor A": [7.5, 8.8, 8.5, 9.2]
            },
            variant="filled",
            max_value=10,
            theme=theme.__dict__
        )
        
        await radar_chart.render(slide, left=1.5, top=1.2, width=7.0, height=4.0)
        
        prs.save("outputs/test_radar_advanced.pptx")
        print("✅ Advanced radar chart created successfully")
        
    except Exception as e:
        print(f"❌ Advanced radar chart error: {e}")
        import traceback
        traceback.print_exc()


async def test_sparkline_chart():
    """Test the sparkline chart with custom colors."""
    
    print("\nTesting SparklineChart with custom colors...")
    
    prs = Presentation()
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)
    
    try:
        sparkline = SparklineChart(
            values=[42, 45, 48, 52, 49, 55, 58, 62, 60, 65, 68, 72, 75, 78],
            color="#3b82f6",
            theme=theme.__dict__
        )
        
        await sparkline.render(slide, left=3.5, top=2.0, width=4.0, height=0.6)
        
        prs.save("outputs/test_sparkline.pptx")
        print("✅ Sparkline chart created successfully")
        
    except Exception as e:
        print(f"❌ Sparkline chart error: {e}")
        import traceback
        traceback.print_exc()


async def test_multiple_sparklines():
    """Test multiple sparklines on the same slide (like slide 14)."""
    
    print("\nTesting multiple SparklineCharts on same slide...")
    
    prs = Presentation()
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)
    
    try:
        sparkline_data = [
            {"label": "Revenue", "values": [42, 45, 48, 52, 49, 55], "color": "#3b82f6"},
            {"label": "Growth", "values": [25, 28, 32, 35, 38, 42], "color": "#10b981"}
        ]
        
        for i, spark_info in enumerate(sparkline_data):
            sparkline = SparklineChart(
                values=spark_info["values"],
                color=spark_info["color"],
                theme=theme.__dict__
            )
            await sparkline.render(slide, left=3.5, top=2.0 + i * 0.8, width=4.0, height=0.6)
        
        prs.save("outputs/test_multiple_sparklines.pptx")
        print("✅ Multiple sparklines created successfully")
        
    except Exception as e:
        print(f"❌ Multiple sparklines error: {e}")
        import traceback
        traceback.print_exc()


async def main():
    """Test all advanced charts."""
    await test_bubble_chart_advanced()
    await test_radar_chart_advanced() 
    await test_sparkline_chart()
    await test_multiple_sparklines()


if __name__ == "__main__":
    asyncio.run(main())