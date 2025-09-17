#!/usr/bin/env python3
"""Test slides 9-15 to find the cumulative issue."""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from chuk_mcp_pptx.components.charts import (
    DoughnutChart, ScatterChart, BubbleChart, RadarChart, SparklineChart
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager
from chuk_mcp_pptx.components import MetricCard


def add_slide_title_and_description(slide, title: str, description: str, features: list, theme_obj):
    """Add title and description to a slide."""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
    title_para.alignment = PP_ALIGN.CENTER


async def test_slides_9_to_15():
    """Test slides 9-15 to find cumulative issues."""
    
    print("Testing slides 9-15 to find cumulative issues...")
    
    # Create presentation with slides 1-8 first (basic ones)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    theme_manager = ThemeManager()
    theme_obj = theme_manager.get_theme("dark-blue")
    
    # Add 8 simple slides first
    for i in range(8):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
        title_box.text_frame.text = f"Simple Slide {i+1}"
    
    print("✅ Added 8 simple slides")
    
    try:
        # SLIDE 9: DOUGHNUT CHART
        print("Creating slide 9 (Doughnut)...")
        slide9 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide9)
        add_slide_title_and_description(slide9, "09. Doughnut Chart", "Test", [], theme_obj)
        
        doughnut_chart = DoughnutChart(
            categories=["North America", "Europe", "Asia Pacific", "Latin America"],
            values=[40, 30, 25, 5],
            hole_size=60,
            theme=theme_obj.__dict__
        )
        await doughnut_chart.render(slide9, left=2.5, top=1.2, width=5.0, height=4.0)
        print("✅ Slide 9 completed")
        
        # Save at this point
        prs.save("outputs/test_slides_1_9.pptx")
        print("✅ Saved slides 1-9 successfully")
        
        # SLIDE 10: SCATTER CHART
        print("Creating slide 10 (Scatter)...")
        slide10 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide10)
        add_slide_title_and_description(slide10, "10. Scatter Chart", "Test", [], theme_obj)
        
        scatter_chart = ScatterChart(
            series_data=[
                {
                    "name": "Performance Analysis",
                    "x_values": [15, 25, 35, 45, 55, 65],
                    "y_values": [20, 30, 40, 50, 60, 70]
                }
            ],
            theme=theme_obj.__dict__
        )
        await scatter_chart.render(slide10, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 10 completed")
        
        # SLIDE 11: BUBBLE CHART
        print("Creating slide 11 (Bubble)...")
        slide11 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide11)
        add_slide_title_and_description(slide11, "11. Bubble Chart", "Test", [], theme_obj)
        
        bubble_chart = BubbleChart(
            series_data=[
                {
                    "name": "Product Portfolio",
                    "points": [
                        [25, 35, 18],
                        [45, 55, 28],
                        [35, 45, 22]
                    ]
                }
            ],
            theme=theme_obj.__dict__  # Removed size_scale and transparency
        )
        await bubble_chart.render(slide11, left=1.0, top=1.2, width=8.0, height=4.0)
        print("✅ Slide 11 completed")
        
        # Save at this point
        prs.save("outputs/test_slides_1_11.pptx")
        print("✅ Saved slides 1-11 successfully")
        
        # SLIDE 12: RADAR CHART
        print("Creating slide 12 (Radar)...")
        slide12 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide12)
        add_slide_title_and_description(slide12, "12. Radar Chart", "Test", [], theme_obj)
        
        radar_chart = RadarChart(
            categories=["Performance", "Reliability", "Security", "Usability"],
            series={
                "Our Product": [9.2, 9.5, 9.0, 8.8],
                "Competitor A": [7.5, 8.8, 8.5, 9.2]
            },
            theme=theme_obj.__dict__  # Removed variant and max_value
        )
        await radar_chart.render(slide12, left=1.5, top=1.2, width=7.0, height=4.0)
        print("✅ Slide 12 completed")
        
        # SLIDE 13: METRIC CARDS (simplified)
        print("Creating slide 13 (Metrics)...")
        slide13 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide13)
        add_slide_title_and_description(slide13, "13. KPI Metrics", "Test", [], theme_obj)
        
        # Just add one metric card instead of 6
        metric_card = MetricCard(
            label="Customer Satisfaction",
            value="87%",
            theme=theme_obj.__dict__
        )
        metric_card.render(slide13, left=3.0, top=2.0, width=2.0, height=1.6)
        print("✅ Slide 13 completed")
        
        # Save at this point
        prs.save("outputs/test_slides_1_13.pptx")
        print("✅ Saved slides 1-13 successfully")
        
        # SLIDE 14: SPARKLINES (simplified)
        print("Creating slide 14 (Sparklines)...")
        slide14 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide14)
        add_slide_title_and_description(slide14, "14. Sparklines", "Test", [], theme_obj)
        
        # Just one sparkline instead of 4
        sparkline = SparklineChart(
            values=[42, 45, 48, 52, 49, 55, 58, 62],
            theme=theme_obj.__dict__  # Removed custom color
        )
        await sparkline.render(slide14, left=3.5, top=2.0, width=4.0, height=0.6)
        print("✅ Slide 14 completed")
        
        # SLIDE 15: SUMMARY
        print("Creating slide 15 (Summary)...")
        slide15 = prs.slides.add_slide(prs.slide_layouts[5])
        theme_obj.apply_to_slide(slide15)
        
        title_box = slide15.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "✅ All 15 Slides Created!"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_obj.get_color("foreground.DEFAULT")
        title_para.alignment = PP_ALIGN.CENTER
        print("✅ Slide 15 completed")
        
        # Final save
        prs.save("outputs/test_slides_1_15.pptx")
        print("✅ Saved all 15 slides successfully")
        
    except Exception as e:
        print(f"❌ Error during slide creation: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    asyncio.run(test_slides_9_to_15())