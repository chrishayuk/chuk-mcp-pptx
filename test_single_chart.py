#!/usr/bin/env python3
"""Test individual chart types to find the issue."""

import asyncio
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'src'))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.charts import RadarChart
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


async def test_radar_chart():
    """Test radar chart specifically."""
    
    print("Testing RadarChart...")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-blue")
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)
    
    try:
        # Simple radar chart
        radar_chart = RadarChart(
            categories=["Speed", "Quality", "Cost", "Reliability"],
            series={
                "Product A": [8, 7, 6, 9],
                "Product B": [6, 9, 8, 7]
            },
            title="Test Radar Chart",
            theme=theme.__dict__
        )
        
        await radar_chart.render(slide, left=2.0, top=2.0, width=6.0, height=4.0)
        print("✅ Radar chart created successfully")
        
        # Save
        prs.save("outputs/test_radar.pptx")
        print("✅ File saved: outputs/test_radar.pptx")
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    asyncio.run(test_radar_chart())